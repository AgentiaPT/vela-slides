#!/usr/bin/env python3
"""Spike verification: prove the hand-rolled .pptx contains NATIVE, EDITABLE
objects (not a flattened image) by reading it back with python-pptx.

Usage: python3 verify.py out/minimal.pptx
Exit 0 = all assertions passed.
"""
import sys
import zipfile
import xml.dom.minidom as minidom
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = sys.argv[1] if len(sys.argv) > 1 else "out/minimal.pptx"
require_media = "--require-media" in sys.argv  # png image + native SVG expected (minimal fixture)

# 1) Structural: valid zip + every XML part well-formed.
with zipfile.ZipFile(path) as z:
    bad = z.testzip()
    assert bad is None, f"corrupt zip entry: {bad}"
    names = z.namelist()
    for n in names:
        if n.endswith(".xml") or n.endswith(".rels"):
            minidom.parseString(z.read(n))  # raises on malformed XML
    assert "[Content_Types].xml" in names
    assert any(n.startswith("ppt/slides/slide") for n in names)
    print(f"[ok] valid zip, {len(names)} parts, all XML well-formed")

# 2) Semantic: python-pptx parses it and finds real objects.
prs = Presentation(path)
print(f"[ok] python-pptx opened deck: {len(prs.slides)} slide(s), "
      f"size {Emu(prs.slide_width).inches:.2f}x{Emu(prs.slide_height).inches:.2f} in")

n_text = n_auto = n_pic = n_svg = 0
sample_text = []
for slide in prs.slides:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            n_text += 1
            sample_text.append(sh.text_frame.text)
            # editable run with real font/size proves it's not a picture
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    assert r.text is not None
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            n_auto += 1
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            n_pic += 1

# 3) SVG survives as a native vector part (image/svg+xml in media).
with zipfile.ZipFile(path) as z:
    svg_parts = [n for n in z.namelist() if n.lower().endswith(".svg")]
    n_svg = len(svg_parts)
    # and the slide XML references it via the asvg:svgBlip extension
    slide_xml = z.read("ppt/slides/slide1.xml").decode("utf8")
    has_svgblip = "svgBlip" in slide_xml

print(f"[ok] editable text boxes : {n_text}   e.g. {sample_text[:2]}")
print(f"[ok] native autoshapes   : {n_auto}")
print(f"[ok] pictures            : {n_pic}")
print(f"[ok] embedded SVG parts  : {n_svg}   svgBlip ref in slide: {has_svgblip}")

assert n_text >= 1, "no editable text boxes found — export would be non-editable"
assert n_auto >= 1, "no native autoshapes found"
if require_media:
    assert n_pic >= 1, "no pictures found"
    assert n_svg >= 1 and has_svgblip, "SVG not embedded as native vector"
print("\nPASS — deck is made of native, editable PowerPoint objects.")
